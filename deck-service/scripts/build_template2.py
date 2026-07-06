# -*- coding: utf-8 -*-
"""Build template2.pptx — a Superba-branded gallery of centered MBB consulting slides.
Composes on the real Superba masters (dark '1_Blank' gradient / light 'White'), so brand
background + footer logos are inherited. Headings Exo 2 italic, body Manrope, accents ruby/teal."""
import os
from pptx import Presentation
from pptx.util import Inches as In, Pt
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from PIL import Image

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
TEMPLATE = os.path.join(ROOT, "template.pptx")
ASSETS = os.path.join(ROOT, "assets")
OUT = os.path.join(ROOT, "template2.pptx")

def C(h): return RGBColor.from_string(h.lstrip('#'))
DEEP=C('163536'); POLAR=C('E9F7F8'); RUBY=C('E50A1A'); TEAL=C('60A09B'); NAVY=C('003462')
PALE=C('A9DBD5'); BRICK=C('BC393F'); WHITE=C('FFFFFF'); DKTEAL=C('185968'); CARD_DK=C('12333F')
HEAD='Exo 2'; BODY='Manrope'
CTR = PP_ALIGN.CENTER; LEFT = PP_ALIGN.LEFT; MID = MSO_ANCHOR.MIDDLE; TOP = MSO_ANCHOR.TOP
CX = 6.667  # slide centre x (13.333/2)

def get_layout(prs, mi, name):
    for lay in prs.slide_masters[mi].slide_layouts:
        if lay.name == name: return lay
    raise ValueError(f"layout {name!r} not on master {mi}")

def new_slide(prs, dark=True):
    lay = get_layout(prs, 0, '1_Blank') if dark else get_layout(prs, 1, 'White')
    s = prs.slides.add_slide(lay)
    for ph in list(s.placeholders):
        ph._element.getparent().remove(ph._element)
    return s

def tb(slide, l, t, w, h, text, *, font=BODY, size=13, color=POLAR, bold=False, italic=False,
       align=CTR, anchor=MID, spacing=1.0):
    box = slide.shapes.add_textbox(In(l), In(t), In(w), In(h)); tf = box.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = In(0.03); tf.margin_top = tf.margin_bottom = In(0.02)
    paras = text if isinstance(text, list) else [text]
    for i, p in enumerate(paras):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = align
        if spacing: para.line_spacing = spacing
        r = para.add_run(); r.text = p
        f = r.font; f.name = font; f.size = Pt(size); f.bold = bold; f.italic = italic; f.color.rgb = color
    return box

def tb_rich(slide, l, t, w, h, runs, *, font=HEAD, size=34, color=POLAR, bold=False, italic=True,
            align=CTR, anchor=MID, spacing=1.05):
    box = slide.shapes.add_textbox(In(l), In(t), In(w), In(h)); tf = box.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = In(0.03); tf.margin_top = tf.margin_bottom = In(0.02)
    p = tf.paragraphs[0]; p.alignment = align; p.line_spacing = spacing
    for text, col in runs:
        r = p.add_run(); r.text = text; f = r.font
        f.name = font; f.size = Pt(size); f.bold = bold; f.italic = italic; f.color.rgb = col or color
    return box

def kicker(slide, y, text, w=11.5):
    tb(slide, CX - w/2, y, w, 0.32, text.upper(), font=HEAD, size=12, color=RUBY, bold=True, italic=True)

def rrect(slide, l, t, w, h, *, fill=None, line=None, lw=0.75, radius=0.06):
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, In(l), In(t), In(w), In(h))
    sp.shadow.inherit = False
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb = line; sp.line.width = Pt(lw)
    try: sp.adjustments[0] = radius
    except Exception: pass
    return sp

def rect(slide, l, t, w, h, *, fill=None, line=None, lw=0.75):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, In(l), In(t), In(w), In(h)); sp.shadow.inherit = False
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb = line; sp.line.width = Pt(lw)
    return sp

def disc(slide, cx, cy, d, *, fill=RUBY, text=None, tcolor=WHITE, size=15, line=None, lw=1.5, italic=True):
    sp = slide.shapes.add_shape(MSO_SHAPE.OVAL, In(cx-d/2), In(cy-d/2), In(d), In(d)); sp.shadow.inherit = False
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb = line; sp.line.width = Pt(lw)
    if text is not None:
        tf = sp.text_frame; tf.margin_left=tf.margin_right=tf.margin_top=tf.margin_bottom=In(0)
        p = tf.paragraphs[0]; p.alignment = CTR
        r = p.add_run(); r.text = text; r.font.name = HEAD; r.font.size = Pt(size)
        r.font.bold = True; r.font.italic = italic; r.font.color.rgb = tcolor
    return sp

def hline(slide, cx, y, w, color=TEAL, weight=1.5, dash=None):
    ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, In(cx-w/2), In(y), In(cx+w/2), In(y))
    ln.line.color.rgb = color; ln.line.width = Pt(weight); ln.shadow.inherit = False
    if dash:
        d = ln.line._get_or_add_ln(); pd = d.makeelement(qn('a:prstDash'), {'val': dash}); d.append(pd)
    return ln

def vline(slide, x, y0, y1, color=TEAL, weight=1.0):
    ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, In(x), In(y0), In(x), In(y1))
    ln.line.color.rgb = color; ln.line.width = Pt(weight); ln.shadow.inherit = False
    return ln

def icon(slide, kw, cx, cy, size):
    p = os.path.join(ASSETS, f"generic_{kw}.png")
    if not os.path.exists(p): p = os.path.join(ASSETS, f"icon_{kw}.png")
    if not os.path.exists(p): return None
    with Image.open(p) as im: iw, ih = im.size
    if iw >= ih: dw, dh = size, size*ih/iw
    else: dh, dw = size, size*iw/ih
    return slide.shapes.add_picture(p, In(cx-dw/2), In(cy-dh/2), In(dw), In(dh))

def title(slide, y, text, *, dark=True, size=30, w=11.7):
    tb(slide, CX - w/2, y, w, 1.0, text, font=HEAD, size=size, italic=True,
       color=POLAR if dark else DEEP, anchor=TOP)

def body_color(dark): return POLAR if dark else DEEP
def mute_color(dark): return PALE if dark else TEAL

# ---------------------------------------------------------------- slides ----
def s_cover(prs):
    s = new_slide(prs, dark=True)
    kicker(s, 2.15, "Growth Strategy")
    tb(s, CX-6, 2.55, 12, 1.5, "Unlocking the Next Wave\nof Superba Krill".split("\n"),
       font=HEAD, size=46, italic=True, color=POLAR, spacing=1.05)
    hline(s, CX, 4.35, 2.4, TEAL, 1.75)
    tb(s, CX-5, 4.5, 10, 0.5, "Global market expansion strategy for premium krill oil",
       font=BODY, size=19, color=PALE)
    tb(s, CX-5, 5.05, 10, 0.4, "Prepared for Superba Krill AS   ·   July 2026",
       font=BODY, size=13, color=TEAL)

def s_exec(prs):
    s = new_slide(prs, dark=False)
    kicker(s, 0.6, "The Answer Up Front")
    title(s, 0.95, "Superba can triple krill-oil revenue by 2030 by\nleading three high-growth segments", dark=False, size=27)
    cards = [("1", "Own the premium supplement shelf",
              "Antarctic-sourced purity commands a 40% price premium in US and DACH nutraceutical channels."),
             ("2", "Convert the pet & aqua-feed pipeline",
              "Omega-3 feed demand is compounding at 12% a year with few certified suppliers."),
             ("3", "Move up the value chain",
              "Shift from bulk oil to branded consumer formulations to capture 3x the margin.")]
    n=3; gap=0.5; cw=(11.53-gap*(n-1))/n; y=2.7; ch=3.1; x0=CX-11.53/2
    for i,(num,ttl,bd) in enumerate(cards):
        x=x0+i*(cw+gap)
        rrect(s, x, y, cw, ch, fill=WHITE, line=TEAL, lw=1)
        disc(s, x+cw/2, y+0.55, 0.62, fill=RUBY, text=num, tcolor=WHITE, size=20)
        tb(s, x+0.2, y+1.0, cw-0.4, 0.9, ttl, font=HEAD, size=16, italic=True, color=DEEP, anchor=TOP)
        tb(s, x+0.25, y+1.85, cw-0.5, 1.1, bd, font=BODY, size=12.5, color=DKTEAL, anchor=TOP, spacing=1.05)

def s_agenda(prs):
    s = new_slide(prs, dark=True)
    kicker(s, 0.85, "What We'll Cover")
    title(s, 1.15, "Agenda", size=34)
    items = ["Market context & krill-oil demand drivers",
             "Segment sizing: supplements, pet, aqua-feed",
             "Competitive landscape & Superba's edge",
             "Growth strategy & prioritized moves",
             "Roadmap, investment & next steps"]
    bw=6.4; x=CX-bw/2; y0=2.55; rh=0.72
    for i,it in enumerate(items):
        y=y0+i*rh
        tb(s, x, y, 0.9, rh, f"{i+1:02d}", font=HEAD, size=22, italic=True, color=RUBY, align=LEFT)
        hline(s, x+1.35, y+rh/2, 0.0001, TEAL)  # placeholder no-op
        vline(s, x+1.05, y+0.12, y+rh-0.12, TEAL, 2)
        tb(s, x+1.3, y, bw-1.3, rh, it, font=HEAD, size=18, italic=True, color=POLAR, align=LEFT)

def s_section(prs):
    s = new_slide(prs, dark=True)
    tb(s, CX-3, 1.4, 6, 1.8, "03", font=HEAD, size=120, italic=True, color=DKTEAL)
    kicker(s, 3.35, "Section Three")
    title_txt = "Competitive Landscape"
    tb(s, CX-6, 3.65, 12, 0.9, title_txt, font=HEAD, size=40, italic=True, color=POLAR)
    hline(s, CX, 4.7, 1.6, RUBY, 2.5)

def s_matrix(prs):
    s = new_slide(prs, dark=False)
    kicker(s, 0.55, "Where To Play")
    title(s, 0.9, "Prioritizing Superba's Growth Moves", dark=False, size=26)
    # plot
    px0,py0,pw,ph = 4.55,1.95,4.2,3.75
    rect(s, px0, py0, pw, ph, fill=C('F2F8F9'), line=PALE, lw=1)
    vline(s, px0+pw/2, py0, py0+ph, PALE, 1)
    hline(s, px0+pw/2, py0+ph/2, pw, PALE, 1)
    # axes labels
    tb(s, px0-1.7, py0, 1.6, ph, "Strategic value  →", font=BODY, size=11, color=TEAL, align=CTR).rotation=270
    tb(s, px0, py0+ph+0.05, pw, 0.3, "Ease of execution  →", font=BODY, size=11, color=TEAL)
    quad=[("Big bets", px0+0.15, py0+0.12, LEFT), ("Quick wins", px0+pw-1.65, py0+0.12, PP_ALIGN.RIGHT),
          ("Question marks", px0+0.15, py0+ph-0.4, LEFT), ("Fill-ins", px0+pw-1.65, py0+ph-0.4, PP_ALIGN.RIGHT)]
    for t,x,y,al in quad:
        tb(s, x, y, 1.6, 0.3, t, font=HEAD, size=11.5, italic=True, color=DKTEAL, align=al)
    bub=[("A",7.75,2.75),("B",5.35,2.65),("C",6.55,4.05),("D",5.75,3.2),("E",7.6,4.95)]
    for L,x,y in bub: disc(s, x, y, 0.5, fill=RUBY, text=L, tcolor=WHITE, size=14)
    leg="A Premium US supplements   B DACH pharmacy   C Pet omega-3 feed   D Branded consumer line   E Aqua-feed contracts"
    tb(s, CX-6, 6.05, 12, 0.4, leg, font=BODY, size=11, color=DKTEAL)

def s_market(prs):
    s = new_slide(prs, dark=True)
    kicker(s, 0.6, "How Big Is The Prize")
    title(s, 0.95, "Sizing the Krill-Oil Opportunity", size=28)
    rings=[(5.2,"$2.1B","TAM","Global omega-3 & krill-oil market",C('1C4A57')),
           (3.5,"$780M","SAM","Premium human-grade krill oil, US + EU",C('2C6E74')),
           (1.9,"$190M","SOM","Winnable by Superba by 2030",TEAL)]
    cy=3.9
    for d,val,lab,desc,col in rings:
        o=s.shapes.add_shape(MSO_SHAPE.OVAL, In(CX-d/2), In(cy-d/2), In(d), In(d)); o.shadow.inherit=False
        o.fill.solid(); o.fill.fore_color.rgb=col; o.line.color.rgb=POLAR; o.line.width=Pt(0.75)
    labs=[("TAM","$2.1B","Global omega-3 & krill-oil market",cy-2.35),
          ("SAM","$780M","Premium human-grade, US + EU",cy-1.5),
          ("SOM","$190M","Winnable by Superba by 2030",cy-0.05)]
    for lab,val,desc,y in labs:
        tb(s, CX-2.0, y, 4.0, 0.34, f"{lab}   {val}", font=HEAD, size=17, italic=True, color=POLAR)
        tb(s, CX-2.6, y+0.32, 5.2, 0.3, desc, font=BODY, size=11.5, color=PALE)
    rrect(s, CX-5.6, 6.05, 11.2, 0.62, fill=RUBY, radius=0.12)
    tb(s, CX-5.4, 6.05, 10.8, 0.62, "A $190M obtainable pool = roughly 3x Superba's current krill-oil revenue on the table.",
       font=BODY, size=13, color=WHITE, bold=True)

def s_pillars3(prs):
    s = new_slide(prs, dark=False)
    kicker(s, 0.6, "Our Strategic Pillars")
    title(s, 0.95, "Three Foundations for Sustainable Growth", dark=False, size=26)
    cards=[("purity","Purity & provenance","Antarctic MSC-certified sourcing as the trust anchor for premium positioning."),
           ("molecule","Branded formulations","Move beyond bulk oil into consumer capsules, gummies and functional blends."),
           ("sustainability","Sustainable scale","Grow harvest capacity and traceability without eroding the eco-certification story.")]
    n=3; gap=0.55; cw=(11.4-gap*(n-1))/n; y=2.4; ch=3.3; x0=CX-11.4/2
    for i,(ic,ttl,bd) in enumerate(cards):
        x=x0+i*(cw+gap)
        rrect(s, x, y, cw, ch, fill=WHITE, line=TEAL, lw=1)
        rect(s, x, y, cw, 0.09, fill=RUBY)
        icon(s, ic, x+cw/2, y+1.0, 1.1)
        tb(s, x+0.2, y+1.65, cw-0.4, 0.6, ttl, font=HEAD, size=17, italic=True, color=DEEP, anchor=TOP)
        tb(s, x+0.3, y+2.35, cw-0.6, 0.9, bd, font=BODY, size=12.5, color=DKTEAL, anchor=TOP, spacing=1.05)

def s_pillars4(prs):
    s = new_slide(prs, dark=True)
    kicker(s, 0.5, "What It Takes To Win")
    title(s, 0.82, "Four Elements of a Winning Krill Strategy", size=26)
    hubx,huby=CX,4.15
    sats=[("sourcing","Supply security","Locked-in Antarctic quota & fleet",CX-3.55,3.05),
          ("award","Brand equity","Superba as the purity gold standard",CX+3.55,3.05),
          ("global","Channel access","Direct routes into US & DACH retail",CX-3.55,5.25),
          ("science","Innovation","New delivery formats & claims",CX+3.55,5.25)]
    for ic,ttl,sub,x,y in sats:
        ln=s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, In(hubx), In(huby), In(x), In(y))
        ln.line.color.rgb=TEAL; ln.line.width=Pt(1.25); ln.shadow.inherit=False
    for ic,ttl,sub,x,y in sats:
        disc(s, x, y, 1.15, fill=CARD_DK, line=TEAL, lw=1.25)
        icon(s, ic, x, y-0.12, 0.5)
        tb(s, x-1.5, y+0.62, 3.0, 0.3, ttl, font=HEAD, size=13.5, italic=True, color=POLAR)
        tb(s, x-1.6, y+0.92, 3.2, 0.3, sub, font=BODY, size=10.5, color=PALE)
    disc(s, hubx, huby, 1.5, fill=RUBY, text="Market\nleadership", tcolor=WHITE, size=13)

def s_valuechain(prs):
    s = new_slide(prs, dark=False)
    kicker(s, 0.6, "From Ocean To Shelf")
    title(s, 0.95, "The Superba Krill Value Chain", dark=False, size=26)
    stages=[("ocean","Harvest","Sustainable Antarctic krill under MSC quota"),
            ("process","Onboard processing","Fresh cold-extraction at sea"),
            ("science","Refining","Concentration to pharma-grade oil"),
            ("molecule","Formulation","Capsules, gummies & functional blends"),
            ("quality","Brand & sell","Premium products & margin capture")]
    n=5; gap=0.18; cw=(11.6-gap*(n-1))/n; y=2.55; ch=1.75; x0=CX-11.6/2
    for i,(ic,ttl,cap) in enumerate(stages):
        x=x0+i*(cw+gap)
        col = RUBY if i==n-1 else TEAL
        ch_shape=s.shapes.add_shape(MSO_SHAPE.CHEVRON, In(x), In(y), In(cw+0.28), In(ch)); ch_shape.shadow.inherit=False
        ch_shape.fill.solid(); ch_shape.fill.fore_color.rgb = C('EAF3F4') if i<n-1 else C('FBEAEA')
        ch_shape.line.color.rgb=col; ch_shape.line.width=Pt(1.25)
        icon(s, ic, x+cw/2-0.05, y+0.42, 0.55)
        tb(s, x+0.05, y+0.72, cw-0.1, 0.35, ttl, font=HEAD, size=12.5, italic=True, color=DEEP)
        tb(s, x+0.1, y+1.05, cw-0.15, 0.6, cap, font=BODY, size=9.5, color=DKTEAL, spacing=1.0)

def s_roadmap(prs):
    s = new_slide(prs, dark=True)
    kicker(s, 0.55, "How We Get There")
    title(s, 0.88, "Three-Horizon Growth Roadmap", size=26)
    spine_y=4.05; x0=1.5; x1=11.83
    hline(s, CX, spine_y, x1-x0, TEAL, 2.5)
    ph=[("2026","Secure the base","Lock quota · MSC re-cert · stabilize premium supply", True),
        ("2027-28","Build the brand","Launch US supplement line · enter DACH pharmacy", False),
        ("2029","Scale formats","Pet & aqua-feed contracts · gummy & functional SKUs", True),
        ("2030","Lead the category","Triple revenue · global purity standard", False)]
    xs=[x0+ (x1-x0)*k/3 for k in range(4)]
    for (tf_,ttl,bd,above),x in zip(ph,xs):
        disc(s, x, spine_y, 0.28, fill=RUBY)
        cy = spine_y-1.55 if above else spine_y+0.35
        rrect(s, x-1.35, cy, 2.7, 1.2, fill=CARD_DK, line=TEAL, lw=1)
        tb(s, x-1.3, cy+0.08, 2.6, 0.3, tf_, font=HEAD, size=13, italic=True, color=RUBY)
        tb(s, x-1.3, cy+0.38, 2.6, 0.3, ttl, font=HEAD, size=12.5, italic=True, color=POLAR)
        tb(s, x-1.28, cy+0.68, 2.56, 0.5, bd, font=BODY, size=9.5, color=PALE, spacing=1.0)

def s_compare(prs):
    s = new_slide(prs, dark=False)
    kicker(s, 0.55, "Why We Win")
    title(s, 0.88, "Superba vs. the Competition", dark=False, size=26)
    cols=["","Superba Krill","Aker","Rimfrost","Fish oil"]
    rows=[("Antarctic MSC certification",["●","◐","○","○"]),
          ("Cold-extraction freshness",["●","●","◐","○"]),
          ("Phospholipid omega-3 potency",["●","●","◐","○"]),
          ("Branded consumer presence",["◐","●","○","●"]),
          ("Value per gram omega-3",["●","◐","●","◐"])]
    x0=1.15; tw=11.03; c0=3.9; cw=(tw-c0)/4; y0=1.9; hh=0.55; rh=0.72
    rect(s, x0+c0, y0, cw*4, hh, fill=RUBY)
    rrect(s, x0+c0, y0+hh, cw, hh*0+rh*5, fill=C('EDF6F7'), radius=0.02)  # superba column highlight
    for j,cname in enumerate(cols):
        if j==0: continue
        tb(s, x0+c0+(j-1)*cw, y0, cw, hh, cname, font=HEAD, size=12.5, italic=True, color=WHITE)
    for i,(rn,vals) in enumerate(rows):
        y=y0+hh+i*rh
        tb(s, x0, y, c0-0.1, rh, rn, font=BODY, size=12, color=DKTEAL, align=LEFT)
        for j,v in enumerate(vals):
            col = TEAL if v=="●" else (PALE if v=="◐" else C('C9D6DA'))
            tb(s, x0+c0+j*cw, y, cw, rh, v, font=BODY, size=18, color=col)
        hline(s, x0+tw/2, y+rh, tw, C('DDE9EC'), 0.75)
    tb(s, CX-5.5, 6.25, 11, 0.4, "Superba leads on purity and potency; the gap to close is downstream brand presence.",
       font=BODY, size=11.5, italic=True, color=DKTEAL)

def s_funnel(prs):
    s = new_slide(prs, dark=True)
    kicker(s, 0.6, "The GTM Funnel")
    title(s, 0.95, "From Market to Sale", size=27)
    bands=[("Addressable buyers","48M omega-3 consumers (US+EU)",7.2,C('1C4A57')),
           ("Reachable via target channels","14M  ·  29%",5.6,C('2C6E74')),
           ("Premium-willing buyers","4.2M  ·  30%",4.0,TEAL),
           ("Superba conversions","0.9M  ·  21%",2.5,RUBY)]
    y=2.15; bh=0.82
    for i,(lab,val,w,col) in enumerate(bands):
        rrect(s, CX-w/2, y, w, bh, fill=col, radius=0.03)
        tb(s, CX-w/2, y+0.06, w, 0.42, lab, font=HEAD, size=14.5, italic=True, color=WHITE)
        tb(s, CX-w/2, y+0.44, w, 0.34, val, font=BODY, size=12, color=WHITE)
        y+=bh+0.13
    tb(s, CX-5.5, 6.15, 11, 0.4, "A 2-point lift in premium-buyer conversion adds ~$18M in annual revenue.",
       font=BODY, size=12.5, italic=True, color=PALE)

def s_waterfall(prs):
    s = new_slide(prs, dark=False)
    kicker(s, 0.55, "The Path To 3x")
    title(s, 0.88, "Bridging Today's Revenue to the 2030 Target", dark=False, size=25)
    bars=[("2026\nbaseline",0,65,TEAL),("Premium\nUS",65,113,C('7FC0B8')),
          ("DACH\npharmacy",113,144,C('7FC0B8')),("Pet &\naqua-feed",144,170,C('7FC0B8')),
          ("Legacy\nerosion",155,170,RUBY),("2030\ntarget",0,155,NAVY)]
    n=len(bars); x0=1.4; plotw=10.5; gap=0.35; bw=(plotw-gap*(n-1))/n
    base=5.7; top=2.35; vmax=180.0
    def yv(v): return base-(base-top)*v/vmax
    prev_x=None; prev_top=None
    for i,(lab,lo,hi,col) in enumerate(bars):
        x=x0+i*(bw+gap)
        y_hi=yv(hi); y_lo=yv(lo); h=y_lo-y_hi
        rect(s, x, y_hi, bw, h, fill=col)
        val = hi if lo==0 else (hi-lo)
        sign = "" if (lo==0 or hi>lo and i not in()) else ""
        disp = f"${hi}M" if lo==0 else (f"+${hi-lo}M" if col!=RUBY else f"-${hi-lo}M")
        tb(s, x-0.15, y_hi-0.32, bw+0.3, 0.3, disp, font=HEAD, size=11.5, italic=True, color=DEEP)
        tb(s, x-0.2, base+0.05, bw+0.4, 0.5, lab.split("\n"), font=BODY, size=9.5, color=DKTEAL, spacing=0.95)
        if prev_x is not None:
            hline(s, (prev_x+bw+x)/2, prev_top, (x)-(prev_x+bw), C('C9D6DA'), 1, dash='dash')
        prev_x=x; prev_top=y_hi if col!=RUBY else yv(155)
    tb(s, CX-5.5, 6.2, 11, 0.4, "New branded segments contribute >70% of the incremental $90M.",
       font=BODY, size=11.5, italic=True, color=DKTEAL)

def s_reco(prs):
    s = new_slide(prs, dark=True)
    kicker(s, 1.85, "Our Recommendation")
    tb_rich(s, CX-5.6, 2.5, 11.2, 1.9, [
        ("Superba should pivot from bulk supplier to ", POLAR),
        ("branded krill-oil leader", RUBY),
        (" — and invest now to ", POLAR),
        ("triple revenue by 2030.", RUBY)], font=HEAD, size=34, italic=True, spacing=1.1)
    hline(s, CX, 4.75, 2.0, TEAL, 2)
    tb(s, CX-5, 4.95, 10, 0.5, "Requires a $40M three-year investment in brand, capacity and channel; payback by 2029.",
       font=BODY, size=15, color=PALE)

def s_quote(prs):
    s = new_slide(prs, dark=False)
    tb(s, CX-1, 1.15, 2, 1.0, "“", font=HEAD, size=90, italic=True, color=RUBY, anchor=TOP)
    tb(s, CX-5.2, 2.35, 10.4, 2.0,
       "Consumers no longer just buy omega-3 — they buy provenance, purity and a story they can trust. That is exactly where Superba wins.",
       font=HEAD, size=25, italic=True, color=DEEP, spacing=1.1)
    hline(s, CX, 4.85, 1.6, TEAL, 2)
    tb(s, CX-5, 5.05, 10, 0.4, "Kari Nordahl — VP Global Marketing, Superba Krill AS",
       font=BODY, size=14, color=DKTEAL)

def main():
    prs = Presentation(TEMPLATE)
    lst = prs.slides._sldIdLst
    for sldId in list(lst):
        rId = sldId.get(qn('r:id'))
        if rId: prs.part.drop_rel(rId)
        lst.remove(sldId)
    for fn in [s_cover, s_exec, s_agenda, s_section, s_matrix, s_market, s_pillars3,
               s_pillars4, s_valuechain, s_roadmap, s_compare, s_funnel, s_waterfall,
               s_reco, s_quote]:
        fn(prs)
    prs.save(OUT)
    print("wrote", OUT, "slides:", len(prs.slides._sldIdLst))

if __name__ == "__main__":
    main()
