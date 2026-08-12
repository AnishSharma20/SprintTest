"""Visual QA gate — the last line of defence the schema can't provide.

The planner + validator guarantee the plan is well-formed and within length limits, and the
renderer guarantees on-brand styling. But nothing else LOOKS at the finished slides. This gate
does: it rasterises the rendered deck to images, sends them to one vision call that flags only
objective visual defects (text overflow / collision / truncation / an icon that doesn't match its
text), and returns short, actionable fixes. `pipeline.generate` then revises just the flagged
slides and re-renders.

Deliberately conservative: it reports on what's VISIBLE, never critiques wording/tone, and if no
rasteriser is available (no LibreOffice, not on Windows) it degrades to a no-op so generation never
breaks. Enabled only in "polished" mode (fast mode skips it).

Rasteriser: LibreOffice headless (portable — add `soffice` to the Docker image to enable on Render)
→ PDF → PyMuPDF; falls back to PowerPoint COM on a Windows dev box.
"""
from __future__ import annotations

import base64
import io
import os
import re
import shutil
import subprocess
import sys
import tempfile
from pathlib import Path

from . import config

GATE_MODEL = os.environ.get("DECK_GATE_MODEL", "").strip() or config.MODEL
_MAX_W = 1200  # downscale slides before the vision call to keep tokens/cost sane

ISSUE_ENUM = ["text_overflow", "collision", "truncation", "icon_mismatch",
              "empty_or_broken", "off_brand", "other"]

_SCHEMA = {
    "type": "object", "additionalProperties": False, "required": ["slides"],
    "properties": {"slides": {"type": "array", "items": {
        "type": "object", "additionalProperties": False, "required": ["slide", "ok"],
        "properties": {
            "slide": {"type": "integer", "description": "1-based slide number as labelled."},
            "ok": {"type": "boolean", "description": "true if the slide has NO visual defect."},
            "issues": {"type": "array", "items": {"enum": ISSUE_ENUM}},
            "fix": {"type": "string", "maxLength": 240,
                    "description": "If not ok: one short, concrete instruction the writer can apply "
                                   "(shorten a heading, change/drop a mismatched icon, trim body). Empty if ok."},
        }}}},
}

_SYSTEM = """You are a meticulous slide-design QA reviewer for on-brand corporate decks. You are shown
each rendered slide as an image, labelled 'Slide N' with a note of its layout and any icons it uses.

Flag ONLY objective, visible defects:
- text_overflow: text spills out of its box, off the slide edge, or overlaps the footer/logo.
- collision: two elements overlap or a title runs into the body/an image.
- truncation: a word or label is visibly cut off (e.g. ends mid-word or on a dangling "&"/"and").
- icon_mismatch: an icon's meaning clearly contradicts its heading/text (e.g. a heart icon on a
  liver point, a random object unrelated to the words). Judge against the labelled icon + what you see.
- empty_or_broken: an empty picture box, a missing image, obvious render corruption.
- off_brand: a jarring colour/element clearly outside the deep-sea / brand look.

Do NOT critique wording, tone, persuasiveness, or content choices — only what is visibly wrong.
Most slides should be ok:true. Be strict but not fussy: a slide that merely looks plain is fine.
For every not-ok slide give a SHORT, concrete `fix` the writer can apply. Report every slide."""


def _render_pngs(pptx: Path, out_dir: Path) -> bool:
    """LibreOffice (portable) → PDF → PyMuPDF; else PowerPoint COM (Windows). Same approach as
    scripts/qa.py. Returns False if no rasteriser is available."""
    out_dir.mkdir(parents=True, exist_ok=True)
    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    if soffice:
        # An isolated --env:UserInstallation profile per attempt, not soffice's shared default
        # one: without it, an invocation that overlaps with another (a concurrent request, or a
        # background generation job's own visual-QA rasterise running in another thread) can
        # silently fail to acquire the shared profile lock and exit 0 with NO pdf written. This
        # was a real, previously-latent bug: the old 4 MB upload cap made this endpoint rarely
        # exercised in production, so it never surfaced until larger files started actually
        # reaching LibreOffice on Render.
        #
        # Even isolated, a from-scratch profile bootstrap is itself a known source of headless
        # LibreOffice flakiness under a constrained container (seen live on Render: exit 0 but an
        # `SfxBaseModel::impl_store` write failure on the very first attempt) — so retry once
        # with a brand new temp dir + profile before giving up; a transient failure typically
        # does not repeat.
        last_error = None
        for attempt in range(2):
            with tempfile.TemporaryDirectory() as tmp:
                profile = Path(tmp) / "lo_profile"
                result = subprocess.run(
                    [soffice, "--headless", "--norestore",
                     f"-env:UserInstallation=file://{profile.as_posix()}",
                     "--convert-to", "pdf:impress_pdf_Export", "--outdir", tmp, str(pptx)],
                    capture_output=True, text=True,
                )
                pdfs = list(Path(tmp).glob("*.pdf"))
                if not pdfs:
                    detail = result.stderr.strip() or result.stdout.strip() or "no output"
                    # A write failure could plausibly be disk exhaustion on a constrained
                    # container (a large embedded image inflates during PDF re-encoding) — log
                    # free space so a recurrence is diagnosable instead of another guess.
                    try:
                        free_mb = shutil.disk_usage(tmp).free // (1024 * 1024)
                        detail += f" [free space at {tmp}: {free_mb} MB]"
                    except OSError:
                        pass
                    last_error = f"LibreOffice produced no PDF (exit {result.returncode}): {detail}"
                    continue
                import fitz  # PyMuPDF
                doc = fitz.open(str(pdfs[0]))
                for n, page in enumerate(doc, 1):
                    page.get_pixmap(dpi=110).save(str(out_dir / f"slide{n:03d}.png"))
                doc.close()
                return True
        raise RuntimeError(f"{last_error} (failed on both attempts)")
    if sys.platform.startswith("win"):
        ps = (f'$pp=New-Object -ComObject PowerPoint.Application;'
              f'$pres=$pp.Presentations.Open("{pptx}",$true,$true,$false);'
              f'$pres.Export("{out_dir}","PNG",1600,900);$pres.Close();$pp.Quit()')
        subprocess.run(["powershell", "-NoProfile", "-Command", ps], check=True, capture_output=True)
        return True
    return False


def _natkey(p: Path):
    m = re.findall(r"\d+", p.stem)
    return int(m[-1]) if m else 0


def _has_video(prs) -> bool:
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    return any(shape.shape_type == MSO_SHAPE_TYPE.MEDIA
               for slide in prs.slides for shape in slide.shapes)


def _extract_image_previews(prs, target_w: int = 1280) -> list[bytes] | None:
    """Fallback preview for a deck containing embedded video, which LibreOffice/PowerPoint's PDF
    export path cannot be trusted to handle (see rasterize() for the full story: profile isolation,
    an explicit export filter, a JRE, and shrinking a separate oversized image all failed to fix a
    real, reproducible write failure on Render caused by a slide's embedded mp4 — and three
    different attempts to strip/replace the video shape via python-pptx, verified locally, each
    left the FILE ITSELF something PowerPoint calls corrupted [0x80070570], even the most
    conservative one that only removed the shape and re-added a brand new plain picture via
    add_picture(). Isolated the trigger precisely: removing this ONE shape (regardless of what, if
    anything, replaces it) breaks python-pptx's own save/serialise round trip for this file;
    removing an unrelated shape on the same file does not. Given three independent surgical
    approaches all failed the same way, the safer engineering call is to stop trying to resave a
    modified copy of the pptx at all when video is present.

    For each slide, letterbox its LARGEST embedded picture (a video's own poster frame counts, and
    typically wins — it's usually the biggest thing on the slide) onto a plain canvas at the deck's
    own aspect ratio. Not a full slide render (no text, no layout), but a real, reliable image good
    enough for the 'pick which slide to add' picker — and this whole path only ever READS the
    original bytes, never re-saves them, so none of the resave fragility above can apply. Returns
    None if the deck has no picture content at all to show."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.oxml.ns import qn
    from PIL import Image

    sw, sh = prs.slide_width, prs.slide_height
    canvas_h = round(target_w * sh / sw)
    pngs = []
    any_image = False
    for slide in prs.slides:
        best_blob, best_area = None, 0
        for shape in slide.shapes:
            blob = None
            try:
                if shape.shape_type == MSO_SHAPE_TYPE.MEDIA:
                    blip = shape._element.find(f".//{qn('a:blip')}")
                    if blip is not None:
                        rid = blip.get(qn("r:embed"))
                        blob = slide.part.related_part(rid).blob
                elif hasattr(shape, "image"):
                    blob = shape.image.blob
            except Exception:  # noqa: BLE001 — one bad shape must not break the whole preview
                blob = None
            area = (shape.width or 0) * (shape.height or 0)
            if blob and area > best_area:
                best_blob, best_area = blob, area
        canvas = Image.new("RGB", (target_w, canvas_h), (255, 255, 255))
        if best_blob:
            try:
                im = Image.open(io.BytesIO(best_blob)).convert("RGB")
                ratio = min(target_w / im.width, canvas_h / im.height)
                new_size = (max(1, round(im.width * ratio)), max(1, round(im.height * ratio)))
                im = im.resize(new_size)
                canvas.paste(im, ((target_w - new_size[0]) // 2, (canvas_h - new_size[1]) // 2))
                any_image = True
            except Exception:  # noqa: BLE001 — fall through to the blank canvas for this slide
                pass
        buf = io.BytesIO()
        canvas.save(buf, "PNG")
        pngs.append(buf.getvalue())
    return pngs if any_image else None


_SHRINK_MAX_DIM = 1920  # comfortably above any on-screen preview need


def _shrink_pptx_images(data: bytes) -> bytes:
    """Downscale any oversized embedded image before handing the file to LibreOffice — what
    costs memory during PDF conversion is the DECODED bitmap (proportional to pixel count), not
    the compressed file size, and Render's free-tier 512 MB instance can run out of room holding
    just one large photo (confirmed live: a ~9 MB single-image .pptx consistently failed to
    convert there — SfxBaseModel::impl_store write failures — with disk space and every LibreOffice
    invocation flag already ruled out as the cause). Returns the ORIGINAL bytes unchanged if
    nothing needs shrinking or if anything about this step goes wrong — this is a memory
    optimisation for rasterising a temporary copy, never something that should affect what's
    actually stored or generated from."""
    import zipfile

    try:
        from PIL import Image

        src = zipfile.ZipFile(io.BytesIO(data))
        shrunk_any = False
        media_seen = []
        out_buf = io.BytesIO()
        with zipfile.ZipFile(out_buf, "w", zipfile.ZIP_DEFLATED) as dst:
            for item in src.infolist():
                content = src.read(item.filename)
                ext = item.filename.lower().rsplit(".", 1)[-1] if "." in item.filename else ""
                if item.filename.startswith("ppt/media/"):
                    if ext in ("jpg", "jpeg", "png", "bmp", "tiff", "gif"):
                        try:
                            im = Image.open(io.BytesIO(content))
                            media_seen.append(f"{item.filename} {im.size[0]}x{im.size[1]} {len(content)}b")
                            if max(im.size) > _SHRINK_MAX_DIM:
                                ratio = _SHRINK_MAX_DIM / max(im.size)
                                new_size = (max(1, round(im.width * ratio)), max(1, round(im.height * ratio)))
                                fmt = "JPEG" if ext in ("jpg", "jpeg") else im.format or "PNG"
                                if fmt == "JPEG" and im.mode not in ("RGB", "L"):
                                    im = im.convert("RGB")
                                im = im.resize(new_size)
                                buf = io.BytesIO()
                                im.save(buf, fmt, quality=85) if fmt == "JPEG" else im.save(buf, fmt)
                                content = buf.getvalue()
                                shrunk_any = True
                        except Exception as e:  # noqa: BLE001 — one bad image must not break the whole pass
                            media_seen.append(f"{item.filename} UNREADABLE ({e})")
                    else:
                        media_seen.append(f"{item.filename} {len(content)}b (not a raster type this pass handles)")
                dst.writestr(item, content)
        print(f"[qa-gate] shrink pass: {len(media_seen)} media item(s) — " + "; ".join(media_seen),
              file=sys.stderr)
        return out_buf.getvalue() if shrunk_any else data
    except Exception as e:  # noqa: BLE001 — this is an optimisation, never a hard requirement
        print(f"[qa-gate] image shrink skipped ({e}); rasterising the file as-is", file=sys.stderr)
        return data


def rasterize(pptx_bytes: bytes) -> list[bytes] | None:
    """Render a deck to one PNG per slide, in order. None if no rasteriser is available (gate off)."""
    try:
        from pptx import Presentation

        prs = Presentation(io.BytesIO(pptx_bytes))
        if _has_video(prs):
            # LibreOffice/PowerPoint's PDF export cannot be trusted with embedded video on this
            # environment (see _extract_image_previews' docstring for the full investigation) —
            # use the read-only image-extraction fallback instead of ever resaving this file.
            print("[qa-gate] embedded video detected; using the image-extraction fallback "
                  "preview instead of LibreOffice/PowerPoint", file=sys.stderr)
            return _extract_image_previews(prs)

        orig_size = len(pptx_bytes)
        pptx_bytes = _shrink_pptx_images(pptx_bytes)
        print(f"[qa-gate] pptx size before/after shrink: {orig_size}b / {len(pptx_bytes)}b", file=sys.stderr)
        with tempfile.TemporaryDirectory() as tmp:
            deck = Path(tmp) / "deck.pptx"
            deck.write_bytes(pptx_bytes)
            out = Path(tmp) / "png"
            if not _render_pngs(deck, out):
                return None
            # dedupe by lowercased name — Windows' filesystem is case-insensitive, so *.png and
            # *.PNG would each match the same COM-exported files and double the list.
            uniq = {p.name.lower(): p for p in [*out.glob("*.png"), *out.glob("*.PNG")]}
            pngs = sorted(uniq.values(), key=_natkey)
            return [p.read_bytes() for p in pngs] or None
    except Exception as e:  # noqa: BLE001 — the gate must never break generation
        print(f"[qa-gate] rasterise failed ({e}); skipping visual QA", file=sys.stderr)
        return None


def _jpeg_b64(png_bytes: bytes) -> str:
    from PIL import Image
    im = Image.open(io.BytesIO(png_bytes)).convert("RGB")
    if im.width > _MAX_W:
        im = im.resize((_MAX_W, round(im.height * _MAX_W / im.width)))
    buf = io.BytesIO()
    im.save(buf, "JPEG", quality=80)
    return base64.b64encode(buf.getvalue()).decode("ascii")


def _digest(i: int, spec: dict) -> str:
    parts = [f"Slide {i}: layout={spec.get('layout')}"]
    b = spec.get("benefit")
    if b and b != "none":
        parts.append(f"benefit-icon={b}")
    cols = spec.get("columns") or []
    if cols:
        ic = [f"{(c.get('heading') or '?')[:18]}→{c.get('icon') or c.get('icon_generic') or 'none'}"
              for c in cols]
        parts.append("column icons: " + "; ".join(ic))
    return " | ".join(parts)


def review(client, images: list[bytes], plan: dict, *, model: str | None = None) -> list[dict]:
    """One vision call over all slide images → list of per-slide findings (dicts with slide/ok/
    issues/fix). Returns [] on any error (gate must not break generation)."""
    slides = plan.get("slides", [])
    n = min(len(images), len(slides))
    if n == 0:
        return []
    content: list[dict] = [{"type": "text", "text":
        "Review each of these rendered slides for visual defects and report via report_qa."}]
    for i in range(n):
        content.append({"type": "text", "text": _digest(i + 1, slides[i])})
        content.append({"type": "image", "source": {"type": "base64",
                        "media_type": "image/jpeg", "data": _jpeg_b64(images[i])}})
    try:
        msg = client.messages.create(
            model=model or GATE_MODEL, max_tokens=2000, system=_SYSTEM,
            tools=[{"name": "report_qa", "description": "Report per-slide visual QA findings.",
                    "input_schema": _SCHEMA}],
            tool_choice={"type": "tool", "name": "report_qa"},
            messages=[{"role": "user", "content": content}],
        )
        for block in msg.content:
            if block.type == "tool_use" and isinstance(block.input, dict):
                return block.input.get("slides", [])
    except Exception as e:  # noqa: BLE001
        print(f"[qa-gate] vision review failed ({e}); accepting deck as-is", file=sys.stderr)
    return []


def flagged(findings: list[dict]) -> list[dict]:
    """Findings that are not ok and carry an actionable fix (1-based slide + fix text)."""
    out = []
    for f in findings:
        if not f.get("ok", True) and (f.get("fix") or "").strip() and isinstance(f.get("slide"), int):
            out.append(f)
    return out
