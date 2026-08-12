"""Stage 2 — deterministic renderer.

JSON plan -> python-pptx fills the real Superba template. All styling (fonts, colours,
backgrounds, logos, bullet formatting) is inherited from the template's slide layouts; this
code only chooses which layout to instantiate, drops text into placeholders by index, and
inserts a photo. It never sets a font, colour, or position.

Key disciplines that make template-fill clean (the earlier attempt failed on these):
- strip the 64 example slides first;
- pick the dark (master #0) or light (master #1) variant of a layout by `background`;
- fill placeholders by their INDEX from the layout catalog (not by guessing);
- DELETE any content placeholder the plan didn't fill, so no empty prompt text or tinted
  picture boxes survive (footer/date/slide-number placeholders are left to inherit).
"""
from __future__ import annotations

import copy
import datetime
import io
import math
import re
import sys
import tempfile
import threading
from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls, qn
from pptx.util import Emu, Inches, Pt

from . import brand as _brand_theme
from . import config

CHROME_IDX = {10, 11, 12}   # date / footer / slide-number — never fill, never remove

# Serialize render_deck access: the renderer modifies module-level globals (font, size, margins,
# custom photos) at call time. Without a lock, concurrent renders stomp each other's settings.
_RENDER_LOCK = threading.Lock()

# trailing connector words a truncation must not end on (English + Norwegian)
_ORPHANS = {"and", "og", "of", "the", "to", "for", "with", "og", "&", "i", "på", "med", "av"}


def _delete_example_slides(prs) -> None:
    lst = prs.slides._sldIdLst
    for sldId in list(lst):
        rId = sldId.get(qn("r:id"))
        if rId:
            prs.part.drop_rel(rId)
        lst.remove(sldId)


def _master_indices(brand: str | None = None):
    """(dark, light) master indices for a brand's template.

    Superba ships a dark master and its light twin, both on the theme's major font, next to
    leftover default Office masters — so the light one is "another master using the SAME major
    font as the primary", which finds it without naming a typeface. A brand whose template has
    only ONE master (Revervia) gets that master for both, i.e. its single background is used
    whichever theme a slide asks for."""
    inv = config.inventory(brand)
    dark = inv.get("primary_master_index", inv.get("superba_master_index", 0))
    primary_font = next((m.get("major_font") for m in inv["masters"] if m["index"] == dark), None)
    light = next((m["index"] for m in inv["masters"]
                  if m["index"] != dark and primary_font
                  and (m.get("major_font") or "").lower() == primary_font.lower()), dark)
    return dark, light


# Type scale — exactly 3 text sizes, deck-wide, no exceptions (hero figures and footer chrome
# fold into this same scale, and native template placeholders are forced onto it too, not just
# code-built layouts). Hierarchy is expressed through WEIGHT, COLOUR and CAPS, not extra sizes.
# Defined early: default args below (_set_text/_set_lines) are evaluated at def-time.
_SZ_HERO = 60                          # hero DATA FIGURES (stat values, KPI numbers)
_SZ_COVER = 60                         # the "big statement" title layouts (cover, agenda,
                                        # highlight). Superba uses one number for both roles, which
                                        # is why they were a single constant until multi-brand:
                                        # Revervia's cover title box is 4.6in wide, so a 60pt cover
                                        # title fits ~10 characters. Split so a brand can keep big
                                        # stat figures without an unusable cover title limit.
_SZ_TITLE = 32                         # regular slide titles (every other layout, native or code-built)
_SZ_SUBTITLE = 16                      # cover subtitle, column/text-picture headings
_SZ_BODY = 14                          # body, bullets, table cells, items
_SZ_SMALL = 12                         # eyebrows, captions, notes, footnotes, axis/step labels, footer
                                        # chrome — the floor; nothing in the deck goes smaller
_SZ_AGENDA_ITEMS = 16                  # DELIBERATE exception to the scale above (client spec): the
                                        # agenda's contents list, always Manrope BOLD at this size,
                                        # never the deck-wide body size/weight.

# When the About page's design settings override the body font, native template placeholders must
# be forced onto it too (they otherwise inherit the template theme's font and the deck would mix
# two body typefaces). None = no override, inherit as designed. Set by apply_design().
_FORCE_BODY_FONT: str | None = None
# Same idea for line spacing: code-built text always draws at _LINE_SPACING, but native
# placeholders inherit the template's own spacing UNLESS the user overrode it on the About page.
_FORCE_LINE_SPACING: bool = False
# Footer controls (About page design settings): page numbers on/off, an optional standing
# footer line (e.g. "Confidential, for internal use") and an optional render-date stamp —
# all drawn by _stamp_footer in one centred line. The AI disclaimer is NOT controlled here
# (house rule: every generated deck carries it).
_PAGE_NUMBERS: bool = True
_FOOTER_TEXT: str = ""
_DATE_STAMP: bool = False
# icon_level "none": the team turned brand icons off entirely. Gated where icon paths are
# resolved (_icon_path/_generic_icon_path), so every icon call site inherits the switch.
_ICONS_OFF: bool = False
# The team's own uploaded photos, registered per render: asset key -> temp file path. Kept
# alongside a TemporaryDirectory handle so the files outlive the render that uses them.
# Which brand the in-flight render belongs to. A module global rather than a parameter threaded
# through 31 _fill_* functions, for the same reason _ICONS_OFF and _CUSTOM_PHOTO_PATHS are:
# _RENDER_LOCK serializes the whole of render_deck, so only one brand is ever in flight.
_BRAND: str | None = None
_LIGHT_ONLY: bool = False    # brand's template has no dark master; see apply_brand/_make_slide
_HAS_BENEFITS_SLIDE: bool = True   # brand's template carries the verbatim benefits overview slide
_CUSTOM_PHOTO_PATHS: dict[str, "Path"] = {}
_CUSTOM_PHOTO_DIR = None


def _find_layout(prs, name, master_index):
    master = prs.slide_masters[master_index]
    for lay in master.slide_layouts:
        if lay.name == name:
            return lay
    for m in prs.slide_masters:                 # fallback: any master
        for lay in m.slide_layouts:
            if lay.name == name:
                return lay
    raise ValueError(f"Layout '{name}' not found in template")


def _autofit(tf, *, shrink: bool = True) -> None:
    """Text-frame fit policy. Short single-value labels (title/heading) keep shrink-to-fit as a
    graceful last-resort guard against a 1-char overflow. Multi-line bodies/lists use `shrink=False`
    (MSO_AUTO_SIZE.NONE) so text stays at its fixed size — content is CAPPED by the schema char
    limits rather than shrunk to cram more in (client typography rule). Font/colour inherit."""
    tf.word_wrap = True
    try:
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE if shrink else MSO_AUTO_SIZE.NONE
    except Exception:  # noqa: BLE001
        pass


def _shrink_to_fit(ph, text: str, base_pt: float, min_pt: float = 18) -> None:
    """Deterministically shrink a placeholder's font so `text` fits its box (width x height). The
    template's own shrink-to-fit is unreliable in headless render, so a long title would otherwise
    overflow past its box (into the footer, or onto the content below). Estimates lines from the box
    width and steps the size down until the wrapped text fits the box height. Only ever shrinks."""
    w_in = (ph.width or 0) / 914400.0
    h_in = (ph.height or 0) / 914400.0
    if not text or w_in <= 0 or h_in <= 0:
        return
    pt = base_pt
    while pt > min_pt:
        cpl = max(1.0, (w_in * 72.0) / (pt * 0.52))     # chars per line at this size
        lines = max(1, math.ceil(len(text) / cpl))
        if lines * pt * 1.2 / 72.0 <= h_in:             # wrapped height fits the box
            break
        pt -= 1
    for p in ph.text_frame.paragraphs:
        for r in p.runs:
            r.font.size = Pt(pt)


def _set_ph_box(ph, l, t, w, h) -> None:
    """Reposition/resize a placeholder by writing its full a:xfrm directly. Safer than the python-pptx
    left/top setters, which raise on placeholders that inherit their geometry (no spPr/xfrm yet).
    l, t, w, h are EMU (ints, e.g. Inches(...))."""
    sp = ph._element
    spPr = sp.find(qn("p:spPr"))
    if spPr is None:
        spPr = sp.makeelement(qn("p:spPr"), {})
        nv = sp.find(qn("p:nvSpPr")) or sp.find(qn("p:nvPicPr"))
        (nv.addnext(spPr) if nv is not None else sp.insert(0, spPr))
    for x in spPr.findall(qn("a:xfrm")):
        spPr.remove(x)
    xfrm = spPr.makeelement(qn("a:xfrm"), {})
    off = xfrm.makeelement(qn("a:off"), {"x": str(int(l)), "y": str(int(t))})
    ext = xfrm.makeelement(qn("a:ext"), {"cx": str(int(w)), "cy": str(int(h))})
    xfrm.append(off); xfrm.append(ext)
    spPr.insert(0, xfrm)


def _set_run_size(tf, size: float, *, bold=None, italic=None, font=None) -> None:
    """Force every run in a text frame to one explicit size, overriding whatever the template's
    own layout/master style would otherwise inherit (54-60pt titles, 13-16pt bodies, etc.) — part
    of keeping the deck to exactly 3 text sizes total, not just within the code-built layouts.
    bold/italic/font are left alone (None) unless explicitly given — used for the title role, whose
    NATIVE placeholders otherwise inherit the theme's italic major font at a light weight, while
    code-built layouts explicitly force bold, non-italic Exo 2. Same slide role, different look."""
    for p in tf.paragraphs:
        for r in p.runs:
            r.font.size = Pt(size)
            if bold is not None:
                r.font.bold = bold
            if italic is not None:
                r.font.italic = italic
            if font is not None:
                r.font.name = font


def _set_text(ph, text: str, size: float | None = None, **style) -> None:
    """Single logical value into a placeholder. Size is forced explicitly (not inherited from the
    layout) so native template placeholders stay on the same 3-size scale as code-built layouts.
    size=None means "the current body size" — resolved at CALL time, not def time, so the About
    page's design overrides (apply_design) actually reach these defaults."""
    if size is None:
        size = _SZ_BODY
    if _FORCE_BODY_FONT and "font" not in style:
        style["font"] = _FORCE_BODY_FONT
    ph.text_frame.text = text
    _autofit(ph.text_frame)
    if _FORCE_LINE_SPACING:
        for p in ph.text_frame.paragraphs:
            p.line_spacing = _LINE_SPACING
    _set_run_size(ph.text_frame, size, **style)


def _set_lines(ph, lines: list[str], bullet_rid: str | None = None, size: float | None = None, **style) -> None:
    """Multiple paragraphs (bullets / agenda items). If bullet_rid is given, each paragraph gets
    the brand's PICTURE bullet (the teal figure embedded in the template master) — this overrides
    the content placeholders' `buNone`. Otherwise each paragraph inherits the layout's list format.
    Size is forced explicitly, same reasoning (and same call-time default) as _set_text."""
    if size is None:
        size = _SZ_BODY
    if _FORCE_BODY_FONT and "font" not in style:
        style["font"] = _FORCE_BODY_FONT
    lines = [ln for ln in (l.strip() for l in lines) if ln]
    tf = ph.text_frame
    tf.text = lines[0] if lines else ""
    for ln in lines[1:]:
        tf.add_paragraph().text = ln
    if bullet_rid:
        for para in tf.paragraphs:
            _apply_picture_bullet(para._p, bullet_rid)
    if _FORCE_LINE_SPACING:
        for para in tf.paragraphs:
            para.line_spacing = _LINE_SPACING
    _autofit(tf, shrink=False)   # bodies/lists: cap content, do not shrink
    _set_run_size(tf, size, **style)


# The brand bullet is a small teal PNG embedded in the template master (as a picture bullet at
# body level 1). Content placeholders switch it off with buNone, so for real bullet lists we set
# the picture bullet explicitly on each paragraph, matching the master's indent metrics.
_BULLET_MARL, _BULLET_INDENT, _BULLET_SZ = "342900", "-342900", "100000"


def _bullet_rid(slide) -> str | None:
    """Embed the brand bullet image in the slide part (idempotent) and return its relationship id."""
    path = config.assets_dir(_BRAND) / "bullet.png"
    if not path.exists():
        return None
    _, rid = slide.part.get_or_add_image_part(str(path))
    return rid


def _apply_picture_bullet(p, rid: str) -> None:
    """Set the brand picture bullet on one <a:p>, replacing any inherited buNone/buChar."""
    pPr = p.find(qn("a:pPr"))
    if pPr is None:
        pPr = p.makeelement(qn("a:pPr"), {})
        p.insert(0, pPr)
    pPr.set("marL", _BULLET_MARL)
    pPr.set("indent", _BULLET_INDENT)
    for tag in ("a:buClr", "a:buSzPct", "a:buSzPts", "a:buFont", "a:buFontTx",
                "a:buNone", "a:buChar", "a:buAutoNum", "a:buBlip"):
        for el in pPr.findall(qn(tag)):
            pPr.remove(el)
    # Order matters in the schema: bullet size (buSz*) before the bullet itself (bu*).
    buSz = pPr.makeelement(qn("a:buSzPct"), {"val": _BULLET_SZ})
    buBlip = pPr.makeelement(qn("a:buBlip"), {})
    blip = buBlip.makeelement(qn("a:blip"), {qn("r:embed"): rid})
    buBlip.append(blip)
    pPr.append(buSz)
    pPr.append(buBlip)


DISCLAIMER = ("AI generated draft from the source material. Review all content, claims and figures, "
              "and edit as needed before use.")


def _add_disclaimer(slide, dark: bool) -> None:
    """Add a small 'AI generated, review before use' note along the bottom of the cover slide.
    A free-standing textbox (no template placeholder exists for it), so we set a subtle size/colour."""
    box = slide.shapes.add_textbox(Inches(0.5), Inches(6.98), Inches(10.5), Inches(0.4))
    tf = box.text_frame
    tf.word_wrap = True
    tf.text = DISCLAIMER
    run = tf.paragraphs[0].runs[0]
    run.font.size = Pt(_SZ_SMALL)
    run.font.italic = True
    run.font.color.rgb = RGBColor(0xBF, 0xE3, 0xEF) if dark else RGBColor(0x6B, 0x8B, 0x95)


def _icon_path(benefit: str):
    """Resolve a health-benefit tag to its staged branded icon (one brand-red line-art colourway
    that reads on both masters). Returns None if there is no icon for that benefit — or if the
    team turned icons off entirely (icon_level "none"): every icon call site resolves through
    here or _generic_icon_path, so gating the two covers the whole deck deterministically."""
    if _ICONS_OFF or not benefit or benefit == "none":
        return None
    entry = config.asset_index(_BRAND).get(f"icon_{benefit}")
    if not entry or not entry.get("path"):
        return None
    p = config.resolve_asset(entry["path"], _BRAND)
    return p if p.exists() else None


def _generic_icon_path(keyword: str):
    """Resolve a generic-library keyword to its staged fallback icon (same brand-red line-art).
    Used only when a slide can't be fully covered by the branded benefit icons. Same
    icons-off gate as _icon_path."""
    if _ICONS_OFF or not keyword or keyword == "none":
        return None
    entry = config.asset_index(_BRAND).get(f"generic_{keyword}")
    if not entry or not entry.get("path"):
        return None
    p = config.resolve_asset(entry["path"], _BRAND)
    return p if p.exists() else None


def register_custom_photos(photos) -> None:
    """Stage the team's uploaded photos (About page photo library) for this render: each
    {key, bytes} is written to a temp file and resolvable through _photo_path just like a
    built-in library photo. Called by render_deck; an empty/None list clears the registry."""
    global _CUSTOM_PHOTO_PATHS, _CUSTOM_PHOTO_DIR
    _CUSTOM_PHOTO_PATHS = {}
    _CUSTOM_PHOTO_DIR = None
    if not photos:
        return
    _CUSTOM_PHOTO_DIR = tempfile.TemporaryDirectory(prefix="team_photos_")
    for p in photos:
        key, data = p.get("key"), p.get("bytes")
        if not key or not data:
            continue
        path = Path(_CUSTOM_PHOTO_DIR.name) / f"{key}.jpg"
        path.write_bytes(data)
        _CUSTOM_PHOTO_PATHS[key] = path


def _photo_path(aid: str):
    """Resolve a photo asset id to a file: the team's uploaded photos first (team_photo_*),
    then the built-in manifest. None when unknown or missing on disk."""
    p = _CUSTOM_PHOTO_PATHS.get(aid)
    if p is not None:
        return p if p.exists() else None
    entry = config.asset_index(_BRAND).get(aid)
    if not entry or not entry.get("path"):
        return None
    p = config.resolve_asset(entry["path"], _BRAND)
    return p if p.exists() else None


def _layout_box(layout_name: str, master_index: int, idx: int):
    """(left, top, width, height) in EMU for a placeholder, from the inventory (dims already
    resolved through master inheritance)."""
    for lay in config.inventory(_BRAND)["layouts"]:
        if lay["name"] == layout_name and lay["master_index"] == master_index:
            for p in lay["placeholders"]:
                if p["idx"] == idx and p["width_emu"] and p["height_emu"]:
                    return p["left_emu"], p["top_emu"], p["width_emu"], p["height_emu"]
    return None


def _place_icon(slide, box, icon_path) -> bool:
    """Add an icon scaled to FIT (letterbox), centred in the box — icons must not be
    crop-to-filled the way insert_picture would. box = (left, top, width, height) in EMU."""
    if not box or not icon_path:
        return False
    left, top, w, h = box
    try:
        from PIL import Image
        with Image.open(icon_path) as im:
            iw, ih = im.size
        if (iw / ih) > (w / h):
            dw, dh = w, int(w * ih / iw)
        else:
            dh, dw = h, int(h * iw / ih)
        slide.shapes.add_picture(str(icon_path), left + (w - dw) // 2, top + (h - dh) // 2, dw, dh)
        return True
    except Exception as e:  # noqa: BLE001 — an icon is decorative; never break the render
        print(f"[renderer] icon placement failed for {icon_path} ({e}); slide loses this icon "
              f"silently — see qa_geometry's asset check", file=sys.stderr)
        return False


def _place_cropped(slide, path, l, t, w, h) -> bool:
    """Add a photo that FILLS the box (cover), centre-cropping the overflow instead of stretching —
    the opposite of _place_icon's letterbox fit. Sizes in inches."""
    try:
        pic = slide.shapes.add_picture(str(path), Inches(l), Inches(t), Inches(w), Inches(h))
        from PIL import Image
        with Image.open(path) as im:
            iw, ih = im.size
        src, box = iw / ih, w / h
        if src > box:                       # source too wide → trim the sides
            keep = box / src
            pic.crop_left = pic.crop_right = (1 - keep) / 2
        elif src < box:                     # source too tall → trim top and bottom
            keep = src / box
            pic.crop_top = pic.crop_bottom = (1 - keep) / 2
        return True
    except Exception as e:  # noqa: BLE001 — fall back to the caller's solid panel
        print(f"[renderer] photo placement failed for {path} ({e}); slide falls back to its "
              f"solid panel silently — see qa_geometry's asset check", file=sys.stderr)
        return False


def _fit(text, limit):
    """Hard word-boundary truncation for collision-prone 1-line label fields (cover/agenda
    title, headings). The planner + validation retry keep these within the limit almost
    always; this is the last-resort guarantee that a stray over-limit label can never wrap
    into a 2nd line and collide with the element below. Not applied to bodies/items (they
    grow into empty space or auto-fit)."""
    if not (text and limit and len(text) > limit):
        return text
    words = text.split()
    if len(words[0]) > limit:      # a single word already exceeds the box — keep it whole (no mid-word cut)
        return words[0]
    out = ""
    for w in words:                # pack whole words up to the limit
        if len(out) + len(w) + (1 if out else 0) > limit:
            break
        out = f"{out} {w}".strip()
    out = out or words[0]
    # A cut can leave a dangling connector or symbol ("Omega-3 EPA &", "Heart and"); drop it so
    # the label never ends on an orphan token.
    out = re.sub(r"[\s&+/,-]+$", "", out)
    while len(out.split()) > 1 and out.rsplit(" ", 1)[1].lower() in _ORPHANS:
        out = re.sub(r"[\s&+/,-]+$", "", out.rsplit(" ", 1)[0])
    return out.strip() or words[0]


def _distinct_col_headings(raws, limit):
    """Fit each column heading to its (narrow, 1-line) box AND guarantee the columns don't render
    the SAME heading. Two parallel headings that share an opening ("What the barrier does" / "…
    needs", or "Superba supports heart" / "… brain") would otherwise both truncate to the shared
    prefix. When that happens we drop the common leading words and re-fit the distinguishing tail
    (which also turns "Superba supports heart/brain" into a clean "Heart"/"Brain")."""
    fitted = [_fit(h or "", limit) for h in raws]
    ne = [h for h in fitted if h]
    if len(set(ne)) == len(ne):
        return fitted
    toks = [(h or "").split() for h in raws]
    present = [t for t in toks if t]
    common = 0
    if len(present) > 1:
        shortest = min(len(t) for t in present)
        for i in range(shortest):
            w = present[0][i].lower()
            if i < shortest - 1 and all(t[i].lower() == w for t in present):  # never strip all words
                common += 1
            else:
                break
    if common:
        stripped = [" ".join(t[common:]) for t in toks]
        stripped = [s[:1].upper() + s[1:] if s else "" for s in stripped]
        refit = [_fit(s, limit) for s in stripped]
        rne = [h for h in refit if h]
        if len(set(rne)) == len(rne):
            return refit
    return fitted


def _fill_slide(slide, spec: dict, cat: dict, master_index: int, dark: bool) -> None:
    fields = cat["fields"]
    lim = cat.get("limits", {})
    layout_name = cat["template_layout"]
    phmap = {ph.placeholder_format.idx: ph for ph in slide.placeholders}
    filled: set[int] = set()
    benefit = spec.get("benefit")
    benefit = None if benefit in (None, "none") else benefit

    def put(idx, value, multiline=False, bullets=False, size=_SZ_BODY, **style):
        if idx is None or idx not in phmap or value is None:
            return
        if multiline:
            raw = value if isinstance(value, list) else str(value).split("\n")
            lines = [ln for ln in (str(x).strip() for x in raw) if ln]
            # A list of points (2+ lines) gets the brand picture bullet; a single line stays plain prose.
            rid = _bullet_rid(slide) if (bullets and len(lines) > 1) else None
            _set_lines(phmap[idx], lines, bullet_rid=rid, size=size, **style)
        else:
            _set_text(phmap[idx], str(value), size=size, **style)
        filled.add(idx)

    title = spec.get("title")
    if cat["kind"] in ("title", "agenda"):     # narrow 1-line title box above a neighbour
        title = _fit(title, lim.get("title"))
    # "Big statement" layouts (cover, agenda, highlight) get the hero size, matching what these
    # specific layouts used natively; every other kind gets the regular title size.
    title_size = _SZ_COVER if cat["kind"] in ("title", "agenda", "highlight") else _SZ_TITLE
    # Explicit, not inherited: two native "kind"s (highlight, text_picture) resolve their title
    # placeholder to "Exo 2 Light italic" instead of the template's dominant "Exo 2 italic" — an
    # inconsistency in the template itself. Force the SAME cut everywhere so every title, whatever
    # layout renders it, looks identical.
    put(fields.get("title"), title, size=title_size, bold=False, italic=False, font=_HEAD_TITLE)
    put(fields.get("subtitle"), spec.get("subtitle"), size=_SZ_SUBTITLE)
    put(fields.get("heading"), _fit(spec.get("heading"), lim.get("heading")), size=_SZ_SUBTITLE)
    put(fields.get("body"), spec.get("body"), multiline=True, bullets=True)
    if spec.get("items"):
        if cat["kind"] == "agenda":
            # Client spec: the agenda's contents list is ALWAYS Manrope Bold at _SZ_AGENDA_ITEMS,
            # never the deck-wide body size/weight (only "agenda" maps an "items" field here — see
            # config/layout_catalog.json — so this never reaches another kind).
            put(fields.get("items"), spec["items"], multiline=True, bullets=True,
                size=_SZ_AGENDA_ITEMS, bold=True, font=_BODY)
        else:
            put(fields.get("items"), spec["items"], multiline=True, bullets=True)

    col_head_max = (lim.get("columns") or {}).get("heading_max")
    col_maps = fields.get("columns", [])
    cols = spec.get("columns", [])
    # Icon consistency ACROSS the whole slide (brand rule): every column gets an icon or NONE do,
    # all from ONE source (all AKBM benefit icons OR all generic fallback icons — never mixed),
    # each icon distinct. Prefer the branded benefit icons; fall back to the generic set only
    # when every column can be matched there. If neither source covers all columns, drop icons
    # from the whole slide rather than render a partial / duplicated / mixed set.
    def _consistent(paths):
        strs = [str(p) for p in paths]
        return paths if paths and all(paths) and len(set(strs)) == len(strs) else None
    icons = (_consistent([_icon_path(c.get("icon")) for c in cols])
             or _consistent([_generic_icon_path(c.get("icon_generic")) for c in cols])
             or [None] * len(cols))
    heads = _distinct_col_headings([c.get("heading") for c in cols], col_head_max)
    for col_map, col, icon, head in zip(col_maps, cols, icons, heads):
        put(col_map.get("heading"), head, size=_SZ_SUBTITLE)
        # Content-driven: a column body written as several lines becomes bullets; a single
        # sentence stays prose. Same rule as the main body.
        put(col_map.get("body"), col.get("body"), multiline=True, bullets=True)
        pic = col_map.get("picture")
        if icon and pic is not None:
            _place_icon(slide, _layout_box(layout_name, master_index, pic), icon)

    aid = spec.get("asset_id")
    pic_idx = fields.get("picture")
    if aid:
        if pic_idx is not None and pic_idx in phmap:
            path = _photo_path(aid)
            if path is not None:
                phmap[pic_idx].insert_picture(str(path))
                filled.add(pic_idx)
    elif benefit and pic_idx is not None and pic_idx in phmap:
        # Benefit slide with a picture area but no photo → show the benefit icon there.
        _place_icon(slide, _layout_box(layout_name, master_index, pic_idx), _icon_path(benefit))

    # Benefit icon on text-only benefit slides (highlight / section have open top-left space).
    if benefit and cat["kind"] in ("highlight", "section"):
        _place_icon(slide, (Inches(0.5), Inches(0.42), Inches(0.95), Inches(0.95)), _icon_path(benefit))

    # Text-with-picture: make the title FULL WIDTH (the narrow template title box overflows a long
    # takeaway), drop the unused sub-heading, and place the body (left) and picture (right) BELOW the
    # title with a fixed margin so the body never touches the title.
    if cat["kind"] == "text_picture":
        tph = phmap.get(fields.get("title"))
        if tph is not None:
            _set_ph_box(tph, Inches(_MARGIN), Inches(0.746), Inches(_CONTENT_W), Inches(1.0))
            # base==min: re-assert the fixed title size at the new box (no shrink below the scale).
            _shrink_to_fit(tph, str(title or ""), base_pt=_SZ_TITLE, min_pt=_SZ_TITLE)
        hidx = fields.get("heading")
        if hidx in filled and hidx in phmap:
            phmap[hidx]._element.getparent().remove(phmap[hidx]._element)
            filled.discard(hidx)
        body_top = 2.1                               # title bottom ~1.75 + a fixed ~0.35 margin
        bidx = fields.get("body"); bph = phmap.get(bidx)
        if bph is not None and bidx in filled:
            _set_ph_box(bph, Inches(_MARGIN), Inches(body_top), Inches(4.1), Inches(_BODY_BOTTOM - body_top))
        pidx = fields.get("picture")
        if pidx is not None and pidx in filled:
            # insert_picture replaced the placeholder element, so re-fetch it fresh from the slide.
            pph = next((p for p in slide.placeholders if p.placeholder_format.idx == pidx), None)
            if pph is not None:
                _set_ph_box(pph, Inches(5.0), Inches(body_top), Inches(7.83), Inches(_BODY_BOTTOM - body_top))

    # Section divider: a long section title overflows its short template box down into the footer.
    # Give it a taller box (clear of the footer); font stays fixed at the title size (no shrink
    # below the deck-wide scale) — the taller box is what actually prevents overflow now.
    if cat["kind"] == "section":
        stph = phmap.get(fields.get("title"))
        if stph is not None and title and len(str(title)) > 40:
            _set_ph_box(stph, Inches(0.97), Inches(3.2), Inches(7.30), Inches(3.0))
            _shrink_to_fit(stph, str(title), base_pt=_SZ_TITLE, min_pt=_SZ_TITLE)

    # A takeaway title can run to two lines, but some layouts (Text Slide, Title Only) have a
    # ONE-line title box, so the second line collides with the content below — and the opposite:
    # every "below" box's position was authored assuming the layout's OWN much taller inherited
    # title (up to 60pt), so at the deck-wide _SZ_TITLE a short title leaves a large dead gap above
    # the content (seen concretely on Agenda and multi-column layouts). Both directions are the same
    # fix: snap everything stacked below the title to sit a fixed, designer-preserved gap under the
    # title's REAL rendered height, shifting every "below" box by the SAME delta so their spacing
    # relative to each other (e.g. a column heading above its own body) is unchanged.
    title_idx = fields.get("title")
    tbox = _layout_box(layout_name, master_index, title_idx) if title_idx is not None else None
    # text_picture re-lays out its own title/body/picture above (with its own margin), so the generic
    # push/pull must not re-read the LAYOUT boxes and fight it.
    if cat["kind"] == "text_picture":
        tbox = None
    if tbox and title:
        # How tall will the title REALLY be? Use title_size (the same per-kind value put() just
        # rendered it at — 60 for the hero-tier kinds, 18 otherwise), not a single fixed constant,
        # since the two now differ substantially and the gap math must match reality.
        lines = _est_lines(str(title), tbox[2] / 914400.0, title_size)
        title_h = Inches(lines * title_size * 1.25 / 72.0)  # 1.25 ≈ single line spacing, rounded up
        # Only boxes stacked BELOW the title qualify — never a side-by-side picture that legally
        # starts level with (or above) the title box.
        below = [(i, _layout_box(layout_name, master_index, i)) for i in filled
                 if i != title_idx and i not in CHROME_IDX]
        below = [(i, b) for i, b in below if b and b[1] >= tbox[1] + tbox[3] - Inches(0.05)]
        if below:
            # Preserve the gap the DESIGNER put between the title box and the nearest thing below it.
            top_of_below = min(b[1] for _, b in below)
            gap = max(int(Inches(0.22)), top_of_below - (tbox[1] + tbox[3]))
            safe_top = tbox[1] + int(title_h) + gap
            shift = safe_top - top_of_below     # positive = push down, negative = pull up
            if shift != 0:
                for idx, box in below:
                    ph = phmap.get(idx)
                    if ph is None:
                        continue
                    ph.left, ph.width = Emu(box[0]), Emu(box[2])
                    ph.top = Emu(box[1] + shift)
                    ph.height = Emu(max(int(Inches(0.6)), box[3] - shift))

    # AI-generated disclaimer along the bottom of the cover slide.
    if cat["kind"] == "title":
        _add_disclaimer(slide, dark)

    # Remove every content placeholder we did not fill (prevents empty picture boxes /
    # leftover prompt text). Chrome placeholders (date/footer/number) stay and inherit.
    for ph in list(slide.placeholders):
        idx = ph.placeholder_format.idx
        if idx in CHROME_IDX or idx in filled:
            continue
        ph._element.getparent().remove(ph._element)

_R_ATTRS = (qn("r:embed"), qn("r:link"), qn("r:id"))
_DESIGN_SRC: dict | None = None


def _design_source():
    """Cache-load template.pptx — the SINGLE design file. It holds the master layouts + theme +
    logos AND the verbatim source slides (ingredient, benefits) that the renderer splices in, so all
    design the AI uses lives in one place.

    Keyed by brand: a single cached Presentation would serve the FIRST brand rendered to every
    brand after it, splicing the wrong template's slides."""
    global _DESIGN_SRC
    if _DESIGN_SRC is None:
        _DESIGN_SRC = {}
    if _BRAND not in _DESIGN_SRC:
        _DESIGN_SRC[_BRAND] = Presentation(str(config.template_path(_BRAND)))
    return _DESIGN_SRC[_BRAND]


def _find_design_slide(marker: str):
    """Find a verbatim source slide inside template.pptx by a case-insensitive text marker."""
    m = marker.upper()
    for s in _design_source().slides:
        for sh in s.shapes:
            if sh.has_text_frame and m in sh.text_frame.text.upper():
                return s
    raise ValueError(f"Design source slide not found in template.pptx (marker {marker!r}).")


def _ph_marker(el):
    """The p:ph element of a copied shape (sp OR pic placeholder), or None."""
    for nvPr in el.iter(qn("p:nvPr")):
        ph = nvPr.find(qn("p:ph"))
        if ph is not None:
            return ph
    return None


def _bake_placeholder(el, box, lst) -> None:
    """Turn a COPIED placeholder element into a free-standing shape: explicit geometry from the
    SOURCE file's own inheritance chain, the source layout's list-style text defaults baked into
    its txBody, and the p:ph marker removed — so nothing re-inherits from the DESTINATION blank
    layout (which has no matching placeholder and would collapse it to default position/styling)."""
    spPr = el.find(qn("p:spPr"))
    if spPr is not None and box is not None and spPr.find(qn("a:xfrm")) is None:
        spPr.insert(0, parse_xml(
            f'<a:xfrm {nsdecls("a")}><a:off x="{box[0]}" y="{box[1]}"/>'
            f'<a:ext cx="{box[2]}" cy="{box[3]}"/></a:xfrm>'))
    if (spPr is not None and el.tag == qn("p:sp")
            and spPr.find(qn("a:prstGeom")) is None and spPr.find(qn("a:custGeom")) is None):
        xfrm = spPr.find(qn("a:xfrm"))
        geom = parse_xml(f'<a:prstGeom {nsdecls("a")} prst="rect"><a:avLst/></a:prstGeom>')
        spPr.insert(list(spPr).index(xfrm) + 1 if xfrm is not None else 0, geom)
    txBody = el.find(qn("p:txBody"))
    if txBody is not None and lst is not None and len(lst):
        own = txBody.find(qn("a:lstStyle"))
        if own is None or not len(own):
            if own is not None:
                txBody.remove(own)
            bodyPr = txBody.find(qn("a:bodyPr"))
            txBody.insert(list(txBody).index(bodyPr) + 1 if bodyPr is not None else 0,
                          copy.deepcopy(lst))
    ph = _ph_marker(el)
    if ph is not None:
        ph.getparent().remove(ph)


def _splice_shapes(dst_slide, src_slide, *, strict: bool = False,
                   bake_placeholders: bool = False) -> bool:
    """Deep-copy every shape of src_slide onto dst_slide, re-embedding images and re-linking
    external hyperlinks — the one verbatim-splice mechanism (ingredient, benefits, and the
    About page's team-uploaded slides all ride on it). Internal references that are NOT images
    (an embedded chart, video, OLE object) cannot be carried over: with strict=True the whole
    splice is refused (returns False, dst untouched) so the caller can fall back to a rendered
    PNG; with strict=False just those shapes are skipped. PLACEHOLDER shapes are also unsafe in
    strict mode: their geometry AND text formatting inherit from the source deck's layout, which
    does not travel with them, so a copied placeholder lands at a default position with default
    styling (verified: a spliced two-column slide collapsed to centre-stacked unstyled text).
    bake_placeholders=True (the layout-override path) lifts that restriction by BAKING each
    placeholder instead: explicit geometry resolved through the source's own layout chain, the
    layout's lstStyle copied in, and the p:ph marker removed — chrome placeholders (date/footer/
    slide number) are dropped entirely."""
    rmap = dict(src_slide.part.rels.items())

    baked: dict[int, dict] = {}
    if bake_placeholders:
        lay_lst: dict[int, object] = {}
        try:
            for lp in src_slide.slide_layout.placeholders:
                tx = lp._element.find(qn("p:txBody"))
                lst = tx.find(qn("a:lstStyle")) if tx is not None else None
                if lst is not None and len(lst):
                    lay_lst[lp.placeholder_format.idx] = lst
        except Exception:  # noqa: BLE001 — style baking is best-effort
            lay_lst = {}
        for ph in src_slide.placeholders:
            idx = ph.placeholder_format.idx
            try:
                box = (int(ph.left), int(ph.top), int(ph.width), int(ph.height)) \
                    if None not in (ph.left, ph.top, ph.width, ph.height) else None
            except (TypeError, ValueError):
                box = None
            baked[idx] = {"box": box, "lst": lay_lst.get(idx)}

    def unsafe(element) -> bool:
        if strict and not bake_placeholders and element.find(qn("p:nvSpPr")) is not None:
            nvPr = element.find(qn("p:nvSpPr")).find(qn("p:nvPr"))
            if nvPr is not None and nvPr.find(qn("p:ph")) is not None:
                return True
        for node in element.iter():
            for a in _R_ATTRS:
                rid = node.get(a)
                if rid and rid in rmap:
                    rel = rmap[rid]
                    if not rel.is_external and "image" not in rel.reltype:
                        return True
        return False

    if strict and any(unsafe(shp._element) for shp in src_slide.shapes):
        return False
    spTree = dst_slide.shapes._spTree
    for shp in src_slide.shapes:
        if unsafe(shp._element):
            continue
        el = copy.deepcopy(shp._element)
        if bake_placeholders:
            ph = _ph_marker(el)
            if ph is not None:
                idx = int(ph.get("idx") or 0)
                if idx in CHROME_IDX:
                    continue   # date/footer/slide-number chrome never travels
                info = baked.get(idx) or {}
                _bake_placeholder(el, info.get("box"), info.get("lst"))
        for node in el.iter():                    # remap every relationship reference in the copy
            for a in _R_ATTRS:
                if a in node.attrib:
                    rel = rmap.get(node.get(a))
                    if rel is None:
                        continue
                    if rel.is_external:
                        new = dst_slide.part.rels.get_or_add_ext_rel(rel.reltype, rel.target_ref)
                    else:
                        _, new = dst_slide.part.get_or_add_image_part(io.BytesIO(rel._target.blob))
                    node.set(a, new)
        spTree.append(el)
    return True


def _add_ingredient_slide(prs, master_index: int) -> None:
    """Insert AKBM's standard ingredient slide VERBATIM — the exact slide they always use — by
    splicing its self-contained shape tree into the deck. Fidelity is perfect: the slide carries
    its own full-bleed background, so the host layout (a Blank one) is completely hidden behind
    it. Content is FIXED (the product composition never changes), so nothing here is generated."""
    src_slide = _find_design_slide("Cellular Nutrient")
    slide = prs.slides.add_slide(_find_layout(prs, "Blank", master_index))
    for ph in list(slide.shapes):                 # drop the Blank layout's own placeholders
        ph._element.getparent().remove(ph._element)
    _splice_shapes(slide, src_slide)


def _blank_layout(prs, master_index: int):
    master = prs.slide_masters[master_index]
    for lay in master.slide_layouts:
        if "blank" in lay.name.lower():
            return lay
    return master.slide_layouts[-1]


def _set_bg(slide, rgb: RGBColor) -> None:
    """Force a solid slide background of the given colour, overriding whatever the host master
    would otherwise show (used for the white/pastel light-theme variants and the white benefits
    infographic)."""
    cSld = slide._element.find(qn("p:cSld"))
    for old in cSld.findall(qn("p:bg")):
        cSld.remove(old)
    cSld.insert(0, parse_xml(
        f'<p:bg {nsdecls("p", "a")}><p:bgPr><a:solidFill><a:srgbClr val="{rgb}"/></a:solidFill>'
        f'<a:effectLst/></p:bgPr></p:bg>'))


def _set_white_bg(slide) -> None:
    """Force a solid-white slide background (the benefits slide is a white infographic; the host
    master is the dark deep-sea one)."""
    _set_bg(slide, _WHITE)


def _add_benefits_slide(prs, master_index: int) -> None:
    """Splice AKBM's verbatim benefits-overview slide onto a white background, using the LIGHT master's
    blank layout so the footer logos are the light-background (red/dark) colourway. Content is FIXED."""
    src_slide = _find_design_slide("PROVEN HEALTH BENEFITS")
    slide = prs.slides.add_slide(_blank_layout(prs, master_index))
    for ph in list(slide.shapes):
        ph._element.getparent().remove(ph._element)
    _set_white_bg(slide)
    _splice_shapes(slide, src_slide)


def _copy_slide_bg(dst_slide, src_slide) -> None:
    """Carry a foreign slide's EFFECTIVE background over: its own p:bg if it has one, else its
    layout's, else its master's — re-embedding a picture fill from whichever part owned it. A team
    slide spliced without this would sit on OUR deep-sea master and look nothing like its preview.
    No background found anywhere → dst keeps the host master's (rare; every deck has SOME bg)."""
    for holder in (src_slide, src_slide.slide_layout, src_slide.slide_layout.slide_master):
        cSld = holder._element.find(qn("p:cSld"))
        bg = cSld.find(qn("p:bg")) if cSld is not None else None
        if bg is None:
            continue
        el = copy.deepcopy(bg)
        rmap = dict(holder.part.rels.items())
        for node in el.iter():
            for a in _R_ATTRS:
                if a in node.attrib:
                    rel = rmap.get(node.get(a))
                    if rel is None or rel.is_external or "image" not in rel.reltype:
                        continue
                    _, new = dst_slide.part.get_or_add_image_part(io.BytesIO(rel._target.blob))
                    node.set(a, new)
        dcSld = dst_slide._element.find(qn("p:cSld"))
        for old in dcSld.findall(qn("p:bg")):
            dcSld.remove(old)
        dcSld.insert(0, el)
        return


def _add_custom_slide(prs, master_index: int, pptx_bytes: bytes, slide_index: int,
                      png_bytes: bytes | None = None, skip_number: set | None = None) -> None:
    """Splice one slide of a TEAM-UPLOADED deck in verbatim (the About page's "your slides").
    Shapes are copied with images re-embedded and the source's own background carried over, so
    the slide stays EDITABLE in the generated deck. Slides using features the splice cannot
    carry (embedded charts / media / OLE) fall back to their stored full-slide PNG — pixel
    perfect, just not editable. Never raises: a team slide must not fail the whole deck."""
    try:
        src = Presentation(io.BytesIO(pptx_bytes))
        slides = list(src.slides)
        if not 0 <= slide_index < len(slides):
            print(f"[custom-slide] slide index {slide_index} out of range ({len(slides)} slides); skipped",
                  file=sys.stderr)
            return
        src_slide = slides[slide_index]
        slide = prs.slides.add_slide(_blank_layout(prs, master_index))
        # The source slide carries its OWN chrome (page number, logos) — verbatim means verbatim,
        # so our own page-number pass must skip this slide or the two numbers overprint.
        if skip_number is not None:
            skip_number.add(slide.slide_id)
        for ph in list(slide.shapes):
            ph._element.getparent().remove(ph._element)
        _copy_slide_bg(slide, src_slide)
        if _splice_shapes(slide, src_slide, strict=True):
            return
        if png_bytes:
            # Letterbox the rendered PNG (source decks can be 4:3) on a white background.
            from PIL import Image
            with Image.open(io.BytesIO(png_bytes)) as im:
                iw, ih = im.size
            _set_white_bg(slide)
            sw, sh = prs.slide_width, prs.slide_height
            scale = min(sw / iw, sh / ih)
            w, h = int(iw * scale), int(ih * scale)
            slide.shapes.add_picture(io.BytesIO(png_bytes), (sw - w) // 2, (sh - h) // 2, w, h)
            return
        # No stored preview to fall back on: splice what CAN be carried rather than nothing.
        _splice_shapes(slide, src_slide, strict=False)
    except Exception as e:  # noqa: BLE001
        print(f"[custom-slide] splice failed ({e}); slide skipped", file=sys.stderr)


def _rewrite_txbody(sp, text: str) -> None:
    """Replace a shape's text while keeping its DESIGN: per paragraph, the pPr and the first
    run's formatting survive, surplus runs/breaks go; extra plan lines clone the last original
    paragraph's formatting; surplus original paragraphs are removed (always keeping one). Field
    runs (slide number/date) are preserved untouched."""
    txBody = sp.find(qn("p:txBody"))
    if txBody is None:
        return
    # Stale autofit factors were computed for the ORIGINAL text; drop them so the new text
    # renders at the authored size (PowerPoint recomputes; LibreOffice honours the attributes).
    bodyPr = txBody.find(qn("a:bodyPr"))
    if bodyPr is not None:
        norm = bodyPr.find(qn("a:normAutofit"))
        if norm is not None:
            norm.attrib.pop("fontScale", None)
            norm.attrib.pop("lnSpcReduction", None)

    paras = txBody.findall(qn("a:p"))
    if not paras:
        return
    template = copy.deepcopy(paras[-1])   # pristine formatting for any extra plan lines

    def set_para(p, line: str) -> None:
        first = None
        for child in list(p):
            if child.tag == qn("a:r"):
                if first is None:
                    first = child
                else:
                    p.remove(child)
            elif child.tag == qn("a:br"):
                p.remove(child)
        if first is None:
            end_pr = p.find(qn("a:endParaRPr"))
            first = parse_xml(f'<a:r {nsdecls("a")}><a:t/></a:r>')
            if end_pr is not None:
                r_pr = copy.deepcopy(end_pr)
                r_pr.tag = qn("a:rPr")
                first.insert(0, r_pr)
                p.insert(list(p).index(end_pr), first)
            else:
                p.append(first)
        ts = first.findall(qn("a:t"))
        if not ts:
            t = parse_xml(f'<a:t {nsdecls("a")}/>')
            first.append(t)
            ts = [t]
        ts[0].text = line
        for extra in ts[1:]:
            first.remove(extra)

    lines = text.split("\n") or [text]
    for i, line in enumerate(lines):
        if i < len(paras):
            set_para(paras[i], line)
        else:
            clone = copy.deepcopy(template)
            set_para(clone, line)
            txBody.append(clone)
    for p in paras[len(lines):]:
        txBody.remove(p)


def _refill_slots(slide, slots_meta: list[dict], plan_slots: dict) -> None:
    """Write the plan's per-slot text into the just-spliced override shapes, matched by the
    stable cNvPr@id the splice preserves. A slot the plan left blank keeps its designed text —
    never blank a designed box. Text is capped at the slot's measured budget, never shrunk."""
    from .overrides import truncate_budget
    by_id = {s.get("slot_id"): s for s in (slots_meta or []) if s.get("slot_id")}
    if not by_id or not plan_slots:
        return
    for sp in slide.shapes._spTree.iter(qn("p:sp")):   # iter() reaches shapes inside groups too
        cnv = sp.find(qn("p:nvSpPr") + "/" + qn("p:cNvPr"))
        if cnv is None:
            continue
        slot_id = f"s{cnv.get('id')}"
        meta = by_id.get(slot_id)
        text = str(plan_slots.get(slot_id) or "").strip()
        if meta is None or not text:
            continue
        _rewrite_txbody(sp, truncate_budget(text, int(meta.get("char_budget") or 800)))


def _add_override_slide(prs, master_index: int, ov: dict, spec: dict,
                        skip_number: set | None = None) -> None:
    """Splice a TEAM-REDESIGNED layout (an About page design override) and refill its measured
    text slots with the plan's fresh text — the recipe path: their design, the AI's content.
    Placeholders are baked (bake_placeholders) since a downloaded native layout is placeholder
    built. Never raises: like a team slide, an override must not fail the whole deck."""
    try:
        src = Presentation(io.BytesIO(ov["bytes"]))
        slides = list(src.slides)
        idx = int(ov.get("index") or 0)
        if not 0 <= idx < len(slides):
            print(f"[override] {ov.get('layout')}: slide index {idx} out of range "
                  f"({len(slides)} slides); skipped", file=sys.stderr)
            return
        src_slide = slides[idx]
        slide = prs.slides.add_slide(_blank_layout(prs, master_index))
        # The user's redesign may carry its own chrome — our footer pass must skip it.
        if skip_number is not None:
            skip_number.add(slide.slide_id)
        for ph in list(slide.shapes):
            ph._element.getparent().remove(ph._element)
        _copy_slide_bg(slide, src_slide)
        if _splice_shapes(slide, src_slide, strict=True, bake_placeholders=True):
            _refill_slots(slide, ov.get("slots") or [], spec.get("slots") or {})
            return
        png = ov.get("png")
        if png:
            # Shouldn't happen (inspect-slots pre-screens uploads), but the file may have been
            # edited since: freeze the design as a picture rather than ship nothing.
            print(f"[override] {ov.get('layout')} not refillable (splice refused); "
                  f"shipped as picture", file=sys.stderr)
            from PIL import Image
            with Image.open(io.BytesIO(png)) as im:
                iw, ih = im.size
            _set_white_bg(slide)
            sw, sh = prs.slide_width, prs.slide_height
            scale = min(sw / iw, sh / ih)
            w, h = int(iw * scale), int(ih * scale)
            slide.shapes.add_picture(io.BytesIO(png), (sw - w) // 2, (sh - h) // 2, w, h)
            return
        _splice_shapes(slide, src_slide, strict=False, bake_placeholders=True)
        _refill_slots(slide, ov.get("slots") or [], spec.get("slots") or {})
    except Exception as e:  # noqa: BLE001
        print(f"[override] {ov.get('layout')} splice failed ({e}); slide skipped", file=sys.stderr)


# ---------------------------------------------------------------------------
# Synthetic (code-built) layouts — mechanism B. The renderer reproduces a faithful structure on a
# Blank layout (inheriting the master's background + logos) and fills it from the plan: text into
# slots, AI-picked brand icons into circles, or a native chart. Brand palette / fonts below.
# ---------------------------------------------------------------------------
_RED = RGBColor(0xE5, 0x0A, 0x1A)
_TEAL = RGBColor(0x18, 0x59, 0x68)
_TEAL2 = RGBColor(0x2C, 0x74, 0x82)   # secondary panel teal
_PANEL = RGBColor(0xE4, 0xF1, 0xF1)
_INKC = RGBColor(0x16, 0x35, 0x36)
_LTEAL = RGBColor(0xA9, 0xDB, 0xD5)
_ONTEAL = RGBColor(0xEC, 0xF5, 0xF5)  # body text ON a solid teal panel — high contrast (LTEAL was too dim)
_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
_HEAD, _BODY = "Exo 2", "+mn-lt"   # body = the theme MINOR font (embedded Manrope), referenced the same
                                   # way the template placeholders do; a hard-coded "Manrope" can bind to a
                                   # wrong/cursive installed variant instead of the embedded regular one.
_HEAD_TITLE = "Exo 2 italic"       # slide TITLES ONLY: the template's own theme major font and the
                                   # brand guide's actual Exo 2 cut (guide lists only "Italic, Semibold
                                   # Italic" — there is no non-italic Exo 2 in the brand). Non-bold.
_CHART_COLORS = [_RED, _TEAL2, _LTEAL, RGBColor(0x60, 0xA0, 0x9B)]
_TBL_LINE = "C9D9D9"                  # table row-line colour (hex, for XML)

# ── Consultancy design system — ONE fixed skeleton for every synthetic slide ──────────────
# Canvas is 13.333 x 7.5 in (16:9). The title, eyebrow, body zone and footer occupy identical
# positions on every slide, so nothing shifts when moving between slides. All side-by-side
# boxes use the same gutter; parallel boxes get equal heights. See _synth_slide().
_MARGIN = 0.5                          # page margin — matches the template content-title LEFT (0.5)
_CONTENT_W = 13.333 - 2 * _MARGIN      # 12.333 in usable width
_TITLE_Y, _TITLE_H = 0.746, 1.3        # title TOP matches the template's content layouts exactly, so the
                                        # title never shifts between a synthetic and a template slide.
                                        # _TITLE_H covers a worst-case 2-line title at _SZ_TITLE (32pt
                                        # needs ~1.11in for 2 lines); recalibrate together with
                                        # _EYEBROW_Y/_BODY_TOP below if _SZ_TITLE changes again — this
                                        # skeleton is FIXED, not dynamic (unlike the native-template
                                        # title, which self-adjusts to its own real height).
_EYEBROW_Y = 2.1                       # below the enlarged title zone; see _TITLE_H note
_BODY_TOP, _BODY_BOTTOM = 2.55, 6.7    # fixed body zone (the footer band lives below 6.7)
_BODY_H = _BODY_BOTTOM - _BODY_TOP     # 4.6 in
_GUTTER = 0.3                          # the ONE gutter between all side-by-side boxes
_PAD = 0.22                            # inner padding inside panels
_LINE_SPACING = 1.06                   # fixed line spacing, applied everywhere

_STEP_BADGE = 0.5                      # numbered step / timeline node badge diameter
_ICON_DISC = 0.9                       # icon-circle diameter
_BOX = MSO_SHAPE.RECTANGLE            # one shape style for content boxes (square, consulting look)

# ── User design overrides (the About page's "Design settings") ────────────────────────────
# The brand defaults above stay the single source of truth; this layer lets the TEAM (never the
# LLM) deterministically restyle the drawn output: fonts, the three text sizes, line spacing,
# page margin and box gutter. apply_design() is called once at the top of render_deck() — it
# first RESETS everything to the recorded defaults (so one job's overrides can never leak into
# the next) and then applies the given overrides. Subtitle and hero sizes scale with the title
# size so the type hierarchy keeps its designed proportions.
_DESIGN_DEFAULTS = {
    "_SZ_HERO": _SZ_HERO, "_SZ_COVER": _SZ_COVER, "_SZ_TITLE": _SZ_TITLE,
    "_SZ_SUBTITLE": _SZ_SUBTITLE,
    "_SZ_BODY": _SZ_BODY, "_SZ_SMALL": _SZ_SMALL,
    "_LINE_SPACING": _LINE_SPACING, "_MARGIN": _MARGIN, "_GUTTER": _GUTTER,
    "_HEAD": _HEAD, "_BODY": _BODY, "_HEAD_TITLE": _HEAD_TITLE,
    "_FORCE_BODY_FONT": None, "_FORCE_LINE_SPACING": False,
    "_PAGE_NUMBERS": True, "_FOOTER_TEXT": "", "_DATE_STAMP": False, "_ICONS_OFF": False,
}


def _hex(h: str) -> RGBColor:
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def apply_brand(brand: str | None = None) -> None:
    """Point the palette/typography globals at one brand's theme (src/brand.py).

    Called at the top of every render, BEFORE apply_design(), and it rewrites the font and
    size entries of _DESIGN_DEFAULTS too — otherwise apply_design's reset-to-defaults would
    immediately undo the brand and every deck would come out in Superba's typography. The
    team's own design overrides are applied after this, so they still win.

    Colours are only ever read from these globals by the code-built layouts. Native template
    slides keep inheriting their own .pptx theme, untouched, exactly as before."""
    t = _brand_theme.theme(brand)
    c, f, s = t["colors"], t["fonts"], t["sizes"]
    g = globals()
    g["_BRAND"] = brand
    g["_RED"] = _hex(c["accent"])
    g["_TEAL"] = _hex(c["deep"])
    g["_TEAL2"] = _hex(c["deep2"])
    g["_PANEL"] = _hex(c["panel"])
    g["_INKC"] = _hex(c["ink"])
    g["_LTEAL"] = _hex(c["tint"])
    g["_ONTEAL"] = _hex(c["on_deep"])
    g["_TBL_LINE"] = c["table_line"]
    g["_CHART_COLORS"] = [_hex(h) for h in t["chart_colors"]]
    # Rounded vs square content boxes. Superba's design system is deliberately square (the
    # "consulting look"); Revervia's identity is rounded. _BOX is the single shape token every
    # code-built layout draws its panels with, so this is the whole change.
    g["_BOX"] = MSO_SHAPE.ROUNDED_RECTANGLE if t.get("rounded") else MSO_SHAPE.RECTANGLE
    g["_LIGHT_ONLY"] = bool(t.get("light_only"))
    g["_HAS_BENEFITS_SLIDE"] = bool(t.get("has_benefits_slide"))
    brand_type = {
        "_HEAD": f["head"], "_BODY": f["body"], "_HEAD_TITLE": f["title"],
        "_SZ_TITLE": s["title"], "_SZ_BODY": s["body"], "_SZ_SMALL": s["small"],
        "_SZ_HERO": s["hero"], "_SZ_COVER": s["cover"], "_SZ_SUBTITLE": s["subtitle"],
    }
    g.update(brand_type)
    _DESIGN_DEFAULTS.update(brand_type)


def apply_design(design: dict | None) -> None:
    """Reset the design globals to brand defaults, then lay the given overrides on top.
    Unknown keys and out-of-range values are ignored (the API validates ranges; this is the
    server-side belt to that suspender). Safe to call with None/{} — that IS the reset."""
    g = globals()
    g.update(_DESIGN_DEFAULTS)
    g["_CONTENT_W"] = 13.333 - 2 * g["_MARGIN"]
    if not design:
        return
    d = {k: v for k, v in design.items() if v not in (None, "")}

    def num(key, lo, hi):
        try:
            v = float(d[key])
        except (KeyError, TypeError, ValueError):
            return None
        return v if lo <= v <= hi else None

    if (v := num("size_title", 14, 40)) is not None:
        ratio = v / _DESIGN_DEFAULTS["_SZ_TITLE"]
        g["_SZ_TITLE"] = v
        g["_SZ_SUBTITLE"] = round(_DESIGN_DEFAULTS["_SZ_SUBTITLE"] * ratio)
        g["_SZ_HERO"] = round(_DESIGN_DEFAULTS["_SZ_HERO"] * ratio)
        g["_SZ_COVER"] = round(_DESIGN_DEFAULTS["_SZ_COVER"] * ratio)
    if (v := num("size_body", 9, 24)) is not None:
        g["_SZ_BODY"] = v
    if (v := num("size_small", 8, 18)) is not None:
        g["_SZ_SMALL"] = v
    if (v := num("line_spacing", 0.8, 2.0)) is not None:
        g["_LINE_SPACING"] = v
        g["_FORCE_LINE_SPACING"] = True   # reach native placeholders too, not just code-built text
    if (v := num("margin_in", 0.2, 1.5)) is not None:
        g["_MARGIN"] = v
        g["_CONTENT_W"] = 13.333 - 2 * v
    if (v := num("gutter_in", 0.1, 1.0)) is not None:
        g["_GUTTER"] = v
    tf = str(d.get("title_font") or "").strip()
    if tf:
        g["_HEAD_TITLE"] = tf     # slide titles (native and code-built)
        g["_HEAD"] = tf           # bold headings / eyebrows — the same "display" role
    bf = str(d.get("body_font") or "").strip()
    if bf:
        g["_BODY"] = bf                 # code-built text boxes read _BODY at call time
        g["_FORCE_BODY_FONT"] = bf      # native placeholders get forced in _set_text/_set_lines
    # Footer controls + icon switch (booleans/levels validated by the API; belt here again).
    if d.get("page_numbers") is False:
        g["_PAGE_NUMBERS"] = False
    ft = str(d.get("footer_text") or "").strip()
    if ft:
        g["_FOOTER_TEXT"] = ft[:80]
    if d.get("date_stamp") is True:
        g["_DATE_STAMP"] = True
    if d.get("icon_level") == "none":
        g["_ICONS_OFF"] = True


def _is_num(s: str) -> bool:
    s = (s or "").strip()
    return bool(s) and bool(re.search(r"\d", s)) and bool(re.fullmatch(r"[0-9.,%×xX+\-/()\s]+", s))


def _hbar_table(tbl) -> None:
    """Consulting-style table borders: horizontal row lines only, no vertical gridlines."""
    def _ln(tag, solid):
        inner = f'<a:solidFill><a:srgbClr val="{_TBL_LINE}"/></a:solidFill>' if solid else '<a:noFill/>'
        w = ' w="9525" cap="flat"' if solid else ''
        return parse_xml(f'<a:{tag} {nsdecls("a")}{w}>{inner}</a:{tag}>')
    for row in tbl.rows:
        for cell in row.cells:
            tcPr = cell._tc.get_or_add_tcPr()
            for tag in ("lnL", "lnR", "lnT", "lnB"):
                for el in tcPr.findall(qn("a:" + tag)):
                    tcPr.remove(el)
            for tag in ("lnB", "lnT", "lnR", "lnL"):   # insert at 0 -> final order lnL,lnR,lnT,lnB
                tcPr.insert(0, _ln(tag, tag == "lnB"))


def _place_text(slide, l, t, w, h, text, size, color, *, bold=False, font=None,
                align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False,
                line_spacing=None):
    # None defaults resolve at CALL time so apply_design()'s overrides reach them (a def-time
    # default would freeze the brand values at import).
    if font is None:
        font = _BODY
    if line_spacing is None:
        line_spacing = _LINE_SPACING
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    try:
        tf.auto_size = MSO_AUTO_SIZE.NONE      # fixed size — cap content, never shrink text to fit
    except Exception:  # noqa: BLE001
        pass
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    r = p.add_run()
    r.text = text or ""
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = font
    r.font.color.rgb = color
    return tb


def _synth_slide(prs, master_index, *, white=False, pastel=False, title=None, eyebrow=None):
    """Create a blank-layout slide with the ONE fixed skeleton every synthetic layout shares:
    chrome (footer logos / date / page number) is preserved so it sits in an identical position on
    every slide; the takeaway title and optional eyebrow are placed in their fixed boxes. Returns the
    slide with the body zone (_BODY_TOP.._BODY_BOTTOM) free for the builder to fill."""
    slide = prs.slides.add_slide(_blank_layout(prs, master_index))
    for sh in list(slide.shapes):
        if sh.is_placeholder and sh.placeholder_format.idx in CHROME_IDX:
            continue                            # keep date/footer/slide-number → cross-slide consistency
        sh._element.getparent().remove(sh._element)
    if white:
        _set_bg(slide, _LTEAL if pastel else _WHITE)
    if title is not None:
        # Non-bold "Exo 2 italic" — matches the native template's own title placeholders exactly
        # (its dominant cut; the brand guide lists no non-italic Exo 2 weight at all).
        _place_text(slide, _MARGIN, _TITLE_Y, _CONTENT_W, _TITLE_H, title, _SZ_TITLE,
                    _INKC if white else _WHITE, bold=False, italic=False, font=_HEAD_TITLE)
    if eyebrow is not None:
        _place_text(slide, _MARGIN, _EYEBROW_Y, _CONTENT_W, 0.4, eyebrow, _SZ_SMALL,
                    _TEAL if white else _LTEAL, bold=True, font=_HEAD)
    return slide


def _consistent_icons(objs):
    """All-or-nothing, one-source brand icons for a list of objects carrying icon/icon_generic.
    Returns a list of icon paths (one per object) or None if the set can't be cleanly covered — so a
    layout shows a full icon set or none, never a half-empty/mixed ring (the tell-tale AI look).
    Repeats ARE allowed (e.g. three sub-measures of ONE benefit sharing that benefit's icon) — the
    planner is asked to prefer distinct icons across genuinely different topics, but a slide whose
    items are all facets of one theme should still get its icon rather than none at all."""
    def _c(paths):
        return paths if paths and all(paths) else None
    return (_c([_icon_path(o.get("icon")) for o in objs])
            or _c([_generic_icon_path(o.get("icon_generic")) for o in objs]))


def _ink(light: bool) -> RGBColor:
    """Text drawn directly on the bare slide background: white on the dark master, near-black ink
    on the light one — the same swap `_synth_slide(white=True)` already does for the title/eyebrow,
    exposed here so every layout's OWN body text can follow suit when it supports both themes."""
    return _INKC if light else _WHITE


def _muted(light: bool) -> RGBColor:
    """Secondary/muted text directly on the bare background (captions, notes, axis labels)."""
    return _TEAL if light else _LTEAL


def _line_soft(light: bool, pastel: bool = False) -> RGBColor:
    """Connector lines, rules and arrow FILLS drawn directly on the bare background (not on a
    colour panel). `_LTEAL` is tuned for contrast against the dark master; on a light master it is
    nearly indistinguishable from the pale background, so use a darker teal instead - `_TEAL2`
    against white, or plain `_TEAL` (more contrast margin) against the pastel-mint background,
    which IS `_LTEAL` itself."""
    if not light:
        return _LTEAL
    return _TEAL if pastel else _TEAL2


def _chip_bg(light: bool, pastel: bool = False) -> RGBColor:
    """Icon-chip / tile fill meant to contrast against the SLIDE background (built for the dark
    master's near-black bg). On the white master the same near-white `_PANEL` tint would
    disappear, so use a soft teal tint there instead - but on the pastel-mint master the
    background IS that same teal tint (`_LTEAL`), so use plain white chips there instead."""
    if not light:
        return _PANEL
    return _WHITE if pastel else _LTEAL


def _icon_disc(slide, cx, cy, d, icon_path=None, number=None, light=False, pastel=False):
    """A soft light disc centred at (cx, cy) holding the red brand icon (or a red number) — the
    consulting 'icon chip' treatment that replaces ad-hoc accent bars on list slides."""
    disc = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - d / 2), Inches(cy - d / 2), Inches(d), Inches(d))
    disc.fill.solid(); disc.fill.fore_color.rgb = _chip_bg(light, pastel); disc.line.fill.background(); disc.shadow.inherit = False
    if icon_path:
        pad = d * 0.28
        _place_icon(slide, (Inches(cx - d / 2 + pad), Inches(cy - d / 2 + pad), Inches(d - 2 * pad), Inches(d - 2 * pad)), icon_path)
    elif number is not None:
        tf = disc.text_frame; tf.word_wrap = False; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.text = str(number)
        rr = tf.paragraphs[0].runs[0]
        rr.font.size = Pt(_SZ_BODY); rr.font.bold = True; rr.font.color.rgb = _RED; rr.font.name = _HEAD
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    return disc


def _place_bullets(slide, l, t, w, h, lines, size, color, *, font=None,
                   anchor=MSO_ANCHOR.TOP, rid=None):
    """Render lines as a Superba teal picture-bullet list in a synthetic textbox (the standard brand
    bullet). A single line still gets a bullet, so lists read consistently across the deck.
    font=None resolves to the current body font at call time (see apply_design)."""
    if font is None:
        font = _BODY
    lines = [ln.strip() for ln in lines if ln and ln.strip()]
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    try:
        tf.auto_size = MSO_AUTO_SIZE.NONE
    except Exception:  # noqa: BLE001
        pass
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Emu(0)
    if rid is None:
        rid = _bullet_rid(slide)
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = _LINE_SPACING
        p.space_after = Pt(6)
        r = p.add_run(); r.text = ln
        r.font.size = Pt(size); r.font.name = font; r.font.color.rgb = color
        if rid:
            _apply_picture_bullet(p._p, rid)
    return tb


def _fill_key_points(prs, spec: dict, light_index: int) -> None:
    """'Key points' cards: equal-height panels with a brand icon in a circle, a heading and a body,
    with the summary banner CLOSING the slide underneath them. Icons are all-or-nothing from ONE
    source (never a partial/empty set).

    The banner used to sit above the panels; it reads better as the conclusion the four cards add up
    to, and moving it down also reclaims the empty space the panels used to leave at the bottom."""
    pastel = spec.get("background") == "pastel"
    slide = _synth_slide(prs, light_index, white=True, pastel=pastel, title=spec.get("title", ""))
    banner = spec.get("banner")

    items = (spec.get("items") or [])[:4]
    n = len(items)
    if not n:
        return
    icons = _consistent_icons(items)              # all-or-nothing, one source (repeats allowed)
    d = _ICON_DISC
    ban_h = 0.55
    ptop = (_BODY_TOP + d / 2) if icons else _BODY_TOP
    pbot = _BODY_BOTTOM - (ban_h + 0.28 if banner else 0.0)
    if banner:
        by = pbot + 0.28
        ban = slide.shapes.add_shape(_BOX, Inches(_MARGIN), Inches(by), Inches(_CONTENT_W), Inches(ban_h))
        ban.fill.solid(); ban.fill.fore_color.rgb = _TEAL; ban.line.fill.background(); ban.shadow.inherit = False
        tf = ban.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER; p.line_spacing = _LINE_SPACING
        r = p.add_run(); r.text = banner; r.font.size = Pt(_SZ_BODY); r.font.bold = True
        r.font.name = _HEAD; r.font.color.rgb = _WHITE
    pw = (_CONTENT_W - (n - 1) * _GUTTER) / n     # equal panel widths, one standard gutter
    for i, it in enumerate(items):
        x = _MARGIN + i * (pw + _GUTTER); cx = x + pw / 2
        pan = slide.shapes.add_shape(_BOX, Inches(x), Inches(ptop), Inches(pw), Inches(pbot - ptop))
        pan.fill.solid(); pan.fill.fore_color.rgb = _WHITE if pastel else _PANEL
        pan.line.fill.background(); pan.shadow.inherit = False
        hy = ptop + (0.6 if icons else 0.22)
        if icons:
            circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - d / 2), Inches(ptop - d / 2), Inches(d), Inches(d))
            circ.fill.solid(); circ.fill.fore_color.rgb = _WHITE
            circ.line.color.rgb = _RED; circ.line.width = Pt(2.25); circ.shadow.inherit = False
            _place_icon(slide, (Inches(cx - 0.26), Inches(ptop - 0.26), Inches(0.52), Inches(0.52)), icons[i])
        _place_text(slide, x + _PAD, hy, pw - 2 * _PAD, 0.5, it.get("heading", ""), _SZ_BODY, _INKC,
                    bold=True, font=_HEAD, align=PP_ALIGN.CENTER)
        # Body as standard Superba teal bullets (one short point per line).
        _place_bullets(slide, x + _PAD, hy + 0.6, pw - 2 * _PAD, pbot - (hy + 0.6) - _PAD,
                       str(it.get("body", "")).split("\n"), _SZ_SMALL, _INKC)


_CHART_TYPES = {"column": XL_CHART_TYPE.COLUMN_CLUSTERED, "bar": XL_CHART_TYPE.BAR_CLUSTERED,
                "line": XL_CHART_TYPE.LINE, "stacked_column": XL_CHART_TYPE.COLUMN_STACKED,
                "stacked_100": XL_CHART_TYPE.COLUMN_STACKED_100, "doughnut": XL_CHART_TYPE.DOUGHNUT}


def _fill_chart(prs, spec: dict, dark_index: int, light_index: int | None = None) -> None:
    """Native, editable PowerPoint chart from the plan's categories + series, brand-coloured, on the
    deep-sea master (inherits background + logos). Data comes only from the plan (claim fidelity)."""
    light = light_index is not None and spec.get("background") in ("light", "pastel")
    pastel = light and spec.get("background") == "pastel"
    slide = _synth_slide(prs, light_index if light else dark_index, white=light, pastel=pastel, title=spec.get("title", ""),
                         eyebrow=spec.get("caption"))

    cats = spec.get("categories") or []
    series = spec.get("series") or []
    if not cats or not series:
        return
    cd = CategoryChartData()
    cd.categories = cats
    for s in series:
        vals = [(v if isinstance(v, (int, float)) else None) for v in (s.get("values") or [])]
        cd.add_series(s.get("name", ""), vals)
    ctype = _CHART_TYPES.get(spec.get("chart_type", "column"), XL_CHART_TYPE.COLUMN_CLUSTERED)
    gf = slide.shapes.add_chart(ctype, Inches(0.9), Inches(_BODY_TOP), Inches(11.5), Inches(_BODY_H - 0.35), cd)
    chart = gf.chart
    chart.has_title = False                # no auto series-name title (we use the slide title)
    chart.font.color.rgb = _ink(light)
    chart.font.size = Pt(_SZ_SMALL)
    chart.font.name = _BODY
    is_round = spec.get("chart_type") == "doughnut"
    multi = len(series) > 1 or is_round
    chart.has_legend = multi
    if multi:
        from pptx.enum.chart import XL_LEGEND_POSITION
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
    is_bar = spec.get("chart_type", "column") in ("column", "bar")
    if is_round or (is_bar and len(series) == 1):
        # One series, many bars/wedges sharing a single colour reads as a flat, low-contrast block
        # (this is also what a theme-driven "vary colours" default can silently collapse to on some
        # renderers even when a series-level fill IS set) — colour each POINT explicitly instead, the
        # most specific level of formatting, so every bar/wedge gets its own unambiguous brand colour.
        for i, pt in enumerate(chart.plots[0].series[0].points):
            try:
                pt.format.fill.solid(); pt.format.fill.fore_color.rgb = _CHART_COLORS[i % len(_CHART_COLORS)]
            except Exception:  # noqa: BLE001
                pass
    else:
        for i, plot_series in enumerate(chart.plots[0].series):
            try:
                plot_series.format.fill.solid()
                plot_series.format.fill.fore_color.rgb = _CHART_COLORS[i % len(_CHART_COLORS)]
            except Exception:  # noqa: BLE001 — line charts style the line
                plot_series.format.line.color.rgb = _CHART_COLORS[i % len(_CHART_COLORS)]
    try:
        chart.plots[0].vary_by_categories = is_round or (is_bar and len(series) == 1)
    except Exception:  # noqa: BLE001 — not exposed on every chart type
        pass

    # Gridlines/axis chrome default to whatever the renderer's own theme resolves, which on some
    # renderers comes out near-white — indistinguishable from a white/light bar and from the white
    # axis text. Style them explicitly: a muted line for the grid (present but recessive), a bright
    # one for the axis itself, so brand-coloured bars are what actually draws the eye.
    if not is_round:
        try:
            chart.value_axis.has_major_gridlines = True
            gl = chart.value_axis.major_gridlines.format.line
            gl.color.rgb = _TEAL2
            gl.width = Pt(0.5)
        except Exception:  # noqa: BLE001
            pass
        for axis in (chart.category_axis, chart.value_axis):
            try:
                axis.format.line.color.rgb = _ink(light)
                axis.tick_labels.font.color.rgb = _ink(light)
                axis.tick_labels.font.size = Pt(_SZ_SMALL)
                axis.tick_labels.font.name = _BODY
            except Exception:  # noqa: BLE001
                pass

    # Axis titles — mandatory for charts with axes (doughnut has none). The category axis takes the
    # dimension (x_axis), the value axis takes what is measured + units (y_axis); on a bar chart these
    # are visually swapped but stay semantically correct.
    if not is_round:
        def _axis_title(axis, text):
            axis.has_title = True
            tf = axis.axis_title.text_frame
            tf.text = text
            run = tf.paragraphs[0].runs[0]
            run.font.color.rgb = _ink(light)
            run.font.size = Pt(_SZ_SMALL)
            run.font.name = _BODY
            run.font.bold = True
        try:
            if spec.get("x_axis"):
                _axis_title(chart.category_axis, spec["x_axis"])
            if spec.get("y_axis"):
                _axis_title(chart.value_axis, spec["y_axis"])
        except (ValueError, KeyError, IndexError):  # axis absent for this chart type
            pass

    # Line charts: pull the plot to the axis edges (first category flush left) instead of the default
    # half-category padding, so the first point (e.g. "Day 0") sits at the left edge.
    if spec.get("chart_type") == "line":
        try:
            valAx = chart.value_axis._element
            existing = valAx.find(qn("c:crossBetween"))
            if existing is not None:
                existing.set("val", "midCat")
            else:
                anchor = None
                for tag in ("c:crossAx", "c:crosses", "c:crossesAt"):
                    el = valAx.find(qn(tag))
                    if el is not None:
                        anchor = el
                cb = valAx.makeelement(qn("c:crossBetween"), {"val": "midCat"})
                anchor.addnext(cb) if anchor is not None else valAx.append(cb)
        except Exception:  # noqa: BLE001
            pass

    # Think-cell-style delta callout: for a 2-bar single-series column chart, reserve headroom on the
    # value axis, then draw a bracket spanning the two columns with a red delta chip (the % change) —
    # the classic "highlight the difference" annotation.
    if spec.get("chart_type", "column") == "column" and len(cats) == 2 and len(series) == 1:
        vals = [v for v in (series[0].get("values") or []) if isinstance(v, (int, float))]
        if len(vals) == 2 and (vals[0] or vals[1]):
            v0, v1 = vals[0], vals[1]
            label = (f"{'+' if v1 >= v0 else ''}{(v1 - v0) / abs(v0) * 100:.0f}%") if v0 else f"+{v1 - v0:g}"
            arrow = "▲" if v1 >= v0 else "▼"
            try:
                chart.value_axis.minimum_scale = 0
                chart.value_axis.maximum_scale = max(vals) * 1.35     # guaranteed headroom for the callout
            except Exception:  # noqa: BLE001
                pass
            fx, fy, fw, fh = 0.9, _BODY_TOP, 11.5, _BODY_H - 0.35
            plot_l, plot_w = fx + 0.75, fw - 0.95
            plot_top, plot_bot = fy + 0.15, fy + fh - 0.55
            cx0 = plot_l + plot_w * 0.25
            cx1 = plot_l + plot_w * 0.75
            by = max(plot_top + 0.2, plot_bot - (plot_bot - plot_top) / 1.35 - 0.3)   # just above the tall bar
            ln = slide.shapes.add_shape(_BOX, Inches(cx0), Inches(by), Inches(cx1 - cx0), Inches(0.035))
            ln.fill.solid(); ln.fill.fore_color.rgb = _RED; ln.line.fill.background(); ln.shadow.inherit = False
            for cxe in (cx0, cx1):
                tick = slide.shapes.add_shape(_BOX, Inches(cxe - 0.017), Inches(by), Inches(0.035), Inches(0.16))
                tick.fill.solid(); tick.fill.fore_color.rgb = _RED; tick.line.fill.background(); tick.shadow.inherit = False
            cw, ch = 1.5, 0.44
            chip = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches((cx0 + cx1) / 2 - cw / 2),
                                          Inches(by - ch - 0.08), Inches(cw), Inches(ch))
            chip.fill.solid(); chip.fill.fore_color.rgb = _RED; chip.line.fill.background(); chip.shadow.inherit = False
            ctf = chip.text_frame; ctf.word_wrap = False; ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
            ctf.margin_top = ctf.margin_bottom = Emu(0)
            cp = ctf.paragraphs[0]; cp.alignment = PP_ALIGN.CENTER
            cr = cp.add_run(); cr.text = f"{arrow} {label}"; cr.font.size = Pt(_SZ_BODY); cr.font.bold = True
            cr.font.name = _HEAD; cr.font.color.rgb = _WHITE


def _fill_matrix(prs, spec: dict, dark_index: int, light_index: int | None = None) -> None:
    """2x2 matrix: four equal teal quadrants separated by the standard gutter, with axis labels."""
    light = light_index is not None and spec.get("background") in ("light", "pastel")
    pastel = light and spec.get("background") == "pastel"
    slide = _synth_slide(prs, light_index if light else dark_index, white=light, pastel=pastel, title=spec.get("title", ""))
    quads = (spec.get("quadrants") or [])[:4]
    mx, my = 3.0, _BODY_TOP
    mw = 13.333 - _MARGIN - mx
    mh = _BODY_BOTTOM - my - 0.45          # leave room for the x-axis label below
    gap = _GUTTER
    qw, qh = (mw - gap) / 2, (mh - gap) / 2
    pos = [(mx, my), (mx + qw + gap, my), (mx, my + qh + gap), (mx + qw + gap, my + qh + gap)]
    for i, q in enumerate(quads):
        x, y = pos[i]
        pan = slide.shapes.add_shape(_BOX, Inches(x), Inches(y), Inches(qw), Inches(qh))
        pan.fill.solid(); pan.fill.fore_color.rgb = _TEAL; pan.line.fill.background(); pan.shadow.inherit = False
        _place_text(slide, x + _PAD, y + 0.16, qw - 2 * _PAD, 0.45, q.get("heading", ""), _SZ_BODY, _WHITE, bold=True, font=_HEAD)
        _place_text(slide, x + _PAD, y + 0.66, qw - 2 * _PAD, qh - 0.82, q.get("body", ""), _SZ_BODY, _ONTEAL)
    if spec.get("y_axis"):
        _place_text(slide, _MARGIN, my, mx - _MARGIN - 0.15, mh, spec["y_axis"], _SZ_SMALL, _muted(light),
                    bold=True, font=_HEAD, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
    if spec.get("x_axis"):
        _place_text(slide, mx, my + mh + 0.08, mw, 0.35, spec["x_axis"], _SZ_SMALL, _muted(light),
                    bold=True, font=_HEAD, align=PP_ALIGN.CENTER)


_EXEC_SUMMARY_ROWS = (("source", "Source"), ("key_finding", "Key finding"),
                     ("supporting_findings", "Supporting findings"), ("relevance", "Relevance"),
                     ("contents", "Contents"))


def _place_labeled_bullets(slide, l, t, w, h, rows, size, color, *, font=None, anchor=MSO_ANCHOR.TOP, light=False):
    """Render (label, text) rows as the deck's STANDARD teal picture-bullet list — each bullet a
    bold lead-in label followed inline by its sentence(s). One flowing text box (not fixed
    per-row height slots), so a longer row just wraps to another line instead of colliding with
    the next bullet — the failure mode a fixed-slot layout risks when row lengths vary."""
    if font is None:
        font = _BODY
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    try:
        tf.auto_size = MSO_AUTO_SIZE.NONE
    except Exception:  # noqa: BLE001
        pass
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Emu(0)
    rid = _bullet_rid(slide)
    for i, (label, text) in enumerate(rows):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = _LINE_SPACING
        p.space_after = Pt(10)
        r1 = p.add_run(); r1.text = f"{label}: "
        r1.font.size = Pt(size); r1.font.bold = True; r1.font.name = _HEAD; r1.font.color.rgb = _ink(light)
        r2 = p.add_run(); r2.text = text or ""
        r2.font.size = Pt(size); r2.font.bold = False; r2.font.name = font; r2.font.color.rgb = color
        if rid:
            _apply_picture_bullet(p._p, rid)
    return tb


def _fill_exec_summary(prs, spec: dict, dark_index: int, light_index: int | None = None) -> None:
    """The deck's REQUIRED executive summary (slide 2, right after the cover): a FIXED static
    title (never the model's — matches ingredient/agenda's forced-text pattern) and exactly 5
    labelled rows (Source / Key finding / Supporting findings / Relevance / Contents), each a
    bold lead-in label plus 1-2 sentences, rendered as the deck's own standard bulleted list — a
    clean, scannable summary rather than a stacked memo (client feedback on the first cut)."""
    light = light_index is not None and spec.get("background") in ("light", "pastel")
    pastel = light and spec.get("background") == "pastel"
    slide = _synth_slide(prs, light_index if light else dark_index, white=light, pastel=pastel, title="Executive summary")
    rows = [(label, spec.get(key, "")) for key, label in _EXEC_SUMMARY_ROWS if spec.get(key)]
    if not rows:
        return
    _place_labeled_bullets(slide, _MARGIN, _BODY_TOP + 0.15, _CONTENT_W, _BODY_H - 0.3,
                           rows, _SZ_BODY, _muted(light), light=light)


def _fill_comparison(prs, spec: dict, light_index: int) -> None:
    """Comparison table: a native, brand-styled table (teal header, light body) on white."""
    pastel = spec.get("background") == "pastel"
    slide = _synth_slide(prs, light_index, white=True, pastel=pastel, title=spec.get("title", ""))
    headers = spec.get("headers") or []
    rows = spec.get("rows") or []
    ncols = len(headers)
    nrows = len(rows) + 1
    if ncols < 1 or nrows < 2:
        return
    h = min(_BODY_H, 0.5 + 0.62 * (nrows - 1))
    gf = slide.shapes.add_table(nrows, ncols, Inches(_MARGIN), Inches(_BODY_TOP), Inches(_CONTENT_W), Inches(h))
    tbl = gf.table
    tbl.first_row = False; tbl.horz_banding = False
    def _cell(cell, text, *, bold, color, fill, align=PP_ALIGN.LEFT):
        cell.fill.solid(); cell.fill.fore_color.rgb = fill
        cell.margin_left = cell.margin_right = Inches(0.12)
        tf = cell.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = align
        r = p.add_run(); r.text = text or ""
        r.font.size = Pt(_SZ_BODY); r.font.bold = bold; r.font.name = (_HEAD if bold else _BODY); r.font.color.rgb = color
    for j, head in enumerate(headers):
        _cell(tbl.cell(0, j), head, bold=True, color=_WHITE, fill=_TEAL)
    for i, row in enumerate(rows, start=1):
        cells = (row.get("cells") or [])
        band = RGBColor(0xF1, 0xF8, 0xF8) if i % 2 else _WHITE
        for j in range(ncols):
            val = cells[j] if j < len(cells) else ""
            al = PP_ALIGN.RIGHT if (j > 0 and _is_num(val)) else PP_ALIGN.LEFT   # numbers right, text left
            _cell(tbl.cell(i, j), val, bold=(j == 0), color=_INKC, fill=band, align=al)
    _hbar_table(tbl)   # horizontal row lines only


def _fill_stat(prs, spec: dict, dark_index: int, light_index: int | None = None) -> None:
    """Hero stats: 1-3 big red figures with labels (the '50+ / 135+' treatment)."""
    light = light_index is not None and spec.get("background") in ("light", "pastel")
    pastel = light and spec.get("background") == "pastel"
    slide = _synth_slide(prs, light_index if light else dark_index, white=light, pastel=pastel, title=spec.get("title", ""), eyebrow=spec.get("caption"))
    stats = (spec.get("stats") or [])[:3]
    n = max(1, len(stats))
    cw = (_CONTENT_W - (n - 1) * _GUTTER) / n     # equal columns, one standard gutter
    vy = _BODY_TOP + 0.8
    for i, st in enumerate(stats):
        x = _MARGIN + i * (cw + _GUTTER)
        _place_text(slide, x, vy, cw, 1.0, st.get("value", ""), _SZ_HERO, _RED, bold=True, font=_HEAD, align=PP_ALIGN.CENTER)
        _place_text(slide, x, vy + 1.05, cw, 0.5, st.get("label", ""), _SZ_BODY, _ink(light), bold=True, font=_HEAD, align=PP_ALIGN.CENTER)
        if st.get("note"):
            _place_text(slide, x + 0.2, vy + 1.6, cw - 0.4, 1.4, st["note"], _SZ_SMALL, _muted(light), align=PP_ALIGN.CENTER)


def _set_pie_angles(shape, a1_deg, a2_deg):
    """Set a PIE shape's start/end angles (degrees, clockwise from 3 o'clock) via its geometry guides."""
    geom = shape._element.spPr.find(qn("a:prstGeom"))
    av = geom.find(qn("a:avLst"))
    if av is None:
        av = geom.makeelement(qn("a:avLst"), {}); geom.append(av)
    for el in list(av):
        av.remove(el)
    for name, val in (("adj1", a1_deg), ("adj2", a2_deg)):
        av.append(av.makeelement(qn("a:gd"), {"name": name, "fmla": f"val {int(round(val * 60000))}"}))


def _harvey_ball(slide, cx, cy, d, score):
    """A fixed-diameter harvey ball centred at (cx, cy): a red ring always, plus a red pie wedge for the
    filled fraction (score 0..4). Drawn as real shapes so every ball is exactly the same size."""
    ring = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - d / 2), Inches(cy - d / 2), Inches(d), Inches(d))
    ring.fill.background(); ring.line.color.rgb = _RED; ring.line.width = Pt(1.5); ring.shadow.inherit = False
    f = max(0, min(4, int(round(score)))) / 4.0
    if f <= 0:
        return
    if f >= 1:
        disc = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - d / 2), Inches(cy - d / 2), Inches(d), Inches(d))
        disc.fill.solid(); disc.fill.fore_color.rgb = _RED
        disc.line.color.rgb = _RED; disc.line.width = Pt(1.5); disc.shadow.inherit = False
        return
    pie = slide.shapes.add_shape(MSO_SHAPE.PIE, Inches(cx - d / 2), Inches(cy - d / 2), Inches(d), Inches(d))
    pie.fill.solid(); pie.fill.fore_color.rgb = _RED; pie.line.fill.background(); pie.shadow.inherit = False
    _set_pie_angles(pie, 270.0, (270.0 + 360.0 * f) % 360.0)   # fill from 12 o'clock, clockwise


def _fill_harvey_ball(prs, spec: dict, light_index: int) -> None:
    """Harvey-ball rating grid: criteria (rows) x options (columns), each cell a 0-4 filled circle."""
    pastel = spec.get("background") == "pastel"
    slide = _synth_slide(prs, light_index, white=True, pastel=pastel, title=spec.get("title", ""))
    options = spec.get("options") or []
    criteria = spec.get("criteria") or []
    ncols = len(options) + 1
    nrows = len(criteria) + 1
    if ncols < 2 or nrows < 2:
        return
    nopt = len(options)
    label_w = 3.2
    opt_w = (_CONTENT_W - label_w) / nopt
    hr0 = 0.6                                        # header row height
    hr = min(1.0, (_BODY_H - hr0) / (nrows - 1))     # data row height
    h = hr0 + (nrows - 1) * hr
    tbl = slide.shapes.add_table(nrows, ncols, Inches(_MARGIN), Inches(_BODY_TOP), Inches(_CONTENT_W), Inches(h)).table
    tbl.first_row = False; tbl.horz_banding = False
    tbl.columns[0].width = Inches(label_w)           # explicit widths so ball centres are known
    for j in range(1, ncols):
        tbl.columns[j].width = Inches(opt_w)
    tbl.rows[0].height = Inches(hr0)
    for i in range(1, nrows):
        tbl.rows[i].height = Inches(hr)

    def _cell(cell, text, *, bold, color, fill, center=False, font=_BODY):
        cell.fill.solid(); cell.fill.fore_color.rgb = fill
        cell.margin_left = cell.margin_right = Inches(0.12)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf = cell.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]
        if center:
            p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = text or ""
        r.font.size = Pt(_SZ_BODY); r.font.bold = bold; r.font.name = font; r.font.color.rgb = color
    _cell(tbl.cell(0, 0), "", bold=True, color=_WHITE, fill=_TEAL)
    for j, opt in enumerate(options, start=1):
        _cell(tbl.cell(0, j), opt, bold=True, color=_WHITE, fill=_TEAL, center=True, font=_HEAD)
    for i, crit in enumerate(criteria, start=1):
        band = RGBColor(0xF1, 0xF8, 0xF8) if i % 2 else _WHITE
        _cell(tbl.cell(i, 0), crit.get("label", ""), bold=True, color=_INKC, fill=band)
        for j in range(1, ncols):
            _cell(tbl.cell(i, j), "", bold=False, color=_INKC, fill=band)   # empty — ball is drawn on top
    _hbar_table(tbl)   # horizontal row lines only

    # Overlay fixed-size harvey balls at each score cell's centre (so every ball is identical).
    d = min(0.5, hr * 0.55)
    for i, crit in enumerate(criteria):
        scores = crit.get("scores") or []
        cy = _BODY_TOP + hr0 + i * hr + hr / 2
        for j in range(nopt):
            sc = scores[j] if j < len(scores) else 0
            cx = _MARGIN + label_w + j * opt_w + opt_w / 2
            _harvey_ball(slide, cx, cy, d, sc)


def _fill_funnel(prs, spec: dict, dark_index: int, light_index: int | None = None) -> None:
    """Funnel: centred bars that narrow top-to-bottom, one per stage, heading + body."""
    light = light_index is not None and spec.get("background") in ("light", "pastel")
    pastel = light and spec.get("background") == "pastel"
    slide = _synth_slide(prs, light_index if light else dark_index, white=light, pastel=pastel, title=spec.get("title", ""))
    stages = (spec.get("stages") or [])[:5]
    n = len(stages)
    if not n:
        return
    wide, narrow = 9.0, 4.5
    bh = (_BODY_H - (n - 1) * _GUTTER) / n          # bars fill the body zone, one standard gutter
    for i, st in enumerate(stages):
        w = wide - (wide - narrow) * (i / max(1, n - 1))
        x = (13.333 - w) / 2
        y = _BODY_TOP + i * (bh + _GUTTER)
        bar = slide.shapes.add_shape(_BOX, Inches(x), Inches(y), Inches(w), Inches(bh))
        bar.fill.solid(); bar.fill.fore_color.rgb = _TEAL if i % 2 == 0 else _TEAL2
        bar.line.fill.background(); bar.shadow.inherit = False
        tf = bar.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER; p.line_spacing = _LINE_SPACING
        r = p.add_run(); r.text = st.get("heading", ""); r.font.size = Pt(_SZ_BODY); r.font.bold = True
        r.font.name = _HEAD; r.font.color.rgb = _WHITE
        if st.get("body"):
            p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER; p2.line_spacing = _LINE_SPACING
            r2 = p2.add_run(); r2.text = st["body"]; r2.font.size = Pt(_SZ_SMALL); r2.font.name = _BODY
            r2.font.color.rgb = RGBColor(0xE9, 0xF7, 0xF8)


def _fill_closing(prs, spec: dict, dark_index: int, light_index: int | None = None) -> None:
    """Closing / contact: a closing statement, optional tagline, and contact details."""
    light = light_index is not None and spec.get("background") in ("light", "pastel")
    pastel = light and spec.get("background") == "pastel"
    slide = _synth_slide(prs, light_index if light else dark_index, white=light, pastel=pastel)
    cy = 2.6
    _place_text(slide, _MARGIN, cy, _CONTENT_W, 1.6, spec.get("title", ""), _SZ_TITLE, _ink(light),
                bold=False, italic=False, font=_HEAD_TITLE)
    if spec.get("tagline"):
        _place_text(slide, _MARGIN, cy + 1.6, _CONTENT_W, 0.8, spec["tagline"], _SZ_BODY, _muted(light))
    if spec.get("contact"):
        _place_text(slide, _MARGIN, _BODY_BOTTOM - 0.5, _CONTENT_W, 0.5, spec["contact"], _SZ_SMALL, _muted(light), bold=True, font=_HEAD)


# Same-family teal tints (apex/first = darkest, base/last = lightest) — used by layered/sequenced layouts.
_TEAL_TINTS = [_TEAL, RGBColor(0x24, 0x6C, 0x79), RGBColor(0x36, 0x86, 0x90),
               RGBColor(0x50, 0xA0, 0xA7), RGBColor(0x72, 0xB8, 0xBB)]


def _fill_kpi_dashboard(prs, spec: dict, dark_index: int, light_index: int | None = None) -> None:
    """KPI dashboard: a grid of metric tiles (a hero figure + label + optional note) — the MBB scoreboard."""
    light = light_index is not None and spec.get("background") in ("light", "pastel")
    pastel = light and spec.get("background") == "pastel"
    slide = _synth_slide(prs, light_index if light else dark_index, white=light, pastel=pastel, title=spec.get("title", ""), eyebrow=spec.get("caption"))
    metrics = (spec.get("metrics") or [])[:6]
    n = len(metrics)
    if not n:
        return
    cols = n if n <= 3 else 3
    rows = (n + cols - 1) // cols
    tw = (_CONTENT_W - (cols - 1) * _GUTTER) / cols
    th = (_BODY_H - (rows - 1) * _GUTTER) / rows
    for i, m in enumerate(metrics):
        rr, cc = divmod(i, cols)
        x = _MARGIN + cc * (tw + _GUTTER)
        y = _BODY_TOP + rr * (th + _GUTTER)
        tile = slide.shapes.add_shape(_BOX, Inches(x), Inches(y), Inches(tw), Inches(th))
        tile.fill.solid(); tile.fill.fore_color.rgb = _chip_bg(light, pastel); tile.line.fill.background(); tile.shadow.inherit = False
        _place_text(slide, x + _PAD, y, tw - 2 * _PAD, th * 0.5, m.get("value", ""), _SZ_HERO, _RED,
                    bold=True, font=_HEAD, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.BOTTOM)
        _place_text(slide, x + _PAD, y + th * 0.52, tw - 2 * _PAD, 0.5, m.get("label", ""), _SZ_BODY, _INKC,
                    bold=True, font=_HEAD, align=PP_ALIGN.CENTER)
        if m.get("note"):
            _place_text(slide, x + _PAD, y + th * 0.52 + 0.5, tw - 2 * _PAD, 0.5, m["note"], _SZ_SMALL, _INKC,
                        align=PP_ALIGN.CENTER)


def _fill_roadmap(prs, spec: dict, dark_index: int, light_index: int | None = None) -> None:
    """Roadmap: left-to-right interlocking chevron phases, each with a date above, a heading inside, and
    activities below. Same-family teal tints so it reads as one sequence."""
    light = light_index is not None and spec.get("background") in ("light", "pastel")
    pastel = light and spec.get("background") == "pastel"
    slide = _synth_slide(prs, light_index if light else dark_index, white=light, pastel=pastel, title=spec.get("title", ""))
    phases = (spec.get("phases") or [])[:5]
    n = len(phases)
    if not n:
        return
    overlap = 0.35
    cw = (_CONTENT_W + (n - 1) * overlap) / n       # chevrons overlap so the arrows interlock
    ch = 0.95
    cy = _BODY_TOP + 0.45
    for i, ph in enumerate(phases):
        x = _MARGIN + i * (cw - overlap)
        if ph.get("date"):
            _place_text(slide, x, cy - 0.4, cw - overlap, 0.35, ph["date"], _SZ_SMALL, _muted(light),
                        bold=True, font=_HEAD, align=PP_ALIGN.CENTER)
        chev = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(x), Inches(cy), Inches(cw), Inches(ch))
        chev.fill.solid(); chev.fill.fore_color.rgb = _TEAL_TINTS[i % len(_TEAL_TINTS)]
        chev.line.fill.background(); chev.shadow.inherit = False
        tf = chev.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.3)                # the chevron point eats space on the right
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER; p.line_spacing = _LINE_SPACING
        r = p.add_run(); r.text = ph.get("heading", ""); r.font.size = Pt(_SZ_BODY); r.font.bold = True
        r.font.name = _HEAD; r.font.color.rgb = _WHITE
        if ph.get("body"):
            _place_text(slide, x + 0.3, cy + ch + 0.3, cw - overlap - 0.3, _BODY_BOTTOM - (cy + ch + 0.3),
                        ph["body"], _SZ_SMALL, _muted(light))


def _fill_icon_grid(prs, spec: dict, dark_index: int, light_index: int | None = None) -> None:
    """A grid of 3 to 6 icon tiles (icon chip + heading + body), 3 across. Optional teal banner. Icons
    are all-or-nothing from one source. Covers the client's icon-card slides."""
    light = light_index is not None and spec.get("background") in ("light", "pastel")
    pastel = light and spec.get("background") == "pastel"
    slide = _synth_slide(prs, light_index if light else dark_index, white=light, pastel=pastel, title=spec.get("title", ""))
    banner = spec.get("banner")
    top = _BODY_TOP
    if banner:
        ban = slide.shapes.add_shape(_BOX, Inches(_MARGIN), Inches(_EYEBROW_Y), Inches(_CONTENT_W), Inches(0.55))
        ban.fill.solid(); ban.fill.fore_color.rgb = _TEAL; ban.line.fill.background(); ban.shadow.inherit = False
        tf = ban.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER; p.line_spacing = _LINE_SPACING
        r = p.add_run(); r.text = banner; r.font.size = Pt(_SZ_BODY); r.font.bold = True
        r.font.name = _HEAD; r.font.color.rgb = _WHITE
        top = _EYEBROW_Y + 0.55 + 0.35
    items = (spec.get("items") or [])[:6]
    n = len(items)
    if not n:
        return
    icons = _consistent_icons(items)
    cols = 2 if n == 4 else (n if n <= 3 else 3)
    rows = (n + cols - 1) // cols
    cw = (_CONTENT_W - (cols - 1) * _GUTTER) / cols
    ch = (_BODY_BOTTOM - top - (rows - 1) * _GUTTER) / rows
    d = 0.82
    for i, it in enumerate(items):
        rr, cc = divmod(i, cols)
        x = _MARGIN + cc * (cw + _GUTTER)
        y = top + rr * (ch + _GUTTER)
        cx = x + cw / 2
        if icons:
            _icon_disc(slide, cx, y + d / 2 + 0.08, d, icon_path=icons[i], light=light, pastel=pastel)
            hy = y + d + 0.28
        else:
            hy = y + 0.25
        _place_text(slide, x + _PAD, hy, cw - 2 * _PAD, 0.5, it.get("heading", ""), _SZ_BODY, _ink(light),
                    bold=True, font=_HEAD, align=PP_ALIGN.CENTER)
        _place_text(slide, x + _PAD, hy + 0.5, cw - 2 * _PAD, y + ch - (hy + 0.5) - 0.05,
                    it.get("body", ""), _SZ_SMALL, _muted(light), align=PP_ALIGN.CENTER)


def _fill_takeaways(prs, spec: dict, dark_index: int, light_index: int | None = None) -> None:
    """Numbered 'key messages' rows: a red-number chip + a bold statement + supporting detail, one per
    row with a thin divider between. Covers the client's summary / takeaways slides."""
    light = light_index is not None and spec.get("background") in ("light", "pastel")
    pastel = light and spec.get("background") == "pastel"
    slide = _synth_slide(prs, light_index if light else dark_index, white=light, pastel=pastel, title=spec.get("title", ""))
    items = (spec.get("items") or [])[:6]
    n = len(items)
    if not n:
        return
    rh = _BODY_H / n
    d = 0.5
    for i, it in enumerate(items):
        y = _BODY_TOP + i * rh
        _icon_disc(slide, _MARGIN + d / 2, y + rh / 2 - 0.05, d, number=i + 1, light=light, pastel=pastel)
        tx = _MARGIN + d + 0.35
        tw = _CONTENT_W - (d + 0.35)
        if it.get("body"):
            _place_text(slide, tx, y + 0.1, tw, 0.5, it.get("heading", ""), _SZ_BODY, _ink(light), bold=True, font=_HEAD)
            _place_text(slide, tx, y + 0.6, tw, rh - 0.72, it["body"], _SZ_SMALL, _muted(light))
        else:
            _place_text(slide, tx, y, tw, rh - 0.1, it.get("heading", ""), _SZ_BODY, _ink(light), bold=True,
                        font=_HEAD, anchor=MSO_ANCHOR.MIDDLE)
        if i < n - 1:
            ln = slide.shapes.add_shape(_BOX, Inches(_MARGIN), Inches(y + rh - 0.02), Inches(_CONTENT_W), Inches(0.012))
            ln.fill.solid(); ln.fill.fore_color.rgb = _TEAL2; ln.line.fill.background(); ln.shadow.inherit = False


def _fill_from_to(prs, spec: dict, dark_index: int, light_index: int | None = None) -> None:
    """Transformation: a FROM panel, a teal arrow, and a TO panel — each with a small eyebrow, a heading
    and a body. Covers the client's from/to and before/after slides (no red block)."""
    light = light_index is not None and spec.get("background") in ("light", "pastel")
    pastel = light and spec.get("background") == "pastel"
    slide = _synth_slide(prs, light_index if light else dark_index, white=light, pastel=pastel, title=spec.get("title", ""))
    before = spec.get("before") or {}
    after = spec.get("after") or {}
    py = _BODY_TOP + 0.4
    ph = _BODY_H - 0.8
    arrow_w = 1.1
    pw = (_CONTENT_W - arrow_w - 2 * _GUTTER) / 2

    def _panel(x, fill, eyebrow, obj):
        pan = slide.shapes.add_shape(_BOX, Inches(x), Inches(py), Inches(pw), Inches(ph))
        pan.fill.solid(); pan.fill.fore_color.rgb = fill; pan.line.fill.background(); pan.shadow.inherit = False
        _place_text(slide, x + _PAD, py + _PAD, pw - 2 * _PAD, 0.35, eyebrow, _SZ_SMALL, _LTEAL, bold=True, font=_HEAD)
        _place_text(slide, x + _PAD, py + 0.7, pw - 2 * _PAD, 0.9, obj.get("heading", ""), _SZ_BODY, _WHITE, bold=True, font=_HEAD)
        if obj.get("body"):
            _place_text(slide, x + _PAD, py + 1.5, pw - 2 * _PAD, ph - 1.7, obj["body"], _SZ_SMALL, _ONTEAL)

    _panel(_MARGIN, RGBColor(0x22, 0x55, 0x5E), "FROM", before)
    ax = _MARGIN + pw + _GUTTER
    arr = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(ax), Inches(py + ph / 2 - 0.55), Inches(arrow_w), Inches(1.1))
    arr.fill.solid(); arr.fill.fore_color.rgb = _line_soft(light, pastel); arr.line.fill.background(); arr.shadow.inherit = False
    _panel(ax + arrow_w + _GUTTER, _TEAL2, "TO", after)


def _rule(slide, x, y, w, h, color):
    """A thin filled rectangle used as a connector line."""
    ln = slide.shapes.add_shape(_BOX, Inches(x), Inches(y), Inches(max(w, 0.02)), Inches(max(h, 0.02)))
    ln.fill.solid(); ln.fill.fore_color.rgb = color; ln.line.fill.background(); ln.shadow.inherit = False
    return ln


def _fill_pillars(prs, spec: dict, dark_index: int, light_index: int | None = None) -> None:
    """Pillars under a roof: an optional roof banner across the top, then 2 to 5 tall pillars, each with an
    icon chip, a heading and a body. For 'our approach rests on N pillars' framing."""
    light = light_index is not None and spec.get("background") in ("light", "pastel")
    pastel = light and spec.get("background") == "pastel"
    slide = _synth_slide(prs, light_index if light else dark_index, white=light, pastel=pastel, title=spec.get("title", ""))
    items = (spec.get("items") or [])[:5]
    n = len(items)
    if not n:
        return
    icons = _consistent_icons(items)
    top = _BODY_TOP
    banner = spec.get("banner")
    if banner:
        ban = slide.shapes.add_shape(_BOX, Inches(_MARGIN), Inches(_EYEBROW_Y), Inches(_CONTENT_W), Inches(0.55))
        ban.fill.solid(); ban.fill.fore_color.rgb = _TEAL2; ban.line.fill.background(); ban.shadow.inherit = False
        tf = ban.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER; p.line_spacing = _LINE_SPACING
        r = p.add_run(); r.text = banner; r.font.size = Pt(_SZ_BODY); r.font.bold = True
        r.font.name = _HEAD; r.font.color.rgb = _WHITE
        top = _EYEBROW_Y + 0.55 + 0.14
    pw = (_CONTENT_W - (n - 1) * _GUTTER) / n
    ph = _BODY_BOTTOM - top
    d = _ICON_DISC
    for i, it in enumerate(items):
        x = _MARGIN + i * (pw + _GUTTER)
        pan = slide.shapes.add_shape(_BOX, Inches(x), Inches(top), Inches(pw), Inches(ph))
        pan.fill.solid(); pan.fill.fore_color.rgb = _TEAL; pan.line.fill.background(); pan.shadow.inherit = False
        hy = top + 0.3
        if icons:
            _icon_disc(slide, x + pw / 2, top + 0.55, d, icon_path=icons[i], light=light, pastel=pastel); hy = top + 1.15
        _place_text(slide, x + _PAD, hy, pw - 2 * _PAD, 0.5, it.get("heading", ""), _SZ_BODY, _WHITE,
                    bold=True, font=_HEAD, align=PP_ALIGN.CENTER)
        _place_text(slide, x + _PAD, hy + 0.55, pw - 2 * _PAD, top + ph - (hy + 0.55) - _PAD,
                    it.get("body", ""), _SZ_SMALL, _ONTEAL, align=PP_ALIGN.CENTER)


def _fill_team(prs, spec: dict, dark_index: int, light_index: int | None = None) -> None:
    """Team cards: 2 to 4 members, each a person chip (light disc + person icon), name, role and short bio."""
    light = light_index is not None and spec.get("background") in ("light", "pastel")
    pastel = light and spec.get("background") == "pastel"
    slide = _synth_slide(prs, light_index if light else dark_index, white=light, pastel=pastel, title=spec.get("title", ""))
    items = (spec.get("items") or [])[:4]
    n = len(items)
    if not n:
        return
    cw = (_CONTENT_W - (n - 1) * _GUTTER) / n
    top = _BODY_TOP + 0.3
    d = 1.5
    ppl = _generic_icon_path("people")
    for i, m in enumerate(items):
        x = _MARGIN + i * (cw + _GUTTER); cx = x + cw / 2
        disc = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - d / 2), Inches(top), Inches(d), Inches(d))
        disc.fill.solid(); disc.fill.fore_color.rgb = _chip_bg(light, pastel); disc.line.fill.background(); disc.shadow.inherit = False
        if ppl:
            pad = d * 0.26
            _place_icon(slide, (Inches(cx - d / 2 + pad), Inches(top + pad), Inches(d - 2 * pad), Inches(d - 2 * pad)), ppl)
        ny = top + d + 0.25
        _place_text(slide, x, ny, cw, 0.4, m.get("name", ""), _SZ_BODY, _ink(light), bold=True, font=_HEAD, align=PP_ALIGN.CENTER)
        _place_text(slide, x, ny + 0.42, cw, 0.35, m.get("role", ""), _SZ_SMALL, _muted(light), bold=True, align=PP_ALIGN.CENTER)
        if m.get("bio"):
            _place_text(slide, x + _PAD, ny + 0.85, cw - 2 * _PAD, _BODY_BOTTOM - (ny + 0.85),
                        m["bio"], _SZ_SMALL, _muted(light), align=PP_ALIGN.CENTER)


def _fill_metric_bars(prs, spec: dict, dark_index: int, light_index: int | None = None) -> None:
    """Metric bars: rows of label + a horizontal bar (length = pct) + the value on the right (bullet-chart
    style). For KPI rows where the magnitude matters."""
    light = light_index is not None and spec.get("background") in ("light", "pastel")
    pastel = light and spec.get("background") == "pastel"
    slide = _synth_slide(prs, light_index if light else dark_index, white=light, pastel=pastel, title=spec.get("title", ""), eyebrow=spec.get("caption"))
    items = (spec.get("items") or [])[:6]
    n = len(items)
    if not n:
        return
    rh = _BODY_H / n
    label_w, val_w = 3.6, 1.3
    bar_x = _MARGIN + label_w + 0.2
    bar_w = _CONTENT_W - label_w - 0.2 - val_w - 0.2
    for i, it in enumerate(items):
        y = _BODY_TOP + i * rh
        cy = y + rh / 2
        _place_text(slide, _MARGIN, y, label_w, rh, it.get("label", ""), _SZ_BODY, _ink(light), bold=True,
                    font=_HEAD, anchor=MSO_ANCHOR.MIDDLE)
        _rule(slide, bar_x, cy - 0.15, bar_w, 0.3, _TEAL2)
        pct = max(0.0, min(100.0, float(it.get("pct", 0) or 0)))
        if pct > 0:
            _rule(slide, bar_x, cy - 0.15, bar_w * pct / 100.0, 0.3, _RED)
        val = it.get("value") or f"{pct:g}%"
        _place_text(slide, bar_x + bar_w + 0.2, y, val_w, rh, val, _SZ_BODY, _ink(light), bold=True,
                    font=_HEAD, anchor=MSO_ANCHOR.MIDDLE)


def _fill_cause_effect(prs, spec: dict, dark_index: int, light_index: int | None = None) -> None:
    """Cause and effect: parallel rows, each a cause panel, a teal arrow, and the resulting effect text."""
    light = light_index is not None and spec.get("background") in ("light", "pastel")
    pastel = light and spec.get("background") == "pastel"
    slide = _synth_slide(prs, light_index if light else dark_index, white=light, pastel=pastel, title=spec.get("title", ""))
    items = (spec.get("items") or [])[:4]
    n = len(items)
    if not n:
        return
    gap = 0.22
    rh = (_BODY_H - (n - 1) * gap) / n
    cause_w, arrow_w = 4.0, 0.7
    eff_x = _MARGIN + cause_w + arrow_w + 0.5
    eff_w = 13.333 - _MARGIN - eff_x
    for i, it in enumerate(items):
        y = _BODY_TOP + i * (rh + gap)
        pan = slide.shapes.add_shape(_BOX, Inches(_MARGIN), Inches(y), Inches(cause_w), Inches(rh))
        pan.fill.solid(); pan.fill.fore_color.rgb = _TEAL; pan.line.fill.background(); pan.shadow.inherit = False
        _place_text(slide, _MARGIN + _PAD, y, cause_w - 2 * _PAD, rh, it.get("heading", ""), _SZ_BODY, _WHITE,
                    bold=True, font=_HEAD, anchor=MSO_ANCHOR.MIDDLE)
        ar = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(_MARGIN + cause_w + 0.2), Inches(y + rh / 2 - 0.22),
                                    Inches(arrow_w), Inches(0.44))
        ar.fill.solid(); ar.fill.fore_color.rgb = _line_soft(light, pastel); ar.line.fill.background(); ar.shadow.inherit = False
        _place_text(slide, eff_x, y, eff_w, rh, it.get("body", ""), _SZ_BODY, _muted(light), anchor=MSO_ANCHOR.MIDDLE)


def _fill_org_chart(prs, spec: dict, dark_index: int, light_index: int | None = None) -> None:
    """Org chart: a top box connected down to 2 to 4 child boxes with a bus connector."""
    light = light_index is not None and spec.get("background") in ("light", "pastel")
    pastel = light and spec.get("background") == "pastel"
    slide = _synth_slide(prs, light_index if light else dark_index, white=light, pastel=pastel, title=spec.get("title", ""))
    items = (spec.get("items") or [])[:4]
    n = len(items)
    if not n:
        return
    tb_w, tb_h = 3.2, 0.8
    tb_x, tb_y = (13.333 - tb_w) / 2, _BODY_TOP + 0.3
    box = slide.shapes.add_shape(_BOX, Inches(tb_x), Inches(tb_y), Inches(tb_w), Inches(tb_h))
    box.fill.solid(); box.fill.fore_color.rgb = _TEAL2; box.line.fill.background(); box.shadow.inherit = False
    _place_text(slide, tb_x, tb_y, tb_w, tb_h, spec.get("center", ""), _SZ_BODY, _WHITE, bold=True,
                font=_HEAD, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    cw = (_CONTENT_W - (n - 1) * _GUTTER) / n
    cy = tb_y + tb_h + 0.9
    chh = min(_BODY_BOTTOM - cy, 2.4)
    bus_y = tb_y + tb_h + 0.45
    centers = [_MARGIN + i * (cw + _GUTTER) + cw / 2 for i in range(n)]
    _rule(slide, 13.333 / 2 - 0.01, tb_y + tb_h, 0.02, bus_y - (tb_y + tb_h), _line_soft(light, pastel))
    _rule(slide, centers[0], bus_y - 0.01, centers[-1] - centers[0], 0.02, _line_soft(light, pastel))
    for i, it in enumerate(items):
        x = _MARGIN + i * (cw + _GUTTER); ccx = x + cw / 2
        _rule(slide, ccx - 0.01, bus_y, 0.02, cy - bus_y, _line_soft(light, pastel))
        cb = slide.shapes.add_shape(_BOX, Inches(x), Inches(cy), Inches(cw), Inches(chh))
        cb.fill.solid(); cb.fill.fore_color.rgb = _TEAL; cb.line.fill.background(); cb.shadow.inherit = False
        _place_text(slide, x + _PAD, cy + 0.18, cw - 2 * _PAD, 0.5, it.get("heading", ""), _SZ_BODY, _WHITE,
                    bold=True, font=_HEAD, align=PP_ALIGN.CENTER)
        _place_text(slide, x + _PAD, cy + 0.72, cw - 2 * _PAD, chh - 0.9, it.get("body", ""), _SZ_SMALL,
                    _ONTEAL, align=PP_ALIGN.CENTER)


def _fill_decision_tree(prs, spec: dict, dark_index: int, light_index: int | None = None) -> None:
    """Decision tree: a root box on the left branching to 2 to 4 outcome boxes on the right."""
    light = light_index is not None and spec.get("background") in ("light", "pastel")
    pastel = light and spec.get("background") == "pastel"
    slide = _synth_slide(prs, light_index if light else dark_index, white=light, pastel=pastel, title=spec.get("title", ""))
    items = (spec.get("items") or [])[:4]
    n = len(items)
    if not n:
        return
    rb_w, rb_h = 3.0, 1.1
    rb_x = _MARGIN
    rb_y = _BODY_TOP + (_BODY_H - rb_h) / 2
    root = slide.shapes.add_shape(_BOX, Inches(rb_x), Inches(rb_y), Inches(rb_w), Inches(rb_h))
    root.fill.solid(); root.fill.fore_color.rgb = _TEAL2; root.line.fill.background(); root.shadow.inherit = False
    _place_text(slide, rb_x + _PAD, rb_y, rb_w - 2 * _PAD, rb_h, spec.get("center", ""), _SZ_BODY, _WHITE,
                bold=True, font=_HEAD, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    bus_x = rb_x + rb_w + 0.6
    bx = bus_x + 0.6
    bw = 13.333 - _MARGIN - bx
    gap = 0.25
    bh = (_BODY_H - (n - 1) * gap) / n
    root_cy = rb_y + rb_h / 2
    centers = [_BODY_TOP + i * (bh + gap) + bh / 2 for i in range(n)]
    _rule(slide, rb_x + rb_w, root_cy - 0.01, 0.6, 0.02, _line_soft(light, pastel))
    _rule(slide, bus_x - 0.01, min(centers[0], root_cy), 0.02, max(centers[-1], root_cy) - min(centers[0], root_cy), _line_soft(light, pastel))
    for i, it in enumerate(items):
        y = _BODY_TOP + i * (bh + gap); bcy = y + bh / 2
        _rule(slide, bus_x, bcy - 0.01, bx - bus_x, 0.02, _line_soft(light, pastel))
        bb = slide.shapes.add_shape(_BOX, Inches(bx), Inches(y), Inches(bw), Inches(bh))
        bb.fill.solid(); bb.fill.fore_color.rgb = _TEAL; bb.line.fill.background(); bb.shadow.inherit = False
        _place_text(slide, bx + _PAD, y + 0.12, bw - 2 * _PAD, 0.45, it.get("heading", ""), _SZ_BODY, _WHITE,
                    bold=True, font=_HEAD)
        if it.get("body"):
            _place_text(slide, bx + _PAD, y + 0.6, bw - 2 * _PAD, bh - 0.72, it["body"], _SZ_SMALL, _ONTEAL)


def _fill_cycle(prs, spec: dict, dark_index: int, light_index: int | None = None) -> None:
    """Cycle: 3 to 6 labelled nodes arranged in a ring around an optional hub."""
    import math
    light = light_index is not None and spec.get("background") in ("light", "pastel")
    pastel = light and spec.get("background") == "pastel"
    slide = _synth_slide(prs, light_index if light else dark_index, white=light, pastel=pastel, title=spec.get("title", ""))
    items = (spec.get("items") or [])[:6]
    n = len(items)
    if not n:
        return
    cx, cy = 13.333 / 2, _BODY_TOP + _BODY_H / 2 + 0.05
    ring_r, node_d = 1.55, 1.5
    if spec.get("center"):
        hub_d = 1.5
        hub = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - hub_d / 2), Inches(cy - hub_d / 2), Inches(hub_d), Inches(hub_d))
        hub.fill.solid(); hub.fill.fore_color.rgb = _TEAL2; hub.line.fill.background(); hub.shadow.inherit = False
        _place_text(slide, cx - hub_d / 2, cy - hub_d / 2, hub_d, hub_d, spec["center"], _SZ_SMALL, _WHITE,
                    bold=True, font=_HEAD, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    for i, it in enumerate(items):
        ang = -math.pi / 2 + 2 * math.pi * i / n
        nx, ny = cx + ring_r * math.cos(ang), cy + ring_r * math.sin(ang)
        node = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(nx - node_d / 2), Inches(ny - node_d / 2), Inches(node_d), Inches(node_d))
        node.fill.solid(); node.fill.fore_color.rgb = _TEAL_TINTS[i % len(_TEAL_TINTS)]
        node.line.color.rgb = _ink(light); node.line.width = Pt(1.5); node.shadow.inherit = False
        tf = node.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = tf.margin_right = Inches(0.06)
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER; p.line_spacing = 1.0
        r = p.add_run(); r.text = it.get("heading", ""); r.font.size = Pt(_SZ_SMALL); r.font.bold = True
        r.font.name = _HEAD; r.font.color.rgb = _WHITE


def _fill_gantt(prs, spec: dict, dark_index: int, light_index: int | None = None) -> None:
    """Gantt / project schedule: task rows against a period axis. Each task spans `start`..`end`
    columns as a bar; a task flagged `milestone` renders as a red diamond at its `start` period."""
    light = light_index is not None and spec.get("background") in ("light", "pastel")
    pastel = light and spec.get("background") == "pastel"
    slide = _synth_slide(prs, light_index if light else dark_index, white=light, pastel=pastel, title=spec.get("title", ""), eyebrow=spec.get("caption"))
    periods = spec.get("periods") or []
    items = (spec.get("items") or [])[:8]
    npp, n = len(periods), len(items)
    if not n or not npp:
        return
    label_w = 3.4
    grid_x = _MARGIN + label_w + 0.15
    grid_w = _CONTENT_W - label_w - 0.15
    col_w = grid_w / npp
    header_h = 0.42
    top = _BODY_TOP + header_h + 0.12
    rows_h = _BODY_BOTTOM - top
    for j, p in enumerate(periods):
        px = grid_x + j * col_w
        _place_text(slide, px, _BODY_TOP, col_w, header_h, p, _SZ_SMALL, _muted(light), bold=True,
                    font=_HEAD, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        _rule(slide, px, top, 0.015, rows_h, _TEAL2)
    _rule(slide, grid_x + grid_w, top, 0.015, rows_h, _TEAL2)
    rh = rows_h / n
    bar_h = min(0.34, rh * 0.52)
    for i, it in enumerate(items):
        y = top + i * rh
        cy = y + rh / 2
        _place_text(slide, _MARGIN, y, label_w, rh, it.get("label", ""), _SZ_SMALL, _ink(light),
                    bold=True, font=_HEAD, anchor=MSO_ANCHOR.MIDDLE)
        s = max(1, min(int(it.get("start", 1) or 1), npp))
        if it.get("milestone"):
            mx = grid_x + (s - 0.5) * col_w
            dsz = min(0.3, rh * 0.5)
            dia = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, Inches(mx - dsz / 2), Inches(cy - dsz / 2),
                                         Inches(dsz), Inches(dsz))
            dia.fill.solid(); dia.fill.fore_color.rgb = _RED; dia.line.fill.background(); dia.shadow.inherit = False
        else:
            e = max(s, min(int(it.get("end", s) or s), npp))
            bx = grid_x + (s - 1) * col_w + 0.06
            bw = (e - s + 1) * col_w - 0.12
            bar = slide.shapes.add_shape(_BOX, Inches(bx), Inches(cy - bar_h / 2), Inches(bw), Inches(bar_h))
            bar.fill.solid(); bar.fill.fore_color.rgb = _TEAL_TINTS[i % len(_TEAL_TINTS)]
            bar.line.fill.background(); bar.shadow.inherit = False
            if it.get("note") and bw >= 1.0:
                _place_text(slide, bx, cy - bar_h / 2, bw, bar_h, it["note"], _SZ_SMALL, _WHITE,
                            bold=True, font=_HEAD, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def _sine_spine(slide, x0, x1, cy, amp, n_nodes, color, width_pt=2.25, segments=160):
    """A smooth S-curve polyline (approximated with many short segments, since python-pptx freeforms
    have no bezier API) running left→right. The phase is set so node i — sampled at t=(i+0.5)/n —
    lands exactly on a crest or trough, alternating high/low. Returns those node (x, y) points."""
    import math
    span = x1 - x0

    def y_at(t):                      # crest/trough at t = (i+0.5)/n_nodes
        return cy - amp * math.cos((t * n_nodes - 0.5) * math.pi)

    pts = [(x0 + span * i / segments, y_at(i / segments)) for i in range(segments + 1)]
    try:
        b = slide.shapes.build_freeform(Inches(pts[0][0]), Inches(pts[0][1]))
        b.add_line_segments([(Inches(x), Inches(y)) for x, y in pts[1:]], close=False)
        sh = b.convert_to_shape()
        sh.fill.background(); sh.line.color.rgb = color; sh.line.width = Pt(width_pt)
        sh.shadow.inherit = False
    except Exception:  # noqa: BLE001 — the spine is decorative; never break the render
        pass
    return [(x0 + span * (i + 0.5) / n_nodes, y_at((i + 0.5) / n_nodes)) for i in range(n_nodes)]


def _fill_serpentine(prs, spec: dict, dark_index: int, light_index: int | None = None) -> None:
    """Serpentine flow: 3 or 4 stages threaded on an S-curve, numbered discs sitting on the crests and
    each stage's text alternating above / below — text always on the OUTER side of its crest, so the
    two text bands stay clear of the curve. For a narrative sequence of shifts, forces or dated events.

    A per-item `date` makes this the "wavy timeline": the date takes the slot the icon chip would use,
    because at this band height there is only room for one of them and a date is data while the icon is
    decoration. Undated items keep the icon chip as before."""
    light = light_index is not None and spec.get("background") in ("light", "pastel")
    pastel = light and spec.get("background") == "pastel"
    slide = _synth_slide(prs, light_index if light else dark_index, white=light, pastel=pastel, title=spec.get("title", ""), eyebrow=spec.get("caption"))
    items = (spec.get("items") or [])[:4]
    n = len(items)
    if not n:
        return
    dated = any(it.get("date") for it in items)
    icons = None if dated else _consistent_icons(items)
    band_h = 1.45                              # text band at the top and at the bottom
    cy = _BODY_TOP + _BODY_H / 2
    amp = 0.55
    nodes = _sine_spine(slide, _MARGIN + 0.35, _MARGIN + _CONTENT_W - 0.35, cy, amp, n, _line_soft(light, pastel))
    disc, ic = 0.46, 0.52
    cw = _CONTENT_W / n
    for i, (it, (nx, ny)) in enumerate(zip(items, nodes)):
        high = i % 2 == 0                      # even nodes crest UP → their text goes in the top band
        d = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(nx - disc / 2), Inches(ny - disc / 2),
                                   Inches(disc), Inches(disc))
        d.fill.solid(); d.fill.fore_color.rgb = _TEAL2
        d.line.color.rgb = _line_soft(light, pastel); d.line.width = Pt(1.25); d.shadow.inherit = False
        tf = d.text_frame; tf.word_wrap = False; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER; p.line_spacing = 1.0
        r = p.add_run(); r.text = f"{i + 1:02d}"; r.font.size = Pt(_SZ_SMALL); r.font.bold = True
        r.font.name = _HEAD; r.font.color.rgb = _WHITE
        tx = min(max(nx - cw / 2, _MARGIN), _MARGIN + _CONTENT_W - cw)
        if high:                               # top band: heading, body, then the icon nearest the curve
            hy, byy = _BODY_TOP, _BODY_TOP + 0.48
            icy = _BODY_TOP + band_h - ic / 2
        else:                                  # bottom band: icon nearest the curve, then heading, body
            icy = _BODY_BOTTOM - band_h + ic / 2
            hy, byy = _BODY_BOTTOM - band_h + ic + 0.06, _BODY_BOTTOM - band_h + ic + 0.54
        if icons:
            _icon_disc(slide, nx, icy, ic, icon_path=icons[i], light=light, pastel=pastel)
        elif it.get("date"):
            # Same slot as the icon chip: nearest the curve, so the date reads as the node's label.
            _place_text(slide, tx, icy - 0.15, cw, 0.3, it["date"], _SZ_SMALL, _muted(light), bold=True,
                        font=_HEAD, align=PP_ALIGN.CENTER)
        _place_text(slide, tx, hy, cw, 0.46, it.get("heading", ""), _SZ_BODY, _ink(light), bold=True,
                    font=_HEAD, align=PP_ALIGN.CENTER)
        if it.get("body"):
            _place_text(slide, tx + 0.12, byy, cw - 0.24, 0.44, it["body"], _SZ_SMALL,
                        _muted(light), align=PP_ALIGN.CENTER)


def _fill_coverage_matrix(prs, spec: dict, light_index: int) -> None:
    """Coverage matrix: entities as rows, capabilities as columns, a filled tick where the entity
    covers that capability. The right home for BINARY yes/no comparisons (harvey balls misread those)."""
    pastel = spec.get("background") == "pastel"
    slide = _synth_slide(prs, light_index, white=True, pastel=pastel, title=spec.get("title", ""),
                         eyebrow=spec.get("caption"))
    heads = spec.get("headers") or []
    items = (spec.get("items") or [])[:5]
    nc, nr = len(heads), len(items)
    if not nc or not nr:
        return
    label_w = 3.5
    grid_x = _MARGIN + label_w + 0.25
    col_w = (_CONTENT_W - label_w - 0.25) / nc
    head_h = 0.62
    top = _BODY_TOP + head_h + 0.14
    for j, h in enumerate(heads):
        _place_text(slide, grid_x + j * col_w, _BODY_TOP, col_w, head_h, h, _SZ_SMALL, _TEAL,
                    bold=True, font=_HEAD, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.BOTTOM)
    _rule(slide, _MARGIN, top - 0.07, _CONTENT_W, 0.02, _TEAL)
    rh = (_BODY_BOTTOM - top) / nr
    tick = min(0.42, rh * 0.46)
    for i, it in enumerate(items):
        y = top + i * rh
        _place_text(slide, _MARGIN, y + 0.1, label_w, 0.42, it.get("label", ""), _SZ_BODY, _INKC,
                    bold=True, font=_HEAD)
        if it.get("body"):
            _place_text(slide, _MARGIN, y + 0.54, label_w - 0.15, rh - 0.66, it["body"], _SZ_SMALL, _TEAL2)
        if i:
            _rule(slide, _MARGIN, y, _CONTENT_W, 0.01, _WHITE if pastel else _PANEL)
        for j, on in enumerate(list(it.get("marks") or [])[:nc]):
            if not on:
                continue
            cx = grid_x + j * col_w + col_w / 2
            box = slide.shapes.add_shape(_BOX, Inches(cx - tick / 2), Inches(y + rh / 2 - tick / 2),
                                         Inches(tick), Inches(tick))
            box.fill.solid(); box.fill.fore_color.rgb = _TEAL
            box.line.fill.background(); box.shadow.inherit = False
            tf = box.text_frame; tf.word_wrap = False; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER; p.line_spacing = 1.0
            r = p.add_run(); r.text = "✓"; r.font.size = Pt(_SZ_BODY); r.font.bold = True
            r.font.name = _HEAD; r.font.color.rgb = _WHITE


def _fill_photo_stats(prs, spec: dict, dark_index: int, light_index: int | None = None) -> None:
    """Photo-topped stat cards: 2 or 3 cards, each a photo above a solid panel carrying an eyebrow
    label, ONE hero figure and a supporting line. High-impact proof-point opener."""
    light = light_index is not None and spec.get("background") in ("light", "pastel")
    pastel = light and spec.get("background") == "pastel"
    slide = _synth_slide(prs, light_index if light else dark_index, white=light, pastel=pastel, title=spec.get("title", ""), eyebrow=spec.get("caption"))
    items = (spec.get("items") or [])[:3]
    n = len(items)
    if not n:
        return
    cw = (_CONTENT_W - (n - 1) * _GUTTER) / n
    top = _BODY_TOP + 0.18
    total_h = _BODY_BOTTOM - top
    photo_h = total_h * 0.56
    pan_h = total_h - photo_h
    for i, it in enumerate(items):
        x = _MARGIN + i * (cw + _GUTTER)
        placed = False
        aid = it.get("asset_id")
        if aid:
            try:
                path = _photo_path(aid)
                if path is not None:
                    _place_cropped(slide, path, x, top, cw, photo_h); placed = True
            except Exception:  # noqa: BLE001
                placed = False
        if not placed:
            ph = slide.shapes.add_shape(_BOX, Inches(x), Inches(top), Inches(cw), Inches(photo_h))
            ph.fill.solid(); ph.fill.fore_color.rgb = _TEAL2
            ph.line.fill.background(); ph.shadow.inherit = False
        pan = slide.shapes.add_shape(_BOX, Inches(x), Inches(top + photo_h), Inches(cw), Inches(pan_h))
        pan.fill.solid(); pan.fill.fore_color.rgb = _TEAL
        pan.line.fill.background(); pan.shadow.inherit = False
        # label / hero figure / note stacked and centred as a group inside the panel
        note_h = 0.62 if it.get("note") else 0.0
        blk_h = 0.3 + 0.95 + note_h
        py = top + photo_h + max(0.12, (pan_h - blk_h) / 2)
        _place_text(slide, x + _PAD, py, cw - 2 * _PAD, 0.3, (it.get("label") or "").upper(),
                    _SZ_SMALL, _LTEAL, bold=True, font=_HEAD, align=PP_ALIGN.CENTER)
        _place_text(slide, x + _PAD, py + 0.3, cw - 2 * _PAD, 0.95, it.get("value", ""),
                    _SZ_HERO, _WHITE, bold=True, font=_HEAD, align=PP_ALIGN.CENTER,
                    anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
        if it.get("note"):
            _place_text(slide, x + _PAD, py + 1.27, cw - 2 * _PAD, note_h, it["note"],
                        _SZ_SMALL, _ONTEAL, align=PP_ALIGN.CENTER)


def _est_lines(text, width_in, size_pt) -> int:
    """Rough wrapped-line count for a run of text in a box of the given width — enough to size a
    panel to its content instead of stretching it over the whole body zone."""
    if not text:
        return 0
    per_line = max(8, int(width_in * 145.0 / size_pt))   # ~13 chars/inch at 11pt, ~10 at 14pt
    return max(1, -(-len(text) // per_line))


def _fill_numbered_cards(prs, spec: dict, dark_index: int, light_index: int | None = None) -> None:
    """Numbered cards: 2 to 4 equal panels, each with a number badge top-left, an optional icon chip
    top-right, a bold heading and a body. The 'three reasons why' card set. Cards are sized to the
    LONGEST card's content (equal heights) and centred in the body zone, so they never stretch empty."""
    light = light_index is not None and spec.get("background") in ("light", "pastel")
    pastel = light and spec.get("background") == "pastel"
    slide = _synth_slide(prs, light_index if light else dark_index, white=light, pastel=pastel, title=spec.get("title", ""))
    items = (spec.get("items") or [])[:4]
    n = len(items)
    if not n:
        return
    icons = _consistent_icons(items)
    cw = (_CONTENT_W - (n - 1) * _GUTTER) / n
    badge = 0.52
    tw = cw - 2 * _PAD
    need = 0.0
    for it in items:
        h = _PAD + badge + 0.2 + _est_lines(it.get("heading"), tw, _SZ_BODY) * 0.25 + _PAD
        if it.get("body"):
            h += 0.14 + _est_lines(it["body"], tw, _SZ_SMALL) * 0.20
        need = max(need, h)
    ch = min(_BODY_H, max(1.9, need) + 1.1)      # content + breathing room inside the card
    top = _BODY_TOP + (_BODY_H - ch) / 2         # balanced margin above and below
    for i, it in enumerate(items):
        x = _MARGIN + i * (cw + _GUTTER)
        pan = slide.shapes.add_shape(_BOX, Inches(x), Inches(top), Inches(cw), Inches(ch))
        pan.fill.solid(); pan.fill.fore_color.rgb = _TEAL
        pan.line.fill.background(); pan.shadow.inherit = False
        bd = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + _PAD), Inches(top + _PAD),
                                    Inches(badge), Inches(badge))
        bd.fill.solid(); bd.fill.fore_color.rgb = _TEAL2
        bd.line.fill.background(); bd.shadow.inherit = False
        tf = bd.text_frame; tf.word_wrap = False; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER; p.line_spacing = 1.0
        r = p.add_run(); r.text = str(i + 1); r.font.size = Pt(_SZ_BODY); r.font.bold = True
        r.font.name = _HEAD; r.font.color.rgb = _WHITE
        if icons:                              # brand icon chip (light disc), matching the rest of the library
            _icon_disc(slide, x + cw - _PAD - badge / 2, top + _PAD + badge / 2, badge + 0.06,
                       icon_path=icons[i], light=light, pastel=pastel)
        # heading + body as ONE vertically-centred block, so a short card never leaves a dead gap
        ty = top + _PAD + badge + 0.2
        tb = slide.shapes.add_textbox(Inches(x + _PAD), Inches(ty), Inches(cw - 2 * _PAD),
                                      Inches(top + ch - ty - _PAD))
        tf = tb.text_frame
        tf.word_wrap = True
        try:
            tf.auto_size = MSO_AUTO_SIZE.NONE
        except Exception:  # noqa: BLE001
            pass
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = tf.margin_right = Emu(0)
        ph = tf.paragraphs[0]; ph.line_spacing = _LINE_SPACING
        rh_ = ph.add_run(); rh_.text = it.get("heading", "")
        rh_.font.size = Pt(_SZ_BODY); rh_.font.bold = True
        rh_.font.name = _HEAD; rh_.font.color.rgb = _WHITE
        if it.get("body"):
            pb = tf.add_paragraph(); pb.line_spacing = _LINE_SPACING
            pb.space_before = Pt(10)
            rb = pb.add_run(); rb.text = it["body"]
            rb.font.size = Pt(_SZ_SMALL); rb.font.name = _BODY; rb.font.color.rgb = _ONTEAL


def _fill_implications(prs, spec: dict, dark_index: int, light_index: int | None = None) -> None:
    """Trend / overview / implication rows: a numbered label pill, the detail, then a chevron into the
    'so what'. The classic three-column analysis table (mega-trends → market implications)."""
    light = light_index is not None and spec.get("background") in ("light", "pastel")
    pastel = light and spec.get("background") == "pastel"
    slide = _synth_slide(prs, light_index if light else dark_index, white=light, pastel=pastel, title=spec.get("title", ""))
    items = (spec.get("items") or [])[:5]
    n = len(items)
    if not n:
        return
    heads = spec.get("headers") or []
    lab_w, imp_w, arrow_w = 3.0, 3.5, 0.42
    ov_x = _MARGIN + lab_w + 0.25
    ov_w = _CONTENT_W - lab_w - 0.25 - arrow_w - imp_w - 0.4
    ar_x = ov_x + ov_w + 0.2
    imp_x = ar_x + arrow_w + 0.2
    top = _BODY_TOP
    if heads[:3]:
        hh = 0.4
        for hx, hw, ht in ((_MARGIN, lab_w, heads[0]), (ov_x, ov_w, heads[1]), (imp_x, imp_w, heads[2])):
            _place_text(slide, hx, top, hw, hh, ht, _SZ_SMALL, _muted(light), bold=True, font=_HEAD)
            _rule(slide, hx, top + hh, hw, 0.02, _TEAL2)
        top += hh + 0.16
    gap = 0.18
    rh = (_BODY_BOTTOM - top - (n - 1) * gap) / n
    badge = min(0.46, rh * 0.62)
    for i, it in enumerate(items):
        y = top + i * (rh + gap)
        pill = slide.shapes.add_shape(_BOX, Inches(_MARGIN), Inches(y), Inches(lab_w), Inches(rh))
        pill.fill.solid(); pill.fill.fore_color.rgb = _TEAL
        pill.line.fill.background(); pill.shadow.inherit = False
        _icon_disc(slide, _MARGIN + 0.1 + badge / 2, y + rh / 2, badge, number=i + 1, light=light, pastel=pastel)
        _place_text(slide, _MARGIN + 0.2 + badge, y, lab_w - 0.3 - badge, rh, it.get("heading", ""),
                    _SZ_BODY, _WHITE, bold=True, font=_HEAD, anchor=MSO_ANCHOR.MIDDLE)
        _place_text(slide, ov_x, y, ov_w, rh, it.get("body", ""), _SZ_SMALL, _muted(light),
                    anchor=MSO_ANCHOR.MIDDLE)
        ar = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(ar_x), Inches(y + rh / 2 - 0.17),
                                    Inches(arrow_w), Inches(0.34))
        ar.fill.solid(); ar.fill.fore_color.rgb = _line_soft(light, pastel)
        ar.line.fill.background(); ar.shadow.inherit = False
        ip = slide.shapes.add_shape(_BOX, Inches(imp_x), Inches(y), Inches(imp_w), Inches(rh))
        ip.fill.solid(); ip.fill.fore_color.rgb = _TEAL2
        ip.line.fill.background(); ip.shadow.inherit = False
        _place_text(slide, imp_x + _PAD, y, imp_w - 2 * _PAD, rh, it.get("implication", ""),
                    _SZ_SMALL, _ONTEAL, anchor=MSO_ANCHOR.MIDDLE)


def _fill_breakdown(prs, spec: dict, dark_index: int, light_index: int | None = None) -> None:
    """Total broken into shares: a hub circle carrying the total, fanning out via thin connectors to
    one bar per component. Each bar shows its share and label, tinted largest-to-smallest."""
    light = light_index is not None and spec.get("background") in ("light", "pastel")
    pastel = light and spec.get("background") == "pastel"
    slide = _synth_slide(prs, light_index if light else dark_index, white=light, pastel=pastel, title=spec.get("title", ""))
    items = (spec.get("items") or [])[:6]
    n = len(items)
    if not n:
        return
    hub_d = 2.5
    hub_cx = _MARGIN + hub_d / 2
    hub_cy = _BODY_TOP + _BODY_H / 2 - (0.18 if spec.get("caption") else 0.0)
    bars_x = _MARGIN + hub_d + 1.15
    bars_w = 13.333 - _MARGIN - bars_x
    gap = 0.18
    bar_h = (_BODY_H - (n - 1) * gap) / n
    pct_w = min(1.25, bars_w * 0.22)
    # One decimal for every share if ANY of them has a fraction, so the column reads consistently
    # (42.9% / 8.0%, never 42.9% / 8%).
    dec = 1 if any(float(i.get("pct", 0) or 0) % 1 for i in items) else 0
    for i, it in enumerate(items):                    # connectors first, so bars sit on top of them
        bcy = _BODY_TOP + i * (bar_h + gap) + bar_h / 2
        cn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(hub_cx + hub_d / 2),
                                        Inches(hub_cy), Inches(bars_x), Inches(bcy))
        cn.line.color.rgb = _TEAL2; cn.line.width = Pt(1.0); cn.shadow.inherit = False
    hub = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(hub_cx - hub_d / 2), Inches(hub_cy - hub_d / 2),
                                 Inches(hub_d), Inches(hub_d))
    hub.fill.solid(); hub.fill.fore_color.rgb = _TEAL
    hub.line.fill.background(); hub.shadow.inherit = False
    _place_text(slide, hub_cx - hub_d / 2, hub_cy - 0.5, hub_d, 1.0, spec.get("total", ""),
                _SZ_HERO, _WHITE, bold=True, font=_HEAD, align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
    if spec.get("caption"):
        _place_text(slide, hub_cx - hub_d / 2, hub_cy + hub_d / 2 + 0.1, hub_d, 0.4,
                    spec["caption"], _SZ_SMALL, _muted(light), bold=True, font=_HEAD, align=PP_ALIGN.CENTER)
    for i, it in enumerate(items):
        y = _BODY_TOP + i * (bar_h + gap)
        col = _TEAL_TINTS[i % len(_TEAL_TINTS)]
        bar = slide.shapes.add_shape(_BOX, Inches(bars_x), Inches(y), Inches(bars_w), Inches(bar_h))
        bar.fill.solid(); bar.fill.fore_color.rgb = col
        bar.line.fill.background(); bar.shadow.inherit = False
        # The share sits INSIDE the bar with a light hairline separating it from the label. An earlier
        # version used a dark _INKC block here, but that is almost the slide background colour, so the
        # figures looked like they were floating outside the bar.
        pct = float(it.get("pct", 0) or 0)
        _place_text(slide, bars_x + 0.1, y, pct_w - 0.15, bar_h, f"{pct:.{dec}f}%", _SZ_BODY, _WHITE,
                    bold=True, font=_HEAD, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
        _rule(slide, bars_x + pct_w, y + bar_h * 0.22, 0.015, bar_h * 0.56, _ONTEAL)
        _place_text(slide, bars_x + pct_w + _PAD, y, bars_w - pct_w - 2 * _PAD, bar_h,
                    it.get("label", ""), _SZ_BODY, _WHITE, bold=True, font=_HEAD,
                    anchor=MSO_ANCHOR.MIDDLE)


def _fill_chart_bands(prs, spec: dict, dark_index: int, light_index: int | None = None) -> None:
    """Column chart with narrative phase bands: bars drawn as shapes (NOT a native chart) so the
    band pills under the axis line up EXACTLY with their category columns — the whole point of the
    layout. Bands carry a number badge and span `start`..`end` categories."""
    light = light_index is not None and spec.get("background") in ("light", "pastel")
    pastel = light and spec.get("background") == "pastel"
    slide = _synth_slide(prs, light_index if light else dark_index, white=light, pastel=pastel, title=spec.get("title", ""), eyebrow=spec.get("caption"))
    cats = spec.get("categories") or []
    vals = [float(v) for v in (spec.get("values") or [])]
    bands = (spec.get("bands") or [])[:4]
    n = min(len(cats), len(vals))
    if n < 1:
        return
    cats, vals = cats[:n], vals[:n]
    band_h = 0.46
    band_y = _BODY_BOTTOM - band_h
    axis_y = band_y - 0.42                            # baseline: category labels sit just above it
    plot_top = _BODY_TOP + 0.34                       # room for the value labels above the tallest bar
    col_w = _CONTENT_W / n
    peak = max(vals) or 1.0
    span = axis_y - plot_top
    if spec.get("y_axis"):
        _place_text(slide, _MARGIN, _BODY_TOP - 0.04, 5.0, 0.3, spec["y_axis"], _SZ_SMALL,
                    _muted(light), bold=True, font=_HEAD)
    _rule(slide, _MARGIN, axis_y, _CONTENT_W, 0.02, _TEAL2)
    bw = col_w * 0.54
    for i, (c, v) in enumerate(zip(cats, vals)):
        cx = _MARGIN + (i + 0.5) * col_w
        h = max(0.03, span * (v / peak))
        bar = slide.shapes.add_shape(_BOX, Inches(cx - bw / 2), Inches(axis_y - h), Inches(bw), Inches(h))
        bar.fill.solid(); bar.fill.fore_color.rgb = _TEAL2
        bar.line.fill.background(); bar.shadow.inherit = False
        _place_text(slide, cx - col_w / 2, axis_y - h - 0.3, col_w, 0.28, f"{v:g}", _SZ_SMALL,
                    _ink(light), bold=True, font=_HEAD, align=PP_ALIGN.CENTER)
        _place_text(slide, cx - col_w / 2, axis_y + 0.06, col_w, 0.3, c, _SZ_SMALL, _muted(light),
                    bold=True, font=_HEAD, align=PP_ALIGN.CENTER)
    for j, b in enumerate(bands):
        s = max(1, min(int(b.get("start", 1) or 1), n))
        e = max(s, min(int(b.get("end", s) or s), n))
        bx = _MARGIN + (s - 1) * col_w + 0.04
        bwd = (e - s + 1) * col_w - 0.08
        if s > 1:                                     # dotted divider at the band boundary
            dx = _MARGIN + (s - 1) * col_w
            dv = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(dx), Inches(plot_top),
                                            Inches(dx), Inches(axis_y))
            dv.line.color.rgb = _TEAL2; dv.line.width = Pt(1.0)
            try:
                dv.line.dash_style = MSO_LINE_DASH_STYLE.ROUND_DOT
            except Exception:  # noqa: BLE001 — solid is an acceptable fallback
                pass
            dv.shadow.inherit = False
        pill = slide.shapes.add_shape(_BOX, Inches(bx), Inches(band_y), Inches(bwd), Inches(band_h))
        pill.fill.solid(); pill.fill.fore_color.rgb = _TEAL_TINTS[j % len(_TEAL_TINTS)]
        pill.line.fill.background(); pill.shadow.inherit = False
        bd = min(0.3, band_h - 0.14)
        _icon_disc(slide, bx + 0.12 + bd / 2, band_y + band_h / 2, bd, number=j + 1, light=light, pastel=pastel)
        _place_text(slide, bx + 0.2 + bd, band_y, bwd - 0.3 - bd, band_h, b.get("label", ""),
                    _SZ_BODY, _WHITE, bold=True, font=_HEAD, align=PP_ALIGN.CENTER,
                    anchor=MSO_ANCHOR.MIDDLE)


def _nice_max(v: float) -> float:
    """Round an axis maximum up to a readable step, so ticks land on round numbers. The ladder is
    deliberately fine-grained: with only 1/2/2.5/5 available, a max of 29 jumped to 50 and squeezed
    every point into the left half of the plot."""
    if v <= 0:
        return 1.0
    mag = 10.0 ** math.floor(math.log10(v))
    for step in (1.0, 1.2, 1.5, 2.0, 2.5, 3.0, 4.0, 5.0, 6.0, 8.0, 10.0):
        if v <= step * mag * 1.0000001:
            return step * mag
    return 10.0 * mag


def _fmt_num(v: float) -> str:
    return f"{v:g}"


def _fill_chart_takeaways(prs, spec: dict, light_index: int) -> None:
    """Bubble chart beside a key-takeaways column — the classic 'here is the landscape AND what it
    means' slide. Bubbles are drawn as shapes (python-pptx exposes no plot-area geometry for a native
    bubble chart, so labels could not be anchored to their points) with area proportional to `size`."""
    pastel = spec.get("background") == "pastel"
    slide = _synth_slide(prs, light_index, white=True, pastel=pastel, title=spec.get("title", ""))
    bubbles = (spec.get("bubbles") or [])[:12]
    takeaways = [t for t in (spec.get("takeaways") or []) if t and t.strip()][:5]
    if not bubbles:
        return
    heads = spec.get("headers") or []
    note = spec.get("bottom_note")
    left_w = 8.6
    right_x = _MARGIN + left_w + 0.35
    right_w = 13.333 - _MARGIN - right_x
    top = _BODY_TOP
    if heads[:2]:                                       # column headers, each over a rule
        for hx, hw, ht in ((_MARGIN, left_w, heads[0]), (right_x, right_w, heads[1])):
            _place_text(slide, hx, top, hw, 0.32, ht, _SZ_BODY, _INKC, bold=True, font=_HEAD)
            _rule(slide, hx, top + 0.34, hw, 0.025, _TEAL)
        top += 0.46
    bottom = _BODY_BOTTOM - (0.34 if note else 0.0)
    # plot frame: room on the left for y tick labels, below for x ticks + the axis title
    px0 = _MARGIN + 0.68
    pw = left_w - 0.68
    py0 = top + 0.32                                    # the y-axis title sits in this strip
    py1 = bottom - 0.52
    ph = py1 - py0
    if ph < 1.0 or pw < 1.0:
        return
    xs = [float(b.get("x", 0) or 0) for b in bubbles]
    ys = [float(b.get("y", 0) or 0) for b in bubbles]
    sizes = [float(b.get("size", 0) or 0) for b in bubbles]
    x_max = _nice_max(max(xs) * 1.12 if max(xs) > 0 else 1.0)
    y_max = _nice_max(max(ys) * 1.12 if max(ys) > 0 else 1.0)
    s_max = max(sizes) or 1.0
    _place_text(slide, _MARGIN, top, left_w, 0.3, spec.get("y_axis", ""), _SZ_SMALL, _TEAL,
                bold=True, font=_HEAD)
    frame = slide.shapes.add_shape(_BOX, Inches(px0), Inches(py0), Inches(pw), Inches(ph))
    frame.fill.background(); frame.line.color.rgb = _LTEAL; frame.line.width = Pt(0.75)
    frame.shadow.inherit = False
    for k in range(5):                                  # y ticks + faint gridlines
        v = y_max * k / 4
        gy = py1 - ph * k / 4
        if k:
            _rule(slide, px0, gy, pw, 0.008, _WHITE if pastel else _PANEL)
        _place_text(slide, _MARGIN, gy - 0.13, 0.58, 0.26, _fmt_num(v), _SZ_SMALL, _TEAL2,
                    align=PP_ALIGN.RIGHT)
    for k in range(5):                                  # x ticks
        v = x_max * k / 4
        gx = px0 + pw * k / 4
        _place_text(slide, gx - 0.5, py1 + 0.05, 1.0, 0.26, _fmt_num(v), _SZ_SMALL, _TEAL2,
                    align=PP_ALIGN.CENTER)
    _place_text(slide, px0, py1 + 0.3, pw, 0.28, spec.get("x_axis", ""), _SZ_SMALL, _TEAL,
                bold=True, font=_HEAD, align=PP_ALIGN.CENTER)
    r_max = min(0.5, ph * 0.15)
    order = sorted(range(len(bubbles)), key=lambda i: -sizes[i])   # big first, so small sit on top
    geo = []
    for rank, i in enumerate(order):
        bx = px0 + pw * (xs[i] / x_max if x_max else 0)
        by = py1 - ph * (ys[i] / y_max if y_max else 0)
        r = r_max * math.sqrt(sizes[i] / s_max) if s_max and sizes[i] > 0 else r_max * 0.3
        geo.append((i, bx, by, max(r, 0.075), rank))
    for i, bx, by, r, rank in geo:
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(bx - r), Inches(by - r),
                                     Inches(2 * r), Inches(2 * r))
        dot.fill.solid(); dot.fill.fore_color.rgb = _TEAL_TINTS[rank % len(_TEAL_TINTS)]
        dot.line.color.rgb = _WHITE; dot.line.width = Pt(0.75); dot.shadow.inherit = False
        if r >= 0.24 and sizes[i] > 0:                  # value inside the bubble when it fits
            tf = dot.text_frame; tf.word_wrap = False; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER; p.line_spacing = 1.0
            run = p.add_run(); run.text = _fmt_num(sizes[i]); run.font.size = Pt(_SZ_SMALL)
            run.font.bold = True; run.font.name = _HEAD; run.font.color.rgb = _WHITE
    # Labels, placed greedily (largest bubble first) into the first candidate slot that hits neither
    # another label nor another bubble. A fixed above/below alternation collided badly wherever several
    # small bubbles clustered.
    def _ov_area(a, rects):
        """Total area of `a` covered by `rects`. Used instead of a boolean hit test so that when a
        dense cluster leaves NO free slot we can still choose the least-bad one — the earlier version
        fell back to the first candidate, which guaranteed a collision."""
        ax, ay, aw, ah = a
        tot = 0.0
        for bx2, by2, bw2, bh2 in rects:
            ox = min(ax + aw, bx2 + bw2) - max(ax, bx2)
            oy = min(ay + ah, by2 + bh2) - max(ay, by2)
            if ox > 0 and oy > 0:
                tot += ox * oy
        return tot

    bubble_rects = {i: (bx - r, by - r, 2 * r, 2 * r) for i, bx, by, r, _ in geo}
    taken = []
    lh = 0.25
    for i, bx, by, r, rank in geo:
        label = bubbles[i].get("label", "")
        # Size the box to the TEXT (~0.095in per char at _SZ_SMALL bold). A fixed-width box reported
        # collisions that were not real and packed the cluster far looser than it needed to be.
        wide = max(0.45, min(1.7, 0.095 * len(label) + 0.1))
        narrow, step = wide, lh + 0.02
        cands = []
        for k in range(6):                              # widen the search outward in rings
            cands.append((bx - wide / 2, by - r - lh - k * step, wide, PP_ALIGN.CENTER))
            cands.append((bx - wide / 2, by + r + 0.03 + k * step, wide, PP_ALIGN.CENTER))
            dy = 0 if k == 0 else (-step * k if k % 2 else step * k)
            cands.append((bx + r + 0.07, by - lh / 2 + dy, narrow, PP_ALIGN.LEFT))
            cands.append((bx - r - 0.07 - narrow, by - lh / 2 + dy, narrow, PP_ALIGN.RIGHT))
        others = [v for k, v in bubble_rects.items() if k != i]
        pad = 0.07                                      # keep a gap: touching boxes read as one word
        pick, best = None, None
        for lx, ly, lwd, al in cands:
            lx = min(max(lx, _MARGIN), 13.333 - _MARGIN - lwd)
            ly = min(max(ly, py0 + 0.02), py1 - lh)
            area = _ov_area((lx - pad, ly - 0.02, lwd + 2 * pad, lh + 0.04), taken + others)
            if area <= 0:
                pick = (lx, ly, lwd, al)
                break
            if best is None or area < best[0]:
                best = (area, lx, ly, lwd, al)
        if pick is None:                                # no free slot — take the least-overlapping one
            pick = best[1:]
        lx, ly, lwd, al = pick
        taken.append((lx, ly, lwd, lh))
        # A label pushed clear of the cluster needs a leader, or it reads as belonging to a neighbour.
        lcx = lx + (0.12 if al == PP_ALIGN.LEFT else (lwd - 0.12 if al == PP_ALIGN.RIGHT else lwd / 2))
        lcy = ly + lh / 2
        d = math.hypot(lcx - bx, lcy - by)
        if d > r + 0.42:
            ux, uy = (lcx - bx) / d, (lcy - by) / d
            ln = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT, Inches(bx + ux * (r + 0.02)), Inches(by + uy * (r + 0.02)),
                Inches(lcx - ux * 0.06), Inches(lcy - uy * 0.06))
            ln.line.color.rgb = _LTEAL; ln.line.width = Pt(0.75); ln.shadow.inherit = False
        _place_text(slide, lx, ly, lwd, lh, label, _SZ_SMALL, _INKC, bold=True, font=_HEAD, align=al)
    if takeaways:
        _place_bullets(slide, right_x, py0, right_w, bottom - py0, takeaways, _SZ_SMALL, _INKC)
    if note:
        _place_text(slide, _MARGIN, _BODY_BOTTOM - 0.3, _CONTENT_W, 0.28, note, _SZ_SMALL, _TEAL2)


def _slide_has_white_bg(slide) -> bool:
    cSld = slide._element.find(qn("p:cSld"))
    bg = cSld.find(qn("p:bg")) if cSld is not None else None
    return bg is not None and "FFFFFF" in (bg.xml or "")


def _stamp_footer(slide, n: int) -> None:
    """The one centred footer line on every slide (the template carries none): page number by
    default, plus the About page's optional standing footer text and/or render-date stamp,
    joined with middle dots. Coloured for the slide's background; small tier — part of the
    3-size scale, not an exception. Draws nothing when the team turned everything off."""
    parts = []
    if _FOOTER_TEXT:
        parts.append(_FOOTER_TEXT)
    if _DATE_STAMP:
        parts.append(datetime.date.today().strftime("%d %b %Y"))
    if _PAGE_NUMBERS:
        parts.append(str(n))
    if not parts:
        return
    text = "  ·  ".join(parts)
    # _slide_has_white_bg only sees a background the RENDERER set on the slide. A light-only
    # brand's pale background comes from its master instead, so without the first test the footer
    # picks the pale tint and the page number is invisible on it (and fails that brand guide's own
    # contrast matrix).
    color = _TEAL if (_LIGHT_ONLY or _slide_has_white_bg(slide)) else _LTEAL
    w = 9.0 if len(parts) > 1 else 1.0   # a lone number keeps its original narrow box
    tb = slide.shapes.add_textbox(Inches((13.333 - w) / 2), Inches(7.06), Inches(w), Inches(0.3))
    tf = tb.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.size = Pt(_SZ_SMALL)
    r.font.name = _BODY
    r.font.color.rgb = color


_APPENDIX_MAX_PER_STUDY = 4   # a rich paper can have a dozen+ figures/tables; cap so the appendix
_APPENDIX_MAX_TOTAL = 20      # stays a reference section, not a second deck's worth of slides


def _add_appendix_slides(prs, master_index: int, study_meta: list[dict] | None) -> None:
    """Splice one slide per extracted chart/table from the studies actually cited in this deck -
    the reviewer's own source evidence, spliced in verbatim (not planner-authored, so there is no
    hallucination risk on a scientific figure). No-op if none of the cited studies have any
    extracted figures/tables (assets/figures/, built by scripts/extract_figures.py)."""
    if not study_meta:
        return
    index = config.figures_index(_BRAND)
    if not index:
        return
    entries = []  # (pmid, cite, entry)
    for m in study_meta:
        pmid = str(m.get("pmid") or "").strip()
        if not pmid:
            continue
        cite = (m.get("cite") or f"PMID {pmid}").strip()
        for e in index.get(pmid, [])[:_APPENDIX_MAX_PER_STUDY]:
            entries.append((pmid, cite, e))
    entries = entries[:_APPENDIX_MAX_TOTAL]
    if not entries:
        return

    divider = _synth_slide(prs, master_index, title="Appendix",
                           eyebrow="Source charts and tables from the cited studies")
    divider.notes_slide.notes_text_frame.text = (
        "Appendix: the original charts and tables from the cited studies, reproduced unmodified "
        "for reference.")
    for pmid, cite, e in entries:
        path = config.figure_path(pmid, e["file"], _BRAND)
        if not path.exists():
            continue
        kind_label = "Table" if e.get("kind") == "table" else "Figure"
        slide = _synth_slide(prs, master_index, title="Appendix",
                              eyebrow=_fit(f"{cite} — {kind_label}, page {e['page']}", 100))
        box = (Inches(_MARGIN), Inches(_BODY_TOP), Inches(_CONTENT_W), Inches(_BODY_H))
        _place_icon(slide, box, path)  # letterbox-fit — a chart must never be crop-to-filled
        # These deterministic slides get a deterministic note: the exact source of the image.
        slide.notes_slide.notes_text_frame.text = (
            f"{kind_label} reproduced unmodified from the source study: {cite}, page {e['page']}.")


def _make_slide(prs, spec: dict, catalog: dict, dark: int, light: int,
                custom_by_key: dict, placed_custom: set, unnumbered: set,
                overrides_by_key: dict | None = None) -> None:
    """Add ONE slide for a plan spec — the layout dispatch, extracted from render_deck so the
    loop can apply per-slide extras (speaker notes) to whatever slide any branch added."""
    layout_name = spec["layout"]
    if _LIGHT_ONLY and not layout_name.startswith("custom_") and spec.get("background") != "light":
        # A brand with no dark master renders everything on the light path, whatever the plan
        # asked for. Done here, once, rather than in each of the 26 layouts that decide their own
        # `light` flag. Copied, not mutated: the caller's plan is not ours to rewrite. Verbatim
        # `custom_*` slides are exempt — they carry their own background and no light/dark concept.
        spec = {**spec, "background": "light"}
    if layout_name.startswith("custom_"):   # a team slide the planner placed in the storyline
        c = custom_by_key.get(layout_name)
        if c:
            if c.get("slots"):
                # A team slide the team asked us to FILL: same splice as a redesigned built-in
                # layout, with the plan's per-slot text written into the design's own boxes.
                _add_override_slide(prs, dark, {**c, "layout": layout_name}, spec, unnumbered)
            else:
                _add_custom_slide(prs, dark, c["bytes"], c["index"], c.get("png"), unnumbered)
            placed_custom.add(layout_name)
        else:
            print(f"[custom-slide] plan references unknown {layout_name}; skipped", file=sys.stderr)
        return
    ov = (overrides_by_key or {}).get(layout_name)
    if ov:   # the team replaced this layout's DESIGN — splice theirs, refill its text slots
        _add_override_slide(prs, dark, ov, spec, unnumbered)
        return
    if layout_name == "ingredient":   # AKBM's standard slide, spliced in verbatim
        _add_ingredient_slide(prs, dark)
        return
    if layout_name == "key_points":   # code-built 4-icon-card layout (mechanism B)
        _fill_key_points(prs, spec, light)
        return
    if layout_name == "chart":        # native pptx chart (mechanism B)
        _fill_chart(prs, spec, dark, light)
        return
    if layout_name == "matrix":
        _fill_matrix(prs, spec, dark, light); return
    if layout_name == "exec_summary":
        _fill_exec_summary(prs, spec, dark, light); return
    if layout_name == "comparison":
        _fill_comparison(prs, spec, light); return
    if layout_name == "stat":
        _fill_stat(prs, spec, dark, light); return
    if layout_name == "harvey_ball":
        _fill_harvey_ball(prs, spec, light); return
    if layout_name == "funnel":
        _fill_funnel(prs, spec, dark, light); return
    if layout_name == "closing":
        _fill_closing(prs, spec, dark, light); return
    if layout_name == "kpi_dashboard":
        _fill_kpi_dashboard(prs, spec, dark, light); return
    if layout_name == "roadmap":
        _fill_roadmap(prs, spec, dark, light); return
    if layout_name == "icon_grid":
        _fill_icon_grid(prs, spec, dark, light); return
    if layout_name == "takeaways":
        _fill_takeaways(prs, spec, dark, light); return
    if layout_name == "from_to":
        _fill_from_to(prs, spec, dark, light); return
    if layout_name == "pillars":
        _fill_pillars(prs, spec, dark, light); return
    if layout_name == "team":
        _fill_team(prs, spec, dark, light); return
    if layout_name == "metric_bars":
        _fill_metric_bars(prs, spec, dark, light); return
    if layout_name == "cause_effect":
        _fill_cause_effect(prs, spec, dark, light); return
    if layout_name == "org_chart":
        _fill_org_chart(prs, spec, dark, light); return
    if layout_name == "decision_tree":
        _fill_decision_tree(prs, spec, dark, light); return
    if layout_name == "cycle":
        _fill_cycle(prs, spec, dark, light); return
    if layout_name == "gantt":
        _fill_gantt(prs, spec, dark, light); return
    if layout_name == "serpentine":
        _fill_serpentine(prs, spec, dark, light); return
    if layout_name == "coverage_matrix":
        _fill_coverage_matrix(prs, spec, light); return
    if layout_name == "photo_stats":
        _fill_photo_stats(prs, spec, dark, light); return
    if layout_name == "numbered_cards":
        _fill_numbered_cards(prs, spec, dark, light); return
    if layout_name == "implications":
        _fill_implications(prs, spec, dark, light); return
    if layout_name == "breakdown":
        _fill_breakdown(prs, spec, dark, light); return
    if layout_name == "chart_bands":
        _fill_chart_bands(prs, spec, dark, light); return
    if layout_name == "chart_takeaways":
        _fill_chart_takeaways(prs, spec, light); return
    cat = catalog.get(layout_name)
    if not cat:
        raise ValueError(f"Unknown layout '{layout_name}' (not in catalog)")
    want_light = spec.get("background") in ("light", "pastel") and "light" in cat["backgrounds"]
    master_index = light if want_light else dark
    layout = _find_layout(prs, cat["template_layout"], master_index)
    slide = prs.slides.add_slide(layout)
    _fill_slide(slide, spec, cat, master_index, dark=not want_light)
    if want_light and spec.get("background") == "pastel":
        _set_bg(slide, _LTEAL)


def _render_deck_impl(plan: dict, study_meta: list[dict] | None,
                      design: dict | None,
                      custom_slides: list[dict] | None,
                      custom_photos: list[dict] | None,
                      layout_overrides: list[dict] | None,
                      benefits_slot: str | None,
                      source_appendix: bool,
                      return_slide_map: bool,
                      brand: str | None = None) -> bytes | tuple[bytes, list[int | None]]:
    """Unguarded implementation of render_deck. Called from render_deck within a lock.

    benefits_slot: where AKBM's verbatim benefits overview goes ("first"/"second"/"third"/
    "second_to_last"/"last"), or None to leave it out. source_appendix: append the picked studies'
    own charts and tables. Both were unconditional and invisible until the team's structure rules
    could decide them."""
    apply_brand(brand)        # palette + typography for this brand; must precede apply_design
    apply_design(design)      # ...which resets to the brand defaults apply_brand just set
    register_custom_photos(custom_photos)
    custom_by_key = {c["key"]: c for c in (custom_slides or [])}
    overrides_by_key = {o["layout"]: o for o in (layout_overrides or [])}
    placed_custom: set[str] = set()
    unnumbered: set[int] = set()   # slide ids that carry their own baked page number

    prs = Presentation(str(config.template_path(_BRAND)))
    _delete_example_slides(prs)
    catalog = config.catalog(_BRAND)
    dark, light = _master_indices(_BRAND)

    owners: list[int | None] = []  # parallel to prs.slides, kept in sync through every reorder
    for plan_idx, spec in enumerate(plan["slides"]):
        before = len(prs.slides._sldIdLst)
        _make_slide(prs, spec, catalog, dark, light, custom_by_key, placed_custom, unnumbered,
                    overrides_by_key)
        added = len(prs.slides._sldIdLst) - before
        owners.extend([plan_idx] * added)
        # Speaker notes are written HERE, on whatever slide the dispatch just added, so every
        # layout gets them — code-built, native template and verbatim splices alike (an unknown
        # team-slide key adds nothing, hence the count guard). python-pptx creates the notes
        # slide on first access.
        notes = (spec.get("speaker_notes") or "").strip()
        if notes and added:
            prs.slides[-1].notes_slide.notes_text_frame.text = notes

    # Team slides marked "in every deck" that the plan didn't already place — appended after the
    # content, before the benefits overview below, so they read as part of the deck's fixed tail.
    for c in (custom_slides or []):
        if c.get("mode") == "always" and c["key"] not in placed_custom:
            before = len(prs.slides._sldIdLst)
            if c.get("slots"):
                # An "always" slide is never shown to the planner, so there is no per-slot text
                # to write — the design keeps its own words (_refill_slots leaves a slot the plan
                # didn't fill untouched). The About page steers AI-filled designs to "AI decides"
                # for exactly this reason.
                _add_override_slide(prs, dark, {**c, "layout": c["key"]}, {}, unnumbered)
            else:
                _add_custom_slide(prs, dark, c["bytes"], c["index"], c.get("png"), unnumbered)
            owners.extend([None] * (len(prs.slides._sldIdLst) - before))

    # AKBM's standard "Proven Health Benefits" overview, spliced in verbatim as the second-to-last
    # slide of every deck (appended, then moved into place). It is a SUPERBA slide carrying
    # Superba's own trial counts, so a brand whose template has no equivalent skips it — without
    # this the splice raises and takes the whole deck down.
    if benefits_slot and _HAS_BENEFITS_SLIDE:
        _add_benefits_slide(prs, light)
        owners.append(None)
        sldIdLst = prs.slides._sldIdLst
        benefits = list(sldIdLst)[-1]
        n = len(sldIdLst)
        at = {"first": 0, "second": min(1, n - 1), "third": min(2, n - 1),
              "last": n - 1}.get(benefits_slot, max(1, n - 1))   # default: second to last
        sldIdLst.remove(benefits)
        sldIdLst.insert(at, benefits)
        ben_owner = owners.pop()                          # mirrors sldIdLst.remove(benefits) —
        owners.insert(at, ben_owner)                      # benefits was owners' own last entry

    # The reviewer's own source charts/tables, appended after everything else (added here, so it
    # naturally lands after the just-reordered benefits slide, and picks up page numbers below like
    # any other slide). Dark master, like every other synthetic content slide — the white margin
    # baked into each extracted image already gives it a clean card-like frame against that background.
    if source_appendix:
        before = len(prs.slides._sldIdLst)
        _add_appendix_slides(prs, dark, study_meta)
        owners.extend([None] * (len(prs.slides._sldIdLst) - before))

    # Footer line (page number / footer text / date) in a fixed position on every slide (cover
    # excluded; team slides carry their own baked chrome and would double-print), in final order.
    for i, slide in enumerate(prs.slides):
        if i == 0 or slide.slide_id in unnumbered:
            continue
        _stamp_footer(slide, i + 1)

    buf = io.BytesIO()
    prs.save(buf)
    data = buf.getvalue()
    return (data, owners) if return_slide_map else data


def render_deck(plan: dict, study_meta: list[dict] | None = None,
                design: dict | None = None,
                custom_slides: list[dict] | None = None,
                custom_photos: list[dict] | None = None,
                layout_overrides: list[dict] | None = None,
                benefits_slot: str | None = "second_to_last",
                source_appendix: bool = True,
                return_slide_map: bool = False,
                brand: str | None = None) -> bytes | tuple[bytes, list[int | None]]:
    """Serialize render access: design, photos, icons, the brand palette and other rendering
    globals are mutated at the start of each render. Without synchronization, concurrent
    renders stomp each other's settings — including each other's BRAND, now that one process
    serves several. The lock gates the entire render operation."""
    with _RENDER_LOCK:
        return _render_deck_impl(plan, study_meta, design, custom_slides, custom_photos,
                                 layout_overrides, benefits_slot, source_appendix,
                                 return_slide_map, brand)
