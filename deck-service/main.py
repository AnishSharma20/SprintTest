"""Superba Deck Generator — HTTP service (FastAPI).

Thin layer over the two-stage pipeline in `src`: Claude plans a schema-validated slide
plan, python-pptx fills the real Superba template (all design inherited). One summary
returns a .pptx; several return a .zip that also bundles each deck's wording-review doc.

  POST /jobs             multipart "filer" (1+ summaries) -> {job_id}; runs in background
  GET  /jobs/{id}        -> {status, progress, step, filename, error}
  GET  /jobs/{id}/result -> the .pptx (or .zip) once done
  POST /generate         synchronous single request (legacy convenience)
  GET  /health           readiness probe

ANTHROPIC_API_KEY is read from the environment, server-side only.
"""
from __future__ import annotations

import base64
import io
import json
import os
import re
import subprocess
import sys
import tempfile
import threading
import time
import unicodedata
import traceback
import uuid
import zipfile
from pathlib import Path

import anthropic
from fastapi import FastAPI, File, Form, Header, UploadFile
from fastapi.responses import JSONResponse, Response

import src
from src import brand as brand_mod
from src import config, renderer

app = FastAPI(title="Superba Deck Generator")

PPTX_MEDIA = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
DOCX_MEDIA = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"

# Member names inside the delivered zip. Deliberately SHORT and FIXED: Windows Explorer names the
# extraction folder after the zip, so repeating the topic in every member name put it in the path
# twice and tipped Explorer into "the target path is too long" even well under its 260 character
# limit. The folder already says what the document is about, and the designer saves under their own
# name anyway.
_IDML_MEMBER = "whitepaper.idml"
_PREVIEW_MEMBER = "whitepaper.preview.md"


def _mix_readme(pages: list[str], rationale: str, images: list[str],
                photos: dict | None = None) -> str:
    """Designer note for a COMPOSED whitepaper. Images are linked, and a composed document can pull
    pages from more than one brochure, so it may need Links from several of the original packages —
    which the design team already has. We cannot ship them (they run to hundreds of MB)."""
    return (
        "Superba whitepaper, composed from several Superba brochures\n"
        "===========================================================\n\n"
        f"Pages used, in order:\n" + "".join(f"  {i}. {p}\n" for i, p in enumerate(pages, 1)) +
        (f"\nWhy these pages: {rationale}\n" if rationale else "") +
        "\nEach page keeps the exact design it was drawn with. Pages marked verbatim in the tool\n"
        "(the benefit grid, the ingredient and portfolio spread) are placed unchanged because their\n"
        "icons and figures are brand facts.\n\n"
        "TO OPEN\n"
        "  1. Keep the Links folder next to the .idml, then open the .idml in InDesign\n"
        "     (File > Open). Every image this document uses is in that folder and the links point\n"
        "     at it, so the pictures come with the document.\n"
        "  2. If InDesign still reports missing links: an .idml opens as an UNTITLED document, so\n"
        "     save it once (File > Save As, into this same folder), then Window > Links > panel\n"
        "     menu > Relink to Folder and choose Links. That resolves all of them in one go.\n"
        "  3. Fonts: Manrope, Exo 2, Montserrat.\n"
        "  4. Review every page, then export to PDF.\n\n"
        "AI GENERATED DRAFT. Review all content, claims and figures before any use.\n\n"
        + _photo_note(photos) + _links_note(photos)
    )


def _links_note(photos: dict | None) -> str:
    """What is in the Links folder, and the few assets we genuinely cannot supply."""
    if not photos:
        return ""
    bundled = photos.get("bundle") or []
    unresolved = photos.get("unresolved") or []
    out = f"IMAGES INCLUDED IN Links ({len(bundled)}):\n" + "".join(f"  {n}\n" for n in bundled)
    if unresolved:
        out += (f"\nNOT INCLUDED ({len(unresolved)}). These are missing from the source package the\n"
                "design team supplied, so we have no file to send. They will show as missing links;\n"
                "please drop the originals into Links, or delete the frames:\n"
                + "".join(f"  {n}\n" for n in unresolved))
    return out


def _photo_note(photos: dict | None) -> str:
    """Tell the designer which photographs were substituted and which were deliberately kept.

    A frame keeps its original photograph when our library cannot fill it at print resolution, and
    saying so beats leaving the designer to wonder why some pictures changed and others did not.
    """
    if not photos:
        return ""
    lines = [f"PHOTOGRAPHS (subject: {photos.get('theme') or 'kept as designed'})\n"]
    for r in photos.get("replaced", []):
        lines.append(f"  replaced on {r['page']}: {r['file']} (bundled in Links)\n")
    for k in photos.get("kept", []):
        lines.append(f"  kept the designed photo on {k['page']}: {k['reason']}\n")
    return "".join(lines) + "\n" if len(lines) > 1 else ""


def _pptx_shape_lines(shape) -> list[str]:
    """Every piece of human-readable text one shape carries: text frames, table cells (one line
    per row, cells joined with a separator), chart data (best-effort: series names, categories and
    values), and grouped shapes recursively."""
    lines: list[str] = []
    if getattr(shape, "shape_type", None) == 6:  # MSO_SHAPE_TYPE.GROUP
        for sub in shape.shapes:
            lines.extend(_pptx_shape_lines(sub))
        return lines
    if getattr(shape, "has_text_frame", False):
        t = shape.text_frame.text.strip()
        if t:
            lines.append(t)
    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            cells = [(c.text or "").strip() for c in row.cells]
            if any(cells):
                lines.append(" | ".join(cells))
    if getattr(shape, "has_chart", False):
        try:
            chart = shape.chart
            cats = [str(c) for c in chart.plots[0].categories]
            if cats:
                lines.append("Chart categories: " + ", ".join(cats))
            for ser in chart.series:
                vals = ", ".join("" if v is None else f"{v:g}" for v in ser.values)
                lines.append(f"Chart series {ser.name or '?'}: {vals}")
        except Exception:  # noqa: BLE001 — chart internals vary; text extraction must not fail on them
            pass
    return lines


def _pptx_text(name: str, data: bytes) -> str:
    """An existing PowerPoint as source material: all reader-facing text, slide by slide, plus any
    speaker notes — so the tool can regenerate an old/off-brand deck as a proper Superba one. The
    header tells the planner what it is looking at (the old deck's slide split is input, not a
    required structure)."""
    from pptx import Presentation  # python-pptx, already a renderer dependency

    prs = Presentation(io.BytesIO(data))
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        lines: list[str] = []
        for shape in slide.shapes:
            lines.extend(_pptx_shape_lines(shape))
        if slide.has_notes_slide:
            n = slide.notes_slide.notes_text_frame.text.strip()
            if n:
                lines.append(f"Speaker notes: {n}")
        if lines:
            parts.append(f"## Slide {i}\n" + "\n".join(lines))
    if not parts:
        return ""
    return (f"EXISTING PRESENTATION (text extracted from the PowerPoint file \"{name}\"; use its "
            "content as source material — its slide split reflects the OLD deck, so restructure "
            "freely into the best new storyline):\n\n" + "\n\n".join(parts))


def _read_summary(name: str, data: bytes) -> str:
    if name.lower().endswith(".docx"):
        import docx  # python-docx
        return "\n".join(p.text for p in docx.Document(io.BytesIO(data)).paragraphs)
    if name.lower().endswith((".pptx", ".potx")):
        return _pptx_text(name, data)
    return data.decode("utf-8", errors="replace")


# ---------------------------------------------------------------------------
# Async jobs. Generation takes longer than a proxy/gateway will hold a connection,
# so it runs in a background thread reporting progress into an in-memory store; the
# client POSTs to start, polls status, then downloads.
# ---------------------------------------------------------------------------
JOBS: dict[str, dict] = {}
JOB_TTL_SECONDS = 3600

# Concurrency control: semaphore gates the render + QA gate (the heavy-tail part that uses
# lots of memory) to prevent OOM on the free 512 MB tier. Planning (1–3 minutes, network-bound)
# stays fully parallel. The render itself is also locked inside renderer.py to serialize access
# to its module-level globals (design, photos, icons).
_HEAVY_TAIL_SEMAPHORE = threading.BoundedSemaphore(2)

# Temp directory for spilled results (option 2a), cleaned up on _prune_jobs.
_RESULT_TEMP_DIR = tempfile.mkdtemp(prefix="deck_results_")


def _prune_jobs() -> None:
    now = time.time()
    expired = [k for k, v in JOBS.items() if now - v.get("created", now) > JOB_TTL_SECONDS]
    for jid in expired:
        j = JOBS.pop(jid, {})
        # Clean up spilled result files on disk.
        if "result_path" in j:
            try:
                Path(j["result_path"]).unlink(missing_ok=True)
            except Exception:  # noqa: BLE001
                pass


# Letters NFKD cannot fold, because they are distinct letters rather than base + combining mark.
# Without this, a Norwegian title loses them outright: "øker" would slug to "ker".
_TRANSLIT = str.maketrans({
    "ø": "o", "Ø": "O", "æ": "ae", "Æ": "AE", "å": "a", "Å": "A",
    "ß": "ss", "þ": "th", "Þ": "Th", "ð": "d", "Ð": "D",
    "ł": "l", "Ł": "L", "đ": "d", "Đ": "D", "œ": "oe", "Œ": "OE",
})


def _slug(text: str, limit: int = 60) -> str:
    """A filesystem-safe ASCII stem from a GENERATED title, so a download is named after its topic
    instead of the source file (picking studies always produced "Selected-scientific-studies").

    ASCII-only on purpose: Content-Disposition here carries a plain `filename="..."`, and non-ASCII
    there is mangled by some clients. Nordic letters are transliterated and accents folded
    (Omega-3 nivået øker -> Omega-3-nivaet-oker); a title with no usable ASCII at all (e.g. fully
    CJK) yields "" so the caller falls back to the source-file base name.
    """
    s = (text or "").translate(_TRANSLIT)
    s = unicodedata.normalize("NFKD", s)
    s = s.encode("ascii", "ignore").decode("ascii")
    s = re.sub(r"[^A-Za-z0-9]+", "-", s)
    s = re.sub(r"-{2,}", "-", s).strip("-")
    if len(s) <= limit:
        return s
    # Cut on a word boundary. Chopping mid-word ("...-The-Kri") looks like a corrupted download,
    # and the stem also becomes the folder Windows extracts into, so shorter genuinely matters.
    cut = s[:limit]
    boundary = cut.rfind("-")
    return (cut[:boundary] if boundary > limit // 2 else cut).rstrip("-")


def _store_result(data: bytes) -> str:
    """Write result bytes to a temp file and return the path. Files are cleaned up by
    _prune_jobs() when the job expires."""
    p = Path(_RESULT_TEMP_DIR) / f"{uuid.uuid4().hex}.bin"
    p.write_bytes(data)
    return str(p)


def _parse_custom_slides(custom_slides_meta: str,
                         custom_file_blobs: dict[str, bytes]) -> list[dict]:
    """The About page's team slides, matched to their uploaded .pptx blobs. Malformed or
    incomplete entries are dropped silently — a team slide must never fail a whole job."""
    try:
        meta = json.loads(custom_slides_meta) if custom_slides_meta else []
    except Exception:  # noqa: BLE001
        return []
    out = []
    for m in meta if isinstance(meta, list) else []:
        if not isinstance(m, dict):
            continue
        blob = custom_file_blobs.get(str(m.get("file_id") or ""))
        if not blob or not m.get("id"):
            continue
        png = None
        if m.get("preview_b64"):
            try:
                png = base64.b64decode(m["preview_b64"])
            except Exception:  # noqa: BLE001
                png = None
        # slots (optional) mark this as a design the AI FILLS rather than a verbatim slide: the
        # same measured text-slot map a redesigned built-in layout carries. Entries must be
        # well formed or the slide falls back to verbatim — never fail a job over it.
        slots = m.get("slots")
        if not (isinstance(slots, list) and slots
                and all(isinstance(s, dict) and s.get("slot_id") and s.get("char_budget")
                        for s in slots)):
            slots = None
        out.append({"key": f"custom_{m['id']}", "name": str(m.get("name") or "Team slide"),
                    "description": str(m.get("description") or ""),
                    "mode": m.get("mode") if m.get("mode") in ("auto", "always") else "auto",
                    "bytes": blob, "index": int(m.get("slide_index") or 0), "png": png,
                    "slots": slots})
    return out


def _parse_layout_overrides(layout_overrides_meta: str,
                            custom_file_blobs: dict[str, bytes],
                            brand: str | None = None) -> list[dict]:
    """The About page's TEAM REDESIGNED layouts (design overrides), matched to their uploaded
    .pptx blobs — the blobs ride the same custom_files channel as team slides, named
    <file_id>.pptx. Same never-fail contract as _parse_custom_slides: malformed or incomplete
    entries (unknown/fixed-role layout, missing blob, no slots) are dropped silently."""
    try:
        meta = json.loads(layout_overrides_meta) if layout_overrides_meta else []
    except Exception:  # noqa: BLE001
        return []
    from src.overrides import OVERRIDE_EXCLUDED
    try:
        known = set(config.catalog(brand))
    except Exception:  # noqa: BLE001
        known = set()
    out, seen = [], set()
    for m in meta if isinstance(meta, list) else []:
        if not isinstance(m, dict):
            continue
        layout = str(m.get("layout") or "")
        blob = custom_file_blobs.get(str(m.get("file_id") or ""))
        slots = m.get("slots")
        if (not layout or layout in seen or layout in OVERRIDE_EXCLUDED
                or (known and layout not in known) or not blob
                or not isinstance(slots, list) or not slots):
            continue
        if not all(isinstance(s, dict) and s.get("slot_id") and s.get("char_budget")
                   for s in slots):
            continue
        png = None
        if m.get("preview_b64"):
            try:
                png = base64.b64decode(m["preview_b64"])
            except Exception:  # noqa: BLE001
                png = None
        seen.add(layout)
        out.append({"layout": layout, "bytes": blob, "index": int(m.get("slide_index") or 0),
                    "slots": slots, "png": png})
    return out


def _parse_custom_photos(custom_photos_meta: str,
                         photo_blobs: dict[str, bytes]) -> list[dict]:
    """The About page's team photo library, matched to the uploaded image blobs. Same
    never-fail contract as _parse_custom_slides."""
    try:
        meta = json.loads(custom_photos_meta) if custom_photos_meta else []
    except Exception:  # noqa: BLE001
        return []
    out = []
    for m in meta if isinstance(meta, list) else []:
        if not isinstance(m, dict) or not m.get("id"):
            continue
        blob = photo_blobs.get(str(m["id"]))
        if not blob:
            continue
        out.append({"key": f"team_photo_{m['id']}", "name": str(m.get("name") or "Team photo"),
                    "description": str(m.get("description") or ""), "bytes": blob})
    return out


def _run_job(job_id: str, key: str, files: list[tuple[str, bytes]], lengde: str, tone: str,
             kvalitet: str = "polished", instruksjoner: str = "", innholdstype: str = "deck",
             sprak: str = "English", sider: str = "", study_meta: str = "",
             custom_rules: str = "", disabled_layouts: str = "", design_settings: str = "",
             custom_slides_meta: str = "",
             custom_file_blobs: dict[str, bytes] | None = None,
             custom_photos_meta: str = "",
             custom_photo_blobs: dict[str, bytes] | None = None,
             preferred_layouts: str = "",
             disabled_photos: str = "",
             preferred_photos: str = "",
             color_theme: str = "",
             layout_overrides_meta: str = "",
             structure_rules: str = "",
             managed_blocks: str = "",
             brand: str = "") -> None:
    try:
        client = anthropic.Anthropic(api_key=key)

        # Output language is threaded to the planner/blog as a high-priority instruction so the user
        # can pick ANY language; it overrides the "match the source language" default.
        #
        # The rule BRACKETS the user's own text — stated before it AND restated after it as the final
        # word. Putting it only in front lost to recency: an instruction like "a whitepaper for the
        # german audience" made the model write the whole document in German even with English picked,
        # because it read "german audience" last and treated audience as language. The closing rule
        # therefore names that exact confusion and separates WHO reads it from WHAT language it is in.
        lang = (sprak or "").strip()
        if lang:
            user_text = (instruksjoner or "").strip()
            instruksjoner = (
                f"OUTPUT LANGUAGE: {lang}. Write ALL reader-facing text in {lang}, regardless of the "
                f"language of the source material.\n\n"
                "USER INSTRUCTIONS (follow these, except that they can NEVER change the output "
                "language):\n\"\"\"\n" + (user_text or "(none)") + "\n\"\"\"\n\n"
                f"FINAL LANGUAGE RULE — this overrides everything above. Write ALL reader-facing text "
                f"in {lang}. If the instructions mention a country, market, region, nationality or "
                f"audience (for example a German audience, the DACH market, Japanese buyers), that "
                f"tells you WHO will read it and WHAT to emphasise for them — it does NOT change the "
                f"language. Do NOT translate the output into that audience's language. The ONLY thing "
                f"that sets the language is this rule: {lang}. Keep brand names (Superba, Aker "
                f"BioMarine) and study citations intact."
            )

        if innholdstype in ("blog", "whitepaper_mix"):
            # One long-form asset from ALL sources combined (files + picked study summaries).
            parts = [t for (fname, data) in files if (t := _read_summary(fname, data).strip())]
            source = "\n\n".join(parts)
            if not source:
                raise ValueError("No text found in the provided files/studies.")
            base = (files[0][0] if files else innholdstype).rsplit(".", 1)[0] or innholdstype

            if innholdstype == "whitepaper_mix":
                # Assembled from designed pages across SEVERAL Superba brochures, then filled.
                b = src.generate_whitepaper_composed(
                    client, source, base, length=lengde, tone=tone, instructions=instruksjoner,
                    pages=[p for p in (sider or "").split(",") if p.strip()],
                    on_progress=lambda p, s: JOBS[job_id].update(progress=p, step=s))
                stem = _slug(b.get("title", ""), 24) or base
                zbuf = io.BytesIO()
                with zipfile.ZipFile(zbuf, "w", zipfile.ZIP_DEFLATED) as z:
                    z.writestr(_IDML_MEMBER, b["idml"])
                    z.writestr(_PREVIEW_MEMBER, b["markdown"])
                    z.writestr("OPEN_IN_INDESIGN.txt",
                               _mix_readme(b["pages"], b["rationale"], b["images"],
                                           b.get("photos")))
                    # EVERY image the document links travels with it, not just the photos we
                    # swapped. Shipping only the swapped ones left the recipient with 11 of 16
                    # frames empty, and the links still pointed at a designer's own machine.
                    from src.idml_images import shipped_asset
                    for name in (b.get("photos") or {}).get("bundle", []):
                        asset = shipped_asset(name)
                        if asset:
                            z.writestr(f"Links/{name}", asset.read_bytes())
                data = zbuf.getvalue()
                JOBS[job_id].update(status="done", progress=100, step="Done",
                                    result_path=_store_result(data), media_type="application/zip",
                                    filename=stem + ".zip")
                return

            b = src.generate_blog(client, source, base, length=lengde, tone=tone, instructions=instruksjoner,
                                  on_progress=lambda p, s: JOBS[job_id].update(progress=p, step=s))
            stem = _slug(b.get("title", "")) or base
            data = b["markdown"].encode("utf-8")
            JOBS[job_id].update(status="done", progress=100, step="Done",
                                result_path=_store_result(data),
                                media_type="text/markdown; charset=utf-8", filename=stem + ".md")
            return

        # Only the synthesized "picked studies" file (byggKilder() in the frontend) carries a PMID
        # for each source - a raw uploaded doc has none, so the appendix is only ever built from
        # that file's study_meta, never invented for the others.
        try:
            parsed_study_meta = json.loads(study_meta) if study_meta else []
        except Exception:  # noqa: BLE001 — malformed input is never worth failing the whole job over
            parsed_study_meta = []

        # The About page's deterministic design overrides + the team's own verbatim slides.
        try:
            parsed_design = json.loads(design_settings) if design_settings else None
            if not isinstance(parsed_design, dict):
                parsed_design = None
        except Exception:  # noqa: BLE001
            parsed_design = None
        parsed_custom = _parse_custom_slides(custom_slides_meta, custom_file_blobs or {})
        parsed_photos = _parse_custom_photos(custom_photos_meta, custom_photo_blobs or {})
        parsed_overrides = _parse_layout_overrides(layout_overrides_meta, custom_file_blobs or {}, brand or None)
        # None (field absent) = no answer, keep the built in shape; [] = the team removed every
        # structure rule and means it. Malformed input is treated as no answer.
        parsed_blocks = None
        if managed_blocks.strip():
            try:
                loaded = json.loads(managed_blocks)
                parsed_blocks = loaded if isinstance(loaded, dict) else None
            except Exception:  # noqa: BLE001 — malformed input falls back to the defaults
                parsed_blocks = None
        parsed_structure = None
        if structure_rules.strip():
            try:
                loaded = json.loads(structure_rules)
                parsed_structure = loaded if isinstance(loaded, list) else None
            except Exception:  # noqa: BLE001
                parsed_structure = None
        parsed_preferred = [p.strip() for p in (preferred_layouts or "").split(",") if p.strip()]
        parsed_disabled_photos = [p.strip() for p in (disabled_photos or "").split(",") if p.strip()]
        parsed_preferred_photos = [p.strip() for p in (preferred_photos or "").split(",") if p.strip()]

        # One deck from ALL attached sources combined (uploaded files + picked studies +
        # approved claims), same rule as blog/whitepaper_mix above — a run with several sources
        # used to spin up one independent deck PER file and zip them together, which surprised
        # users expecting "one generation = one deck" (client feedback 2026-08-10).
        parts = [t for (fname, data) in files if (t := _read_summary(fname, data).strip())]
        if not parts:
            raise ValueError("No text found in the provided files/studies.")
        source = "\n\n".join(parts)
        base = files[0][0].rsplit(".", 1)[0] if files else "deck"
        has_study_meta = any(fname == "Selected-scientific-studies.txt" for fname, _ in files)

        # Acquire semaphore before rendering to prevent OOM from concurrent heavy renders.
        # Planning (the slow 1–3 minute part) happens first, outside the semaphore, so
        # multiple users' planning still overlaps; only the render/QA gates queue.
        with _HEAVY_TAIL_SEMAPHORE:
            d = src.generate(client, source, base, length=lengde, tone=tone,
                             quality=kvalitet, instructions=instruksjoner,
                             on_progress=lambda p, s: JOBS[job_id].update(progress=p, step=s),
                             study_meta=parsed_study_meta if has_study_meta else None,
                             custom_rules=custom_rules,
                             disabled_layouts=[dl.strip() for dl in disabled_layouts.split(",") if dl.strip()],
                             design=parsed_design, custom_slides=parsed_custom,
                             custom_photos=parsed_photos,
                             preferred_layouts=parsed_preferred,
                             disabled_photos=parsed_disabled_photos,
                             preferred_photos=parsed_preferred_photos,
                             color_theme=color_theme.strip() or None,
                             layout_overrides=parsed_overrides,
                             brand=brand or None)

        # Named after the deck's own generated deck_title (the topic), falling back to the
        # source file stem when the title yields no usable ASCII.
        stem = _slug((d.get("plan") or {}).get("deck_title", "")) or d["filename"].rsplit(".", 1)[0]
        JOBS[job_id].update(status="done", progress=100, step="Done",
                            result_path=_store_result(d["pptx"]), media_type=PPTX_MEDIA, filename=stem + ".pptx")
    except Exception as e:  # noqa: BLE001 — record the failure for the client to read
        # Log the FULL traceback: the client only ever sees str(e), and a bare message like
        # "'str' object has no attribute 'get'" says nothing about where it came from. The message
        # keeps its exception type for the same reason — it costs nothing and names the shape of
        # the failure.
        traceback.print_exc()
        JOBS[job_id].update(status="error", step="Failed",
                            error=f"{type(e).__name__}: {e}")


def _auth_or_error(x_deck_token):
    expected = os.environ.get("DECK_SERVICE_TOKEN")
    if expected and x_deck_token != expected:
        return JSONResponse({"feil": "Unauthorized."}, status_code=401)
    if not os.environ.get("ANTHROPIC_API_KEY"):
        return JSONResponse({"feil": "Missing ANTHROPIC_API_KEY on the server."}, status_code=500)
    return None


@app.get("/health")
def health():
    try:
        layouts = len(config.catalog())          # the default brand; see /brands
        template = config.template_path().exists()
    except Exception:  # noqa: BLE001
        layouts, template = 0, False
    return {"ok": True, "layouts": layouts, "template": template, "model": config.MODEL}


@app.post("/generate")
async def generate(
    filer: list[UploadFile],
    lengde: str = Form(default="standard"),
    tone: str = Form(default="balansert"),
    x_deck_token: str | None = Header(default=None),
):
    err = _auth_or_error(x_deck_token)
    if err:
        return err
    if not filer:
        return JSONResponse({"feil": "No files uploaded."}, status_code=400)

    client = anthropic.Anthropic(api_key=os.environ["ANTHROPIC_API_KEY"])
    try:
        decks = []
        for i, uf in enumerate(filer):
            data = await uf.read()
            text = _read_summary(uf.filename or f"summary-{i}", data).strip()
            if not text:
                return JSONResponse({"feil": f"No text found in {uf.filename}."}, status_code=400)
            base = (uf.filename or f"deck-{i + 1}").rsplit(".", 1)[0]
            decks.append(src.generate(client, text, base, length=lengde, tone=tone))

        if len(decks) == 1:
            d = decks[0]
            return Response(content=d["pptx"], media_type=PPTX_MEDIA,
                            headers={"Content-Disposition": f'attachment; filename="{d["filename"]}"'})

        zbuf = io.BytesIO()
        with zipfile.ZipFile(zbuf, "w", zipfile.ZIP_DEFLATED) as z:
            for d in decks:
                z.writestr(d["filename"], d["pptx"])
                z.writestr(d["filename"].rsplit(".", 1)[0] + ".wording.md", d["wording_md"])
        return Response(content=zbuf.getvalue(), media_type="application/zip",
                        headers={"Content-Disposition": 'attachment; filename="superba-decks.zip"'})
    except Exception as e:  # noqa: BLE001 — surface a clean error to the client
        return JSONResponse({"feil": f"Generation failed: {e}"}, status_code=500)


@app.post("/jobs")
async def create_job(
    filer: list[UploadFile],
    lengde: str = Form(default="standard"),
    tone: str = Form(default="balansert"),
    # Default flipped to "polished" 2026-08-13: the vision QA pass is the highest value quality
    # lever available (it is the only thing that looks at the actual pixels). Callers can still send
    # "fast" per job, and DECK_QA_GATE=off disables it deploy-wide without a code change.
    kvalitet: str = Form(default="polished"),
    instruksjoner: str = Form(default=""),
    innholdstype: str = Form(default="deck"),
    sprak: str = Form(default="English"),
    sider: str = Form(default=""),
    study_meta: str = Form(default=""),
    brand: str = Form(default=""),
    custom_rules: str = Form(default=""),
    disabled_layouts: str = Form(default=""),
    design_settings: str = Form(default=""),
    custom_slides_meta: str = Form(default=""),
    custom_files: list[UploadFile] | None = File(default=None),
    custom_photos_meta: str = Form(default=""),
    custom_photo_files: list[UploadFile] | None = File(default=None),
    preferred_layouts: str = Form(default=""),
    disabled_photos: str = Form(default=""),
    preferred_photos: str = Form(default=""),
    color_theme: str = Form(default=""),
    layout_overrides_meta: str = Form(default=""),
    structure_rules: str = Form(default=""),
    managed_blocks: str = Form(default=""),
    x_deck_token: str | None = Header(default=None),
):
    """Start a deck-generation job in the background and return its id immediately.

    kvalitet: "fast" (default) or "polished" (adds a visual QA pass — needs a rasteriser on the
    server, i.e. LibreOffice installed; degrades to fast if absent).
    sprak: output language for the generated text (any language; defaults to English).
    study_meta: JSON array of {pmid, cite} for the picked studies, so a deck can append an
    appendix of those studies' real charts/tables (see extract_figures.py); ignored for
    non-deck content types.
    custom_rules: the team's standing generation rules from the tool's About page (newline
    separated text), injected into the deck planner's prompt; deck only.
    disabled_layouts: comma separated layout keys turned OFF on the About page — removed from
    the planner's vocabulary and schema so the model cannot pick them; deck only.
    design_settings: JSON object of deterministic design overrides (fonts, sizes, spacing,
    margins) the renderer enforces; deck only.
    custom_slides_meta + custom_files: the team's own slides — meta is a JSON array of
    {id, file_id, slide_index, name, description, mode, preview_b64?, slots?}; each custom_files
    upload is named <file_id>.pptx. Without `slots` the slide is spliced VERBATIM (the AI writes
    nothing on it); with `slots` it is a design the AI FILLS, exactly like a redesigned built-in
    layout (see layout_overrides_meta); deck only.
    custom_photos_meta + custom_photo_files: the team's photo library — meta is a JSON array of
    {id, name, description}; each custom_photo_files upload is named <id>.jpg; deck only.
    preferred_layouts: comma separated layout keys the team starred as house favourites —
    a soft planner preference among equally fitting layouts; deck only.
    disabled_photos: comma separated BUILT-IN photo ids turned OFF on the About page — removed
    from every asset_id enum so the model cannot pick them; deck only.
    preferred_photos: comma separated photo ids (built-in or team_photo_<id>) starred as house
    favourites — a soft planner preference among equally fitting photos; deck only.
    managed_blocks: JSON object {block_key: text} of the WRITING rules the team now owns (see
    GET /rules/builtin). Absent = every built in default applies; a key missing from a supplied
    object means the team switched that rule off; deck only.
    structure_rules: JSON array of the About page's STRUCTURE rules ({slide, action, position}) —
    which slides every deck must have and where they sit. Empty keeps the built in shape (cover
    first, executive summary second, agenda third); deck only.
    layout_overrides_meta: the About page's TEAM REDESIGNED layouts — a JSON array of
    {layout, file_id, slide_index, slots, preview_b64?}; the .pptx blobs ride the SAME
    custom_files channel, named <file_id>.pptx. The design is spliced verbatim; the planner
    writes fresh per-slot text on every use; deck only."""
    err = _auth_or_error(x_deck_token)
    if err:
        return err
    if not filer:
        return JSONResponse({"feil": "No files uploaded."}, status_code=400)

    # Reject an unknown brand rather than falling back to the default. A silent fallback would
    # hand the user a SUPERBA deck under a Revervia label — wrong logo, wrong colours, wrong
    # product — which is far worse than a clear error. known_brands() is filesystem-derived, so
    # this also catches a brand whose pack has not reached this deploy yet.
    brand = (brand or "").strip().lower()
    if brand and brand not in config.known_brands():
        return JSONResponse(
            {"feil": f"Unknown brand '{brand}'. This service can render: "
                     f"{', '.join(config.known_brands())}."}, status_code=400)

    _prune_jobs()
    files = [((uf.filename or f"summary-{i}"), await uf.read()) for i, uf in enumerate(filer)]
    # Each team-slide upload is named <file_id>.pptx by the frontend; index the blobs by that id.
    custom_blobs: dict[str, bytes] = {}
    for uf in custom_files or []:
        fid = (uf.filename or "").rsplit(".", 1)[0]
        if fid:
            custom_blobs[fid] = await uf.read()
    photo_blobs: dict[str, bytes] = {}
    for uf in custom_photo_files or []:
        pid = (uf.filename or "").rsplit(".", 1)[0]
        if pid:
            photo_blobs[pid] = await uf.read()
    job_id = uuid.uuid4().hex
    JOBS[job_id] = {"status": "running", "progress": 0, "step": "Starting", "created": time.time()}
    key = os.environ["ANTHROPIC_API_KEY"]
    threading.Thread(target=_run_job,
                     args=(job_id, key, files, lengde, tone, kvalitet, instruksjoner, innholdstype,
                           sprak, sider, study_meta, custom_rules, disabled_layouts,
                           design_settings, custom_slides_meta, custom_blobs,
                           custom_photos_meta, photo_blobs, preferred_layouts,
                           disabled_photos, preferred_photos, color_theme,
                           layout_overrides_meta, structure_rules, managed_blocks,
                           brand),
                     daemon=True).start()
    return {"job_id": job_id}


# Fixed sample content for the design preview: one NATIVE text slide (title font/size, body
# font, bullets, line spacing) and one CODE-BUILT card slide (headings, gutter, icon chips) —
# together they show every design setting. Content is fixed; no LLM is involved.
_PREVIEW_PLAN = {
    "deck_title": "Design preview", "language": "en", "slides": [
        {"layout": "text", "title": "This is how your slide titles will look",
         "body": ("Body text renders in your chosen body font and size.\n"
                  "Line spacing and bullets follow your design settings.\n"
                  "The footer below shows page number, footer text and date.")},
        {"layout": "key_points", "title": "Code built slides follow the same settings",
         "banner": "Headings, boxes and gutters restyle too",
         "items": [
             {"heading": "Fonts", "body": "Titles and body\nin your fonts", "icon_generic": "science"},
             {"heading": "Sizes", "body": "Three text sizes\nset by you", "icon_generic": "quality"},
             {"heading": "Spacing", "body": "Line spacing\nand gutters", "icon_generic": "proven"}]},
    ],
}


@app.post("/design/preview")
async def design_preview(payload: dict, x_deck_token: str | None = Header(default=None)):
    """Render the two fixed sample slides with the GIVEN design settings and return them as
    JPEGs, so the About page can show the effect before saving. Deterministic, no LLM."""
    expected = os.environ.get("DECK_SERVICE_TOKEN")
    if expected and x_deck_token != expected:
        return JSONResponse({"feil": "Unauthorized."}, status_code=401)
    try:
        from PIL import Image

        from src import qa_gate, renderer
        design = payload.get("settings") if isinstance(payload, dict) else None
        if not isinstance(design, dict):
            design = None
        pptx = renderer.render_deck(_PREVIEW_PLAN, design=design)
        # Same LibreOffice rasteriser a full deck generation uses — must share its OOM-guarding
        # semaphore or this can race a concurrent generation job's own render for memory on the
        # 512 MB instance (see _HEAVY_TAIL_SEMAPHORE's definition).
        with _HEAVY_TAIL_SEMAPHORE:
            images = qa_gate.rasterize(pptx)
        if not images:
            return JSONResponse({"feil": "No slide renderer is available on the server."},
                                status_code=503)
        # render_deck splices the verbatim benefits slide in SECOND-TO-LAST, so a 2-slide plan
        # renders as [text, benefits, key_points] — return indices 0 and 2.
        picks = [images[0], images[2]] if len(images) >= 3 else images[:1]
        out = []
        for png in picks:
            im = Image.open(io.BytesIO(png)).convert("RGB")
            if im.width > 1100:
                im = im.resize((1100, round(im.height * 1100 / im.width)))
            buf = io.BytesIO()
            im.save(buf, "JPEG", quality=82)
            out.append(base64.b64encode(buf.getvalue()).decode("ascii"))
        return {"slides": out}
    except Exception as e:  # noqa: BLE001 — surface a clean error to the client
        return JSONResponse({"feil": f"Preview failed: {e}"}, status_code=500)


def _inspect_previews(data: bytes) -> list[dict]:
    """Every slide of a .pptx as a JPEG preview (≤1280px wide, base64) — the shared core of the
    sync and job-based inspect endpoints. Raises on failure; callers translate.

    Rasterises with the smaller PREVIEW_MAX_DIM image ceiling: the output is a 1280px JPEG, so
    any embedded image above that is memory the deployed 512 MB instance cannot spare (the
    team's own 42-slide template export holds ~122 MB of decoded bitmaps at the generation
    default, ~67 MB here)."""
    from PIL import Image

    from src import qa_gate
    # Shares the generation path's OOM-guarding semaphore — see _HEAVY_TAIL_SEMAPHORE's definition.
    with _HEAVY_TAIL_SEMAPHORE:
        images = qa_gate.rasterize(data, qa_gate.PREVIEW_MAX_DIM)
    if not images:
        raise RuntimeError("No slide renderer is available on the server.")
    out = []
    for i, png in enumerate(images):
        im = Image.open(io.BytesIO(png)).convert("RGB")
        if im.width > 1280:
            im = im.resize((1280, round(im.height * 1280 / im.width)))
        buf = io.BytesIO()
        im.save(buf, "JPEG", quality=82)
        out.append({"index": i, "preview_b64": base64.b64encode(buf.getvalue()).decode("ascii")})
    return out


@app.get("/rules/builtin")
def rules_builtin(brand: str = "", x_deck_token: str | None = Header(default=None)):
    """The WRITING rules the team can see and edit, with the exact text the planner uses today.
    The About page seeds these as ordinary rule rows on first view, so what used to be buried in
    the prompt becomes something the team can reword or switch off. Read only: this endpoint never
    changes anything, it just says what the built in defaults are."""
    expected = os.environ.get("DECK_SERVICE_TOKEN")
    if expected and x_deck_token != expected:
        return JSONResponse({"feil": "Unauthorized."}, status_code=401)
    from src.planner import builtin_blocks
    b = (brand or "").strip().lower()
    if b and b not in config.known_brands():
        return JSONResponse({"feil": f"Unknown brand '{b}'."}, status_code=400)
    return {"blocks": builtin_blocks(b or None)}


@app.post("/slides/inspect")
async def slides_inspect(
    file: UploadFile,
    x_deck_token: str | None = Header(default=None),
):
    """Render every slide of an uploaded .pptx to a preview image, so the About page can let the
    user pick which slides to add as team slides. Pure rasterisation, no LLM — LibreOffice on the
    deployed service, PowerPoint COM in local dev. Returns {slides: [{index, preview_b64}]}
    (JPEG, ≤1280px wide) — the same previews are stored and reused as gallery thumbnails and as
    the pixel-perfect fallback when a slide can't be shape-spliced.

    Called server-to-server from app/api/custom-slides/inspect/route.ts (same X-Deck-Token as
    every other endpoint here) — the browser never calls this directly. A large .pptx avoids
    Vercel's ~4.5 MB serverless body ceiling by living in Supabase Storage first: the route
    downloads it there and forwards it here, so what transits this endpoint's own caller (Vercel)
    is a server-to-server hop, never subject to that browser-facing limit.

    SYNCHRONOUS — fine for a handful of slides, but a many-slide deck (e.g. the full 42-slide
    template export re-uploaded after editing) renders longer than Vercel will hold the proxy
    route open (60s), killing the request mid-render. Use /slides/inspect-job for those."""
    expected = os.environ.get("DECK_SERVICE_TOKEN")
    if expected and x_deck_token != expected:
        return JSONResponse({"feil": "Unauthorized."}, status_code=401)
    data = await file.read()
    if not data:
        return JSONResponse({"feil": "Empty file."}, status_code=400)
    try:
        return {"slides": _inspect_previews(data)}
    except RuntimeError as e:
        return JSONResponse({"feil": str(e)}, status_code=503)
    except Exception as e:  # noqa: BLE001 — surface a clean error to the client
        return JSONResponse({"feil": f"Could not read the presentation: {e}"}, status_code=500)


@app.post("/slides/inspect-job")
async def slides_inspect_job(
    file: UploadFile,
    x_deck_token: str | None = Header(default=None),
):
    """Job-based twin of /slides/inspect: returns {job_id} immediately and renders the previews
    in the background — poll GET /jobs/{id}, then GET /jobs/{id}/result for the same
    {slides: [{index, preview_b64}]} JSON. Exists because rasterising a many-slide upload (the
    full template export is 42 slides) takes minutes on the deployed instance, far past the 60s
    a Vercel proxy route can hold a synchronous request open (seen in production: the route died
    with Vercel's plain-text timeout page while LibreOffice was still rendering)."""
    expected = os.environ.get("DECK_SERVICE_TOKEN")
    if expected and x_deck_token != expected:
        return JSONResponse({"feil": "Unauthorized."}, status_code=401)
    data = await file.read()
    if not data:
        return JSONResponse({"feil": "Empty file."}, status_code=400)

    _prune_jobs()
    job_id = uuid.uuid4().hex
    JOBS[job_id] = {"status": "running", "progress": 5, "step": "Rendering slide previews",
                    "created": time.time()}

    def run() -> None:
        try:
            slides = _inspect_previews_chunked(data, job_id)
            JOBS[job_id].update(status="done", progress=100, step="Done",
                                result=json.dumps({"slides": slides}).encode("utf-8"),
                                media_type="application/json", filename="previews.json")
        except Exception as e:  # noqa: BLE001 — record the failure for the client to read
            JOBS[job_id].update(status="error", step="Failed", error=str(e))

    threading.Thread(target=run, daemon=True).start()
    return {"job_id": job_id}


_INSPECT_CHUNK = 8  # slides per LibreOffice conversion — small enough for real progress steps,
                    # large enough that the per-invocation startup overhead stays a minor cost

# Relationship types python-pptx resolves to MSO_SHAPE_TYPE.MEDIA — a deck carrying one must
# never be trimmed/resaved (see qa_gate._extract_image_previews for that investigation).
_VIDEO_REL_MARKERS = ("relationships/video", "relationships/media", "relationships/audio")
_VIDEO_MEDIA_EXT = (".mp4", ".mov", ".avi", ".wmv", ".m4v", ".mpg", ".mpeg", ".mkv", ".webm")


def _deck_shape(data: bytes) -> tuple[int, bool]:
    """(slide count, carries video/audio) read straight from the .pptx zip — deliberately NOT via
    python-pptx: parsing a 42-slide package costs tens of MB that lxml does not hand back, and on
    a 512 MB instance this function runs right before LibreOffice needs all the room it can get.
    Falls back to the python-pptx parse if anything about the zip is unexpected."""
    try:
        z = zipfile.ZipFile(io.BytesIO(data))
        names = z.namelist()
        total = sum(1 for n in names if re.fullmatch(r"ppt/slides/slide\d+\.xml", n))
        if not total:
            raise ValueError("no slide parts found")
        has_media = any(n.lower().endswith(_VIDEO_MEDIA_EXT) for n in names)
        if not has_media:
            for n in names:
                if n.startswith("ppt/slides/_rels/") and n.endswith(".rels"):
                    rels = z.read(n).decode("utf-8", "ignore")
                    if any(m in rels for m in _VIDEO_REL_MARKERS):
                        has_media = True
                        break
        return total, has_media
    except Exception as e:  # noqa: BLE001
        print(f"[inspect-job] zip inspection failed ({e}); using python-pptx", file=sys.stderr)
        from pptx import Presentation

        from src import qa_gate
        prs = Presentation(io.BytesIO(data))
        return len(prs.slides), qa_gate._has_video(prs)


def _inspect_previews_chunked(data: bytes, job_id: str) -> list[dict]:
    """_inspect_previews with REAL progress for big decks: convert a few slides at a time and
    update the job's progress/step after each run, so the About page can show "12 of 42" instead
    of a silent minutes-long wait (LibreOffice's PDF conversion is one monolithic call per file —
    the only way to get progress out of it is to hand it smaller files).

    Each run happens in a CHILD PROCESS (src/preview_chunk.py). Doing it in-process leaked: the
    per-chunk python-pptx parse, PIL bitmaps and PNG buffers are not returned to the OS between
    chunks, so this function peaked at 337 MB and stayed at 324 MB afterwards — on a 512 MB
    instance that leaves nothing for LibreOffice and the container gets killed mid-render (seen
    live). A child process hands all of it back on exit, so peak stays flat per chunk.

    Chunking needs a python-pptx resave of a trimmed copy, so it is skipped for decks with
    embedded video (rasterize() routes those to a read-only image-extraction fallback that never
    resaves — see the qa_gate saga) and for small decks where one conversion is quick anyway.
    Any chunking failure falls back to the plain whole-file render."""
    from src import qa_gate
    try:
        total, has_video = _deck_shape(data)
        if total <= _INSPECT_CHUNK or has_video:
            return _inspect_previews(data)

        # Shrink the oversized embedded images ONCE, up front, and chunk from the result: the
        # masters' media travel with every chunk (a trimmed copy still inherits them), so
        # shrinking per chunk would redo the same work for each one.
        shrunk = qa_gate._shrink_pptx_images(data, qa_gate.PREVIEW_MAX_DIM)

        root = Path(__file__).resolve().parent
        slides: list[dict] = []
        with tempfile.TemporaryDirectory() as tmp:
            deck = Path(tmp) / "deck.pptx"
            deck.write_bytes(shrunk)
            del shrunk  # the child reads it from disk from here on
            for start in range(0, total, _INSPECT_CHUNK):
                end = min(start + _INSPECT_CHUNK, total)
                out_dir = Path(tmp) / f"chunk{start}"
                run = subprocess.run(
                    [sys.executable, "-m", "src.preview_chunk", str(deck), str(start), str(end),
                     str(out_dir)],
                    cwd=str(root), capture_output=True, text=True,
                    timeout=qa_gate._SOFFICE_TIMEOUT_S + 120,
                )
                if run.returncode != 0:
                    detail = (run.stderr or run.stdout or "no output").strip().splitlines()[-1]
                    if "No slide renderer" in detail:
                        raise RuntimeError("No slide renderer is available on the server.")
                    raise ValueError(f"chunk {start}-{end} failed: {detail}")
                jpgs = sorted(out_dir.glob("slide*.jpg"))
                if len(jpgs) != end - start:
                    # ValueError, not RuntimeError: the latter is reserved for "no rasteriser at
                    # all", which must propagate — this one should fall back to a whole-file pass.
                    raise ValueError(f"chunk rendered {len(jpgs)} of {end - start} slides")
                for j, p in enumerate(jpgs):
                    slides.append({"index": start + j,
                                   "preview_b64": base64.b64encode(p.read_bytes()).decode("ascii")})
                    p.unlink(missing_ok=True)
                JOBS[job_id].update(progress=5 + round(90 * end / total),
                                    step=f"Rendering slide previews ({end} of {total})")
        return slides
    except RuntimeError:
        raise  # "no rasteriser available" — chunking again won't help
    except Exception as e:  # noqa: BLE001 — chunking is an optimisation, never a requirement
        print(f"[inspect-job] chunked render failed ({e}); falling back to one pass", file=sys.stderr)
        JOBS[job_id].update(step="Rendering slide previews")
        return _inspect_previews(data)


@app.post("/slides/inspect-slots")
async def slides_inspect_slots(
    file: UploadFile,
    slide_index: int = Form(default=0),
    layout: str = Form(default=""),
    brand: str = Form(default=""),
    x_deck_token: str | None = Header(default=None),
):
    """Measure the AI-refillable TEXT SLOTS of one slide in an uploaded .pptx — the analysis
    half of a layout design override ("TEAM REDESIGNED" layout). For every text-bearing shape:
    a stable slot id (survives the renderer's splice), the text it says now, and a character
    budget measured from its geometry. The About page stores the result in Supabase and ships
    it back with each generation job (layout_overrides_meta) — this service stays stateless.

    Rejects up front (400, human message) anything the recipe path cannot honour: an excluded
    fixed-role layout, a slide with an embedded chart/video/object (the splice cannot carry
    those parts, and the picture fallback would freeze the text), no editable text at all, or
    more slots than the schema/prompt can reasonably carry.

    { file, slide_index, layout? } -> { slide_index, slots: [...], preview_b64 }"""
    expected = os.environ.get("DECK_SERVICE_TOKEN")
    if expected and x_deck_token != expected:
        return JSONResponse({"feil": "Unauthorized."}, status_code=401)
    data = await file.read()
    if not data:
        return JSONResponse({"feil": "Empty file."}, status_code=400)
    from src.overrides import MAX_SLOTS, OVERRIDE_EXCLUDED, extract_slots, slide_ineligible_reason
    if layout:
        try:
            known = set(config.catalog(brand or None))
        except Exception:  # noqa: BLE001
            known = set()
        if known and layout not in known:
            return JSONResponse({"feil": f'Unknown layout "{layout}".'}, status_code=400)
        if layout in OVERRIDE_EXCLUDED:
            return JSONResponse(
                {"feil": "This slide has a fixed structural role in every deck and its design "
                         "can't be replaced."}, status_code=400)
    try:
        from pptx import Presentation

        prs = Presentation(io.BytesIO(data))
        slides = list(prs.slides)
        if not 0 <= slide_index < len(slides):
            return JSONResponse(
                {"feil": f"slide_index {slide_index} out of range ({len(slides)} slides)."},
                status_code=400)
        src_slide = slides[slide_index]
        reason = slide_ineligible_reason(src_slide)
        if reason:
            return JSONResponse({"feil": f"This slide can't be used as a design: {reason}."},
                                status_code=400)
        slots = extract_slots(src_slide)
        if not slots:
            return JSONResponse(
                {"feil": "No editable text found on this slide — the AI would have nothing to "
                         "write into."}, status_code=400)
        if len(slots) > MAX_SLOTS:
            over = len(slots) - MAX_SLOTS
            return JSONResponse(
                {"feil": f"This slide has {len(slots)} text areas — the limit is {MAX_SLOTS}. "
                         f"Merge {over} of them (a label and the text under it can share one box) "
                         f"and try again."}, status_code=400)

        preview_b64 = None
        try:
            from PIL import Image

            from src import qa_gate
            # Rasterise a SINGLE-SLIDE copy, not the whole file: an override upload can be the
            # full 42-slide template export with one slide edited, and rendering all of it takes
            # minutes on the deployed instance (past what the Vercel proxy will wait). Trimming
            # can fail on exotic files (e.g. the python-pptx video-resave quirk) — fall back to
            # the target slide's index in a full render, and a missing preview never blocks.
            raster_input, raster_index = data, slide_index
            if len(slides) > 1:
                try:
                    trimmed = Presentation(io.BytesIO(data))
                    for i in reversed(range(len(trimmed.slides))):
                        if i != slide_index:
                            _drop_slide(trimmed, i)
                    tbuf = io.BytesIO()
                    trimmed.save(tbuf)
                    raster_input, raster_index = tbuf.getvalue(), 0
                except Exception as e:  # noqa: BLE001
                    print(f"[inspect-slots] single-slide trim failed ({e}); rendering the full file",
                          file=sys.stderr)
            # Shares the generation path's OOM-guarding semaphore — see _HEAVY_TAIL_SEMAPHORE's
            # definition; without this an override-recipe preview can race a concurrent deck
            # generation's own LibreOffice render for memory and crash the single worker.
            with _HEAVY_TAIL_SEMAPHORE:
                images = qa_gate.rasterize(raster_input, qa_gate.PREVIEW_MAX_DIM)
            if images and raster_index < len(images):
                im = Image.open(io.BytesIO(images[raster_index])).convert("RGB")
                if im.width > 1280:
                    im = im.resize((1280, round(im.height * 1280 / im.width)))
                buf = io.BytesIO()
                im.save(buf, "JPEG", quality=82)
                preview_b64 = base64.b64encode(buf.getvalue()).decode("ascii")
        except Exception as e:  # noqa: BLE001 — a missing preview must not block the analysis
            print(f"[inspect-slots] preview rasterisation failed: {e}", file=sys.stderr)
        return {"slide_index": slide_index, "slots": slots, "preview_b64": preview_b64}
    except Exception as e:  # noqa: BLE001 — surface a clean error to the client
        return JSONResponse({"feil": f"Could not read the presentation: {e}"}, status_code=500)


_GALLERY_SAMPLES: dict[str, dict] | None = None


def _gallery_samples() -> dict[str, dict]:
    """Lazily load the same one-slide-per-layout sample content the About page's static preview
    gallery is built from (scripts/build_gallery.py), keyed by layout name — lets a standard
    layout be exported as a real, editable .pptx on demand instead of only ever being a flat PNG."""
    global _GALLERY_SAMPLES
    if _GALLERY_SAMPLES is None:
        import sys
        from pathlib import Path

        scripts_dir = Path(__file__).resolve().parent / "scripts"
        if str(scripts_dir) not in sys.path:
            sys.path.insert(0, str(scripts_dir))
        # build_gallery.py reads sys.argv[1] as an optional output dir at IMPORT time (for its own
        # CLI use as `python scripts/build_gallery.py [outdir]`) and mkdir's it immediately. Inside
        # a live server process argv[1] is whatever the server was launched with (e.g. uvicorn's
        # "main:app"), which isn't a valid path — hide argv for the duration of the import so that
        # branch takes its harmless default (ROOT / "build") instead of crashing on it.
        saved_argv, sys.argv = sys.argv, sys.argv[:1]
        try:
            from build_gallery import SYNTH, TMPL  # type: ignore
        finally:
            sys.argv = saved_argv

        title = {"layout": "title", "title": "Superba by Aker BioMarine",
                 "subtitle": "Science backed krill oil"}
        _GALLERY_SAMPLES = {s["layout"]: s for s in [title] + SYNTH + TMPL}
    return _GALLERY_SAMPLES


def _drop_slide(prs, index: int) -> None:
    """Remove one slide (and its relationship) from an already-built Presentation. python-pptx
    exposes no public API for this; mirrors renderer._delete_example_slides."""
    from pptx.oxml.ns import qn

    lst = prs.slides._sldIdLst
    sld_id = list(lst)[index]
    r_id = sld_id.get(qn("r:id"))
    if r_id:
        prs.part.drop_rel(r_id)
    lst.remove(sld_id)


@app.post("/slides/export")
async def slides_export(payload: dict, x_deck_token: str | None = Header(default=None)):
    """Export ONE standard layout as a real, editable single-slide .pptx — half of the "view it,
    then make it yours" round trip: download this file, edit it freely in PowerPoint, then hand
    the edited file back through the team-slides upload flow to save it as your own version.

    { layout, background? ("dark" | "light" | "pastel") } -> the .pptx bytes.
    """
    expected = os.environ.get("DECK_SERVICE_TOKEN")
    if expected and x_deck_token != expected:
        return JSONResponse({"feil": "Unauthorized."}, status_code=401)
    layout = (payload or {}).get("layout")
    background = (payload or {}).get("background") or "dark"
    brand = ((payload or {}).get("brand") or "").strip().lower() or None
    if brand and brand not in config.known_brands():
        return JSONResponse({"feil": f"Unknown brand '{brand}'."}, status_code=400)
    if not isinstance(layout, str) or not layout:
        return JSONResponse({"feil": "Missing layout key."}, status_code=400)
    if layout == "benefits_verbatim":
        return JSONResponse(
            {"feil": "The Proven Health Benefits slide is a fixed brand asset and can't be exported."},
            status_code=400,
        )
    try:
        from pptx import Presentation

        from src import renderer
        sample = _gallery_samples().get(layout)
        if sample is None:
            return JSONResponse({"feil": f'Unknown layout "{layout}".'}, status_code=404)
        from build_gallery import localize_sample     # same substitutions as the preview gallery
        slide = localize_sample(sample, brand)
        if background != "dark":
            slide["background"] = background
        data = renderer.render_deck({"deck_title": "Sample", "language": "en", "slides": [slide]},
                                    brand=brand)
        # A single-slide plan renders as [that slide, the verbatim benefits overview] for a brand
        # whose template HAS that slide — drop the trailing one so the download is exactly the
        # layout asked for. A brand without it gets one slide and nothing to drop.
        prs = Presentation(io.BytesIO(data))
        if len(prs.slides) > 1:
            _drop_slide(prs, len(prs.slides) - 1)
        buf = io.BytesIO()
        prs.save(buf)
        return Response(
            content=buf.getvalue(),
            media_type=PPTX_MEDIA,
            headers={"Content-Disposition": f'attachment; filename="{layout}.pptx"'},
        )
    except Exception as e:  # noqa: BLE001 — surface a clean error to the client
        return JSONResponse({"feil": f"Could not export the slide: {e}"}, status_code=500)


@app.post("/slides/export-all")
async def slides_export_all(payload: dict, x_deck_token: str | None = Header(default=None)):
    """Export EVERY standard layout (all 42; the fixed benefits_verbatim brand asset excluded, same
    as the single-slide export) as ONE multi-slide .pptx, in catalog order — the bulk version of
    /slides/export. Edit any subset of slides in PowerPoint, then use the existing team-slides
    upload flow (which already lets you tick individual slides out of a multi-slide file and name
    each one) to save just the ones you changed, without re-picking all 42.

    { background? ("dark" | "light" | "pastel") } -> the .pptx bytes.
    """
    expected = os.environ.get("DECK_SERVICE_TOKEN")
    if expected and x_deck_token != expected:
        return JSONResponse({"feil": "Unauthorized."}, status_code=401)
    background = (payload or {}).get("background") or "dark"
    brand = ((payload or {}).get("brand") or "").strip().lower() or None
    if brand and brand not in config.known_brands():
        return JSONResponse({"feil": f"Unknown brand '{brand}'."}, status_code=400)
    try:
        from pptx import Presentation

        from src import renderer
        # Only layouts THIS brand's catalog has: Revervia has no `ingredient`, and offering a
        # layout its template cannot build would fail the whole export.
        known = set(config.catalog(brand))
        from build_gallery import localize_sample
        slides = [localize_sample(s, brand) for key, s in _gallery_samples().items()
                  if key != "benefits_verbatim" and key in known]
        if background != "dark":
            for s in slides:
                s["background"] = background
        n = len(slides)
        product = brand_mod.theme(brand)["product"]
        data = renderer.render_deck(
            {"deck_title": f"{product} slide templates", "language": "en", "slides": slides},
            brand=brand,
        )
        # Same benefits-splice quirk as /slides/export, generalised: render_deck always appends the
        # verbatim benefits overview then moves it to the SECOND-TO-LAST position (max(1, n - 1),
        # n = the number of content slides fed in) — drop exactly that index to get back all n (and
        # only those n) slides, in their original order.
        prs = Presentation(io.BytesIO(data))
        if len(prs.slides) > n:            # only when the benefits slide was actually spliced
            _drop_slide(prs, max(1, n - 1))
        buf = io.BytesIO()
        prs.save(buf)
        fname = f"{(brand or config.DEFAULT_BRAND)}-slide-templates.pptx"
        return Response(
            content=buf.getvalue(),
            media_type=PPTX_MEDIA,
            headers={"Content-Disposition": f'attachment; filename="{fname}"'},
        )
    except Exception as e:  # noqa: BLE001 — surface a clean error to the client
        return JSONResponse({"feil": f"Could not export the slide templates: {e}"}, status_code=500)


@app.post("/slides/extract")
async def slides_extract(
    file: UploadFile,
    slide_index: int = Form(...),
    x_deck_token: str | None = Header(default=None),
):
    """Pull ONE slide out of an uploaded .pptx as its own standalone file — the other half of the
    edit round trip, for a team slide already sitting in the library: download just this slide,
    edit it, re-upload it to replace what's stored. { file, slide_index } -> the .pptx bytes."""
    expected = os.environ.get("DECK_SERVICE_TOKEN")
    if expected and x_deck_token != expected:
        return JSONResponse({"feil": "Unauthorized."}, status_code=401)
    data = await file.read()
    if not data:
        return JSONResponse({"feil": "Empty file."}, status_code=400)
    try:
        from pptx import Presentation

        prs = Presentation(io.BytesIO(data))
        n = len(prs.slides)
        if not (0 <= slide_index < n):
            return JSONResponse(
                {"feil": f"Slide {slide_index + 1} does not exist (the file has {n})."}, status_code=400
            )
        for i in range(n - 1, -1, -1):
            if i != slide_index:
                _drop_slide(prs, i)
        buf = io.BytesIO()
        prs.save(buf)
        return Response(
            content=buf.getvalue(),
            media_type=PPTX_MEDIA,
            headers={"Content-Disposition": 'attachment; filename="slide.pptx"'},
        )
    except Exception as e:  # noqa: BLE001 — surface a clean error to the client
        return JSONResponse({"feil": f"Could not extract the slide: {e}"}, status_code=500)


@app.post("/studies/extract-text")
async def studies_extract_text(
    file: UploadFile,
    x_deck_token: str | None = Header(default=None),
):
    """Extract full text from an uploaded study PDF, plus the paper's OWN abstract (copied
    verbatim, not paraphrased), year and authors — the "Add study" flow on the Scientific
    Studies page, for a study the site otherwise has no way to carry (no PMID, or AKBM never
    supplied it as part of the curated/full text library). Pure text extraction (PyMuPDF) + one
    Claude call to locate/transcribe those three fields; no plan/schema involved, same
    server-to-server shape as /slides/inspect (the browser uploads to Storage first; this
    route's only caller is Vercel).
    { file } -> {full_text, abstract, abstract_verified, year, authors, pages, chars}."""
    expected = os.environ.get("DECK_SERVICE_TOKEN")
    if expected and x_deck_token != expected:
        return JSONResponse({"feil": "Unauthorized."}, status_code=401)
    data = await file.read()
    if not data:
        return JSONResponse({"feil": "Empty file."}, status_code=400)

    try:
        import fitz  # PyMuPDF, already a dependency

        doc = fitz.open(stream=data, filetype="pdf")
        pages = [doc[i].get_text() for i in range(doc.page_count)]
        page_count = doc.page_count
        full_text = "\n\n".join(pages).strip()
    except Exception as e:  # noqa: BLE001 — surface a clean error to the client
        return JSONResponse({"feil": f"Could not read the PDF: {e}"}, status_code=400)
    if len(full_text) < 200:
        return JSONResponse(
            {"feil": "No extractable text found — the PDF may be a scanned image with no text layer."},
            status_code=400,
        )

    abstract, year, authors, abstract_verified = "", "", "", False
    try:
        client = anthropic.Anthropic()
        msg = client.messages.create(
            model=config.MODEL,
            max_tokens=800,
            messages=[{
                "role": "user",
                "content": (
                    "From the study text below, extract three things:\n"
                    "1. \"abstract\": the paper's OWN Abstract section, copied EXACTLY word for "
                    "word, character for character, from the source text. Do not paraphrase, "
                    "summarize, fix typos, or write one yourself. If there is no clearly labelled "
                    "abstract in the text, return an empty string.\n"
                    "2. \"year\": the 4 digit publication year, your best reading of the source. "
                    "Empty string if you cannot find one.\n"
                    "3. \"authors\": the author list as \"Surname Initial, Surname Initial, et "
                    "al.\" (up to 3 named, then \"et al.\" if there are more). Empty string if "
                    "you cannot find one.\n\n"
                    f"<paper>\n{full_text[:60000]}\n</paper>\n\n"
                    "Return ONLY a JSON object: "
                    "{\"abstract\": \"...\", \"year\": \"...\", \"authors\": \"...\"}"
                ),
            }],
        )
        raw = "".join(b.text for b in msg.content if b.type == "text")
        parsed = json.loads(raw[raw.index("{"):raw.rindex("}") + 1])
        abstract = (parsed.get("abstract") or "").strip()
        year = (parsed.get("year") or "").strip()
        authors = (parsed.get("authors") or "").strip()

        # Deterministic verbatim check (same normalized substring test the claims library uses
        # for quotes) — if the model paraphrased anyway, don't pass it off as the paper's own
        # words. The UI still shows it, flagged, so a reviewer can fix it up rather than losing it.
        if abstract:
            norm = lambda s: re.sub(r"\s+", " ", s).lower().strip()  # noqa: E731
            abstract_verified = norm(abstract) in norm(full_text)
    except Exception:  # noqa: BLE001 — text extraction still succeeds without these
        pass

    return {
        "full_text": full_text,
        "abstract": abstract,
        "abstract_verified": abstract_verified,
        "year": year,
        "authors": authors,
        "pages": page_count,
        "chars": len(full_text),
    }


@app.get("/jobs/{job_id}")
def job_status(job_id: str):
    j = JOBS.get(job_id)
    if not j:
        return JSONResponse({"feil": "Unknown or expired job."}, status_code=404)
    return {"status": j["status"], "progress": j.get("progress", 0),
            "step": j.get("step", ""), "filename": j.get("filename"), "error": j.get("error")}


@app.get("/jobs/{job_id}/result")
def job_result(job_id: str):
    """Repeatable, not one-shot: the frontend keeps this URL around as a plain clickable link (so a
    non-technical user can re-open their deck without hunting through a downloads folder), so it must
    survive more than one GET. Files are cleaned up by _prune_jobs()'s JOB_TTL_SECONDS (1h)."""
    j = JOBS.get(job_id)
    if not j:
        return JSONResponse({"feil": "Unknown or expired job."}, status_code=404)
    if j.get("status") != "done":
        return JSONResponse({"feil": f"Job is {j.get('status')}, not ready."}, status_code=409)
    # Try result_path first (new storage to disk), then fall back to legacy in-memory result.
    data = None
    if "result_path" in j:
        try:
            data = Path(j["result_path"]).read_bytes()
        except Exception:  # noqa: BLE001 — file may have been deleted prematurely
            return JSONResponse({"feil": "Result file not found or unreadable."}, status_code=410)
    elif "result" in j:
        data = j["result"]
    if data is None:
        return JSONResponse({"feil": "Result data is missing."}, status_code=500)
    media, filename = j["media_type"], j["filename"]
    return Response(content=data, media_type=media,
                    headers={"Content-Disposition": f'attachment; filename="{filename}"'})


@app.post("/blog/docx")
def blog_docx(
    markdown: str = Form(...),
    filename: str = Form(default="superba-blog-draft"),
    x_deck_token: str | None = Header(default=None),
):
    """Convert a (possibly edited) Markdown blog draft to a Word .docx. Pure conversion, no LLM —
    lets the frontend hand back the reviewed draft and get the Word deliverable back."""
    expected = os.environ.get("DECK_SERVICE_TOKEN")
    if expected and x_deck_token != expected:
        return JSONResponse({"feil": "Unauthorized."}, status_code=401)
    if not (markdown or "").strip():
        return JSONResponse({"feil": "No markdown provided."}, status_code=400)
    name = (filename or "superba-blog-draft").rsplit(".", 1)[0] + ".docx"
    data = src.markdown_to_docx(markdown)
    return Response(content=data, media_type=DOCX_MEDIA,
                    headers={"Content-Disposition": f'attachment; filename="{name}"'})
