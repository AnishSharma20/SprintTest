"""Render one contiguous run of a deck's slides to JPEG previews — in a SEPARATE PROCESS.

    python -m src.preview_chunk <pptx> <start> <end> <outdir>   # end exclusive

Why a subprocess rather than a function call: rendering previews for a big upload means
trimming and rasterising the deck once per chunk, and doing that in the service process leaks
badly — python-pptx/lxml hold on to a 42-slide package's parse tree, PIL holds decoded
bitmaps, and none of it is returned to the OS between chunks (measured: the service process
peaked at 337 MB and stayed at 324 MB after rendering the team's own template export, before
LibreOffice's own 150 MB+ was added on a 512 MB instance — i.e. the chunking meant to REDUCE
peak memory was accumulating it instead). A child process gives that memory back on exit
unconditionally, so peak stays flat no matter how many slides the deck has.

Writes `slide<index>.jpg` (already downscaled, the final preview bytes) so the parent only ever
reads small files and never holds a full-resolution PNG.
"""
from __future__ import annotations

import io
import sys
from pathlib import Path

_MAX_W = 1280  # the preview width the About page shows and stores


def render_chunk(pptx_path: Path, start: int, end: int, out_dir: Path) -> int:
    from pptx import Presentation
    from PIL import Image

    from . import qa_gate

    data = pptx_path.read_bytes()
    prs = Presentation(io.BytesIO(data))
    total = len(prs.slides)
    # Drop every slide outside [start, end) so LibreOffice only converts this run.
    for i in reversed(range(total)):
        if not start <= i < end:
            _drop_slide(prs, i)
    buf = io.BytesIO()
    prs.save(buf)
    del prs, data

    images = qa_gate.rasterize(buf.getvalue(), qa_gate.PREVIEW_MAX_DIM)
    if not images:
        raise RuntimeError("No slide renderer is available on the server.")
    out_dir.mkdir(parents=True, exist_ok=True)
    for n, png in enumerate(images):
        im = Image.open(io.BytesIO(png)).convert("RGB")
        if im.width > _MAX_W:
            im = im.resize((_MAX_W, round(im.height * _MAX_W / im.width)))
        im.save(str(out_dir / f"slide{start + n:04d}.jpg"), "JPEG", quality=82)
        im.close()
    return len(images)


def _drop_slide(prs, index: int) -> None:
    """Remove one slide (and its relationship) — python-pptx exposes no public API for this;
    mirrors main._drop_slide / renderer._delete_example_slides."""
    from pptx.oxml.ns import qn

    lst = prs.slides._sldIdLst
    sld_id = list(lst)[index]
    r_id = sld_id.get(qn("r:id"))
    if r_id:
        prs.part.drop_rel(r_id)
    lst.remove(sld_id)


def main(argv: list[str]) -> int:
    if len(argv) != 4:
        print(__doc__, file=sys.stderr)
        return 2
    pptx, start, end, out_dir = Path(argv[0]), int(argv[1]), int(argv[2]), Path(argv[3])
    try:
        n = render_chunk(pptx, start, end, out_dir)
    except Exception as e:  # noqa: BLE001 — the parent reads the message off stderr
        print(f"preview-chunk failed: {e}", file=sys.stderr)
        return 1
    print(n)
    return 0


if __name__ == "__main__":
    raise SystemExit(main(sys.argv[1:]))
