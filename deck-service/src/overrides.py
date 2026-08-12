"""Layout override recipes — a built-in layout whose DESIGN the team replaced.

The About page lets the team download any standard layout as a real .pptx, restyle it in
PowerPoint, and upload it back. From then on the renderer splices THEIR slide in verbatim
whenever the planner picks that layout key, but the AI still writes fresh text into its text
boxes on every use — so the layout stays a reusable recipe with per-slide content, never a
frozen picture. The pptx analogue of the IDML whitepaper fill (src/idml.py): design untouched,
only measured text slots rewritten.

This module owns the slot mechanics shared by main.py (/slides/inspect-slots), planner.py
(schema + prompt) and renderer.py (refill): which shapes are eligible, their stable ids, and
how much text fits in each.
"""
from __future__ import annotations

from pptx.oxml.ns import qn

# Layouts with a fixed structural ROLE (cover, contents, the client-spec'd summary, the verbatim
# ingredient splice): deterministic pipeline nets write their NORMAL fields, so swapping them to
# slot plans would fight the nets. Enforced at inspect time, parse time and sanitize time.
OVERRIDE_EXCLUDED = {"title", "agenda", "exec_summary", "ingredient"}

# Guard against prompt/schema bloat: each slot adds a described property to the tool schema and
# a guide line to the system prompt.
MAX_SLOTS = 24

_EMU_PER_IN = 914400


def char_budget(w_in: float | None, h_in: float | None, pt: float) -> int:
    """Chars a W x H inch box holds at `pt` — the same geometry formula the schema's per-field
    maxLengths come from (scripts/build_schema.py char_limit, honor_height semantics: bodies
    overflow for real, so lines are capped by the box height). Duplicated here on purpose:
    importing from scripts/ at runtime needs the argv dance main._gallery_samples documents."""
    if not w_in:
        return 120
    cpl = (w_in * 72) / (pt * 0.50)
    lines = max(1, int((h_in * 72) // (pt * 1.20))) if h_in else 1
    return max(8, min(800, int(cpl * lines * 0.85)))


def lines_estimate(h_in: float | None, pt: float) -> int:
    return max(1, int((h_in * 72) // (pt * 1.20))) if h_in else 1


def truncate_budget(text: str, cap: int) -> str:
    """Cap, not shrink: cut at the last sentence end (else word end) under the budget."""
    from .idml import truncate
    return truncate(text, cap)


def slide_ineligible_reason(src_slide) -> str | None:
    """A slide the splice cannot carry faithfully: any INTERNAL relationship that is not an
    image (an embedded chart, video, OLE object) — the same test renderer._splice_shapes uses
    to refuse a strict splice. Rejected up front so a design override can never silently
    degrade to the frozen-picture fallback."""
    rmap = dict(src_slide.part.rels.items())
    r_attrs = (qn("r:embed"), qn("r:link"), qn("r:id"))
    for shp in src_slide.shapes:
        for node in shp._element.iter():
            for a in r_attrs:
                rid = node.get(a)
                if rid and rid in rmap:
                    rel = rmap[rid]
                    if not rel.is_external and "image" not in rel.reltype:
                        return ("the slide contains an embedded chart, video or object that "
                                "cannot be carried over with editable text")
    return None


def _visible_text(txBody) -> str:
    """Paragraph texts joined with newlines, EXCLUDING field runs (slide numbers/dates) — a
    box whose only text is a field is chrome, not content."""
    paras = []
    for p in txBody.findall(qn("a:p")):
        runs = [t.text or "" for r in p.findall(qn("a:r")) for t in r.findall(qn("a:t"))]
        paras.append("".join(runs))
    return "\n".join(paras).strip()


def _max_run_pt(txBody) -> float | None:
    """Largest EXPLICIT run size in the shape (conservative: a bigger font means a smaller
    budget). Falls back to paragraph defRPr. Our own exports always force run sizes."""
    best = None
    for tag in ("a:rPr", "a:defRPr", "a:endParaRPr"):
        for el in txBody.iter(qn(tag)):
            sz = el.get("sz")
            if sz:
                pt = int(sz) / 100
                best = pt if best is None or pt > best else best
    return best


def extract_slots(src_slide) -> list[dict]:
    """Every AI-refillable text slot on the slide: id (stable through the deep-copy splice —
    _splice_shapes preserves cNvPr@id verbatim), what it says now, and how much text fits.
    Walks groups recursively with the group transform applied, so a resized group's children
    get budgets measured at their RENDERED size, not their raw child coordinates."""
    slots: list[dict] = []
    from .renderer import DISCLAIMER

    def geom(shape_el):
        xfrm = shape_el.find(qn("p:spPr") + "/" + qn("a:xfrm"))
        if xfrm is None:
            return None
        off, ext = xfrm.find(qn("a:off")), xfrm.find(qn("a:ext"))
        if off is None or ext is None:
            return None
        return (int(off.get("x")), int(off.get("y")), int(ext.get("cx")), int(ext.get("cy")))

    def walk(shapes, sx: float, sy: float, ox: float, oy: float, in_group: bool):
        # Effective position of a child = (child_pos - chOff) * scale + group_off, composed
        # through nested groups via the accumulated (sx, sy, ox, oy) transform.
        for shp in shapes:
            el = shp._element
            if el.tag == qn("p:grpSp"):
                gx = el.find(qn("p:grpSpPr") + "/" + qn("a:xfrm"))
                if gx is None:
                    walk(shp.shapes, sx, sy, ox, oy, True)
                    continue
                off, ext = gx.find(qn("a:off")), gx.find(qn("a:ext"))
                ch_off, ch_ext = gx.find(qn("a:chOff")), gx.find(qn("a:chExt"))
                if None in (off, ext, ch_off, ch_ext):
                    walk(shp.shapes, sx, sy, ox, oy, True)
                    continue
                cw, chh = max(1, int(ch_ext.get("cx"))), max(1, int(ch_ext.get("cy")))
                nsx = sx * int(ext.get("cx")) / cw
                nsy = sy * int(ext.get("cy")) / chh
                nox = ox + sx * int(off.get("x")) - nsx * int(ch_off.get("x"))
                noy = oy + sy * int(off.get("y")) - nsy * int(ch_off.get("y"))
                walk(shp.shapes, nsx, nsy, nox, noy, True)
                continue
            if el.tag != qn("p:sp"):
                continue
            txBody = el.find(qn("p:txBody"))
            if txBody is None:
                continue
            text = _visible_text(txBody)
            if not text or text == DISCLAIMER:
                continue
            cnv = el.find(qn("p:nvSpPr") + "/" + qn("p:cNvPr"))
            if cnv is None or not cnv.get("id"):
                continue
            g = geom(el)
            if g is None and not in_group:
                # A placeholder inherits geometry from the source file's own layout —
                # python-pptx resolves that chain for slide placeholders.
                try:
                    if None not in (shp.left, shp.top, shp.width, shp.height):
                        g = (int(shp.left), int(shp.top), int(shp.width), int(shp.height))
                except (AttributeError, TypeError, ValueError):
                    g = None
            if g is not None:
                x, y, w, h = (ox + sx * g[0], oy + sy * g[1], sx * g[2], sy * g[3])
                w_in, h_in = w / _EMU_PER_IN, h / _EMU_PER_IN
                pos = {"x": round(x / _EMU_PER_IN, 2), "y": round(y / _EMU_PER_IN, 2),
                       "w": round(w_in, 2), "h": round(h_in, 2)}
            else:
                w_in = h_in = None
                pos = None
            pt = _max_run_pt(txBody) or 14.0
            slots.append({
                "slot_id": f"s{cnv.get('id')}",
                "name": cnv.get("name") or "",
                "original_text": text,
                "char_budget": char_budget(w_in, h_in, pt),
                "lines_estimate": lines_estimate(h_in, pt),
                "font_pt": pt,
                "pos": pos,
                "in_group": in_group,
            })

    walk(src_slide.shapes, 1.0, 1.0, 0.0, 0.0, False)
    return slots
