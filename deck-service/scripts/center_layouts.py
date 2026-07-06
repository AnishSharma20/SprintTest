# -*- coding: utf-8 -*-
"""Build template2.pptx — the Superba template with its LAYOUT placeholders centered and
center-aligned (in the slide masters), so the tool's fill-by-index produces centered decks.

Only geometry + text alignment change; placeholder indices/types are untouched, so the tool's
config-generation maps the same semantic fields (verified). To ADOPT template2 in the tool:
re-run  scripts/inspect_template.py template2.pptx  then  scripts/build_schema.py  (this
recomputes the character limits for the narrower centered boxes) and point the service at it
(DECK_TEMPLATE=template2.pptx or replace template.pptx).

    python scripts/build ... -> run:  python scripts/center_layouts.py
"""
from __future__ import annotations
import os
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
SRC = os.path.join(ROOT, "template.pptx")
OUT = os.path.join(ROOT, "template2.pptx")
SW = 13.333  # slide width in inches

# layout name -> [(placeholder idx, centered-box width in inches | None = center-align only)].
# Symmetric/statement layouts get full centering; asymmetric text-beside-picture layouts are
# left untouched (they read as left-text / right-picture by design).
TREAT = {
    'Title Slide 1':       [(0, 11.5), (1, 9.0)],
    'Section Header 1':    [(0, 11.5)],
    'Highlight Text':      [(0, 11.5)],
    'Title Only 1':        [(0, 11.5)],
    'Agenda 1':            [(0, 11.0), (13, 8.0)],
    'Text Slide 1':        [(0, 11.5), (14, 9.5)],
    'Text With Picture 3': [(0, None)],
    'Picture With Title 1':[(0, None)],
    'Two Columns':         [(0, 11.5), (14, None), (16, None), (22, None), (19, None)],
    'Three Columns':       [(0, 11.5), (14, None), (16, None), (22, None), (19, None), (26, None), (25, None)],
    'Four Columns':        [(0, 11.5), (14, None), (16, None), (22, None), (24, None),
                            (26, None), (28, None), (30, None), (32, None)],
}


def center_ph(ph, width):
    """Center-align the placeholder's text (via its list style, which slides inherit) and,
    if a width is given, centre the box horizontally while preserving its top/height."""
    for p in ph.text_frame.paragraphs:
        p.alignment = PP_ALIGN.CENTER
    txBody = ph._element.find(qn('p:txBody'))
    if txBody is not None:
        ls = txBody.find(qn('a:lstStyle'))
        if ls is None:
            ls = txBody.makeelement(qn('a:lstStyle'), {})
            bp = txBody.find(qn('a:bodyPr'))
            bp.addnext(ls) if bp is not None else txBody.insert(0, ls)
        for lvl in (1, 2, 3):
            e = ls.find(qn(f'a:lvl{lvl}pPr'))
            if e is None:
                e = ls.makeelement(qn(f'a:lvl{lvl}pPr'), {}); ls.append(e)
            e.set('algn', 'ctr')
    if width is not None:
        top, height = ph.top, ph.height
        ph.width = Inches(width); ph.left = Inches((SW - width) / 2)
        ph.top, ph.height = top, height


def main():
    prs = Presentation(SRC)
    n = 0
    for master in prs.slide_masters:              # apply on BOTH masters (dark + light)
        for lay in master.slide_layouts:
            if lay.name in TREAT:
                byidx = {ph.placeholder_format.idx: ph for ph in lay.placeholders}
                for idx, width in TREAT[lay.name]:
                    if idx in byidx:
                        center_ph(byidx[idx], width); n += 1
    prs.save(OUT)
    print(f"centered {n} placeholders across {len(prs.slide_masters)} masters -> {os.path.relpath(OUT, ROOT)}")


if __name__ == "__main__":
    main()
