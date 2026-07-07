"""Build assets/benefits_slide.pptx — AKBM's standard "Multiple, Proven Health Benefits" overview
slide (50+ human clinical trials, benefit hexagon cards), extracted VERBATIM so the renderer can
splice it into every deck as the standard second-to-last slide.

Same idea as build_ingredient_slide.py: splice the slide's shapes (re-embedding images) AND its
background into a fresh 16:9 presentation, so the source deck's ~50 unrelated master images are
dropped. This slide is a WHITE infographic, so we carry its <p:bg> across too.

    python scripts/build_benefits_slide.py [path/to/source.pptx] [slide_number]
"""
from __future__ import annotations

import copy
import io
import sys
from pathlib import Path

from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Inches

ROOT = Path(__file__).resolve().parent.parent
OUT = ROOT / "assets" / "benefits_slide.pptx"

MARKER = "PROVEN HEALTH BENEFITS"
CANDIDATES = [
    ROOT / "reference" / "AKBM_EKSEMPEL_SLIDES_FEED_IN.pptx",
    ROOT / "brand_assets" / "AKBM EKSEMPEL SLIDES FEED IN.pptx",
    Path.home() / "Downloads" / "AKBM EKSEMPEL SLIDES FEED IN.pptx",
]
_R = (qn("r:embed"), qn("r:link"), qn("r:id"))


def splice(dst_prs, dst_layout, src_slide):
    dst_slide = dst_prs.slides.add_slide(dst_layout)
    for ph in list(dst_slide.shapes):
        ph._element.getparent().remove(ph._element)
    # carry the source slide background (this slide is a white infographic)
    dst_cSld = dst_slide._element.find(qn("p:cSld"))
    src_cSld = src_slide._element.find(qn("p:cSld"))
    src_bg = src_cSld.find(qn("p:bg")) if src_cSld is not None else None
    if src_bg is not None:
        for old in dst_cSld.findall(qn("p:bg")):
            dst_cSld.remove(old)
        dst_cSld.insert(0, copy.deepcopy(src_bg))
    spTree = dst_slide.shapes._spTree
    rmap = dict(src_slide.part.rels.items())
    for shp in src_slide.shapes:
        el = copy.deepcopy(shp._element)
        for node in el.iter():
            for a in _R:
                if a in node.attrib:
                    rel = rmap.get(node.get(a))
                    if rel is None:
                        continue
                    if rel.is_external:
                        new = dst_slide.part.rels.get_or_add_ext_rel(rel.reltype, rel.target_ref)
                    else:
                        _, new = dst_slide.part.get_or_add_image_part(io.BytesIO(rel._target.blob))
                    node.set(a, new)
        spTree.append(el)
    return dst_slide


def find_slide(prs, override_no):
    if override_no:
        return prs.slides[override_no - 1]
    for sl in prs.slides:
        for sh in sl.shapes:
            if sh.has_text_frame and MARKER in sh.text.upper():
                return sl
    raise SystemExit(f"Could not find a slide containing {MARKER!r} in the source deck.")


def main(argv):
    src_path = Path(argv[0]) if argv else next((p for p in CANDIDATES if p.exists()), None)
    if not src_path or not src_path.exists():
        raise SystemExit("Source deck not found. Pass it explicitly:\n  python scripts/"
                         "build_benefits_slide.py <source.pptx> [slide_no]\nTried:\n  "
                         + "\n  ".join(str(c) for c in CANDIDATES))
    override_no = int(argv[1]) if len(argv) > 1 else None
    src = Presentation(str(src_path))
    src_slide = find_slide(src, override_no)
    has_bg = src_slide._element.find(qn("p:cSld")).find(qn("p:bg")) is not None
    dst = Presentation()
    dst.slide_width, dst.slide_height = Inches(13.333), Inches(7.5)
    splice(dst, dst.slide_layouts[6], src_slide)  # layout 6 = "Blank"
    OUT.parent.mkdir(parents=True, exist_ok=True)
    dst.save(str(OUT))
    kb = OUT.stat().st_size // 1024
    print(f"Wrote {OUT.relative_to(ROOT)} ({kb} KB) from {src_path.name} — "
          f"{len(Presentation(str(OUT)).slides[0].shapes)} shapes, source has <p:bg>: {has_bg}")


if __name__ == "__main__":
    main(sys.argv[1:])
