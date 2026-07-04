"""Step 1 — template introspection.

Open the Superba PowerPoint template with python-pptx and dump, per slide layout:
the layout name, its master, and every placeholder's idx / type / name / dimensions
(EMU and inches). Also read the theme (fonts + colour scheme) straight from the .pptx,
and measure the text lengths of the bundled example slides per (layout, placeholder) so
Step 3 can derive character limits empirically instead of guessing.

    python scripts/inspect_template.py [path/to/template.pptx]

Writes config/template_inventory.json (the renderer's ground truth) and prints a summary.
Nothing here is Superba-specific: point it at any .pptx to inventory that template.
"""
from __future__ import annotations

import json
import sys
import zipfile
from collections import defaultdict
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.util import Emu

# UTF-8 stdout so layout/font names never crash on a cp1252 console (Windows).
try:
    sys.stdout.reconfigure(encoding="utf-8")
except Exception:  # noqa: BLE001
    pass

ROOT = Path(__file__).resolve().parent.parent
DEFAULT_TEMPLATE = ROOT / "brand_assets" / "05. Superba Brand Identity" / "Superba refresh power point template.pptx"
OUT = ROOT / "config" / "template_inventory.json"

A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"


def _emu_to_in(v):
    return round(Emu(v).inches, 3) if v is not None else None


def _clr(node):
    """Pull a hex colour out of an <a:srgbClr>/<a:sysClr> theme colour node."""
    if node is None:
        return None
    srgb = node.find(f"{A}srgbClr")
    if srgb is not None:
        return "#" + srgb.get("val", "").upper()
    sysclr = node.find(f"{A}sysClr")
    if sysclr is not None:
        return "#" + (sysclr.get("lastClr") or "").upper()
    return None


def master_major_font(master) -> str | None:
    """The major (heading) typeface of a master's theme — used to tell the branded
    Superba master (Exo 2) apart from a leftover default Office (Calibri) master."""
    for rel in master.part.rels.values():
        if rel.reltype.endswith("/theme"):
            root = etree.fromstring(rel.target_part.blob)
            latin = root.find(f".//{A}fontScheme/{A}majorFont/{A}latin")
            return latin.get("typeface") if latin is not None else None
    return None


def read_themes(pptx_path: Path) -> list[dict]:
    """Read every ppt/theme/themeN.xml: colour scheme + major/minor fonts."""
    themes = []
    with zipfile.ZipFile(pptx_path) as z:
        names = sorted(n for n in z.namelist() if n.startswith("ppt/theme/theme") and n.endswith(".xml"))
        for name in names:
            root = etree.fromstring(z.read(name))
            clr = root.find(f".//{A}clrScheme")
            colors = {}
            if clr is not None:
                for child in clr:
                    tag = etree.QName(child).localname
                    colors[tag] = _clr(child)
            fonts = {}
            fs = root.find(f".//{A}fontScheme")
            if fs is not None:
                for role in ("major", "minor"):
                    latin = fs.find(f"{A}{role}Font/{A}latin")
                    if latin is not None:
                        fonts[role] = latin.get("typeface")
            themes.append({"part": name, "scheme": (clr.get("name") if clr is not None else None) if False else None,
                           "fonts": fonts, "colors": colors})
    return themes


def placeholder_rows(container, master_dims=None):
    """One row per placeholder in a layout/master, resolving inherited dims from the master."""
    rows = []
    for ph in container.placeholders:
        pf = ph.placeholder_format
        idx = pf.idx
        l, t, w, h = ph.left, ph.top, ph.width, ph.height
        inherited = False
        if master_dims and (l is None or t is None or w is None or h is None):
            md = master_dims.get(idx)
            if md:
                l = l if l is not None else md[0]
                t = t if t is not None else md[1]
                w = w if w is not None else md[2]
                h = h if h is not None else md[3]
                inherited = True
        rows.append({
            "idx": idx,
            "type": str(pf.type).split(" ")[0].split(".")[-1] if pf.type is not None else None,
            "name": ph.name,
            "left_emu": l, "top_emu": t, "width_emu": w, "height_emu": h,
            "left_in": _emu_to_in(l), "top_in": _emu_to_in(t),
            "width_in": _emu_to_in(w), "height_in": _emu_to_in(h),
            "dims_inherited_from_master": inherited,
        })
    return rows


def main(template: Path) -> None:
    prs = Presentation(str(template))

    master_index = {id(m): i for i, m in enumerate(prs.slide_masters)}
    masters_out = []
    master_dims_by_idx = {}
    for i, m in enumerate(prs.slide_masters):
        dims = {}
        for ph in m.placeholders:
            dims[ph.placeholder_format.idx] = (ph.left, ph.top, ph.width, ph.height)
        master_dims_by_idx[i] = dims
        masters_out.append({"index": i, "name": m.name, "major_font": master_major_font(m),
                            "placeholders": placeholder_rows(m)})

    superba_idx = next((mo["index"] for mo in masters_out
                        if (mo["major_font"] or "").lower().startswith("exo")), 0)

    layouts_out = []
    for i, m in enumerate(prs.slide_masters):
        for layout in m.slide_layouts:
            layouts_out.append({
                "name": layout.name,
                "master_index": i,
                "is_superba_master": i == superba_idx,
                "placeholders": placeholder_rows(layout, master_dims_by_idx.get(i)),
            })

    # Example slides: which layout each uses + text length per placeholder (calibration).
    usage = defaultdict(int)
    layout_master_usage = defaultdict(int)                  # which master the examples actually use
    text_lengths = defaultdict(lambda: defaultdict(list))   # layout -> idx -> [char lengths]
    para_counts = defaultdict(lambda: defaultdict(list))
    for slide in prs.slides:
        lname = slide.slide_layout.name
        usage[lname] += 1
        layout_master_usage[master_index[id(slide.slide_layout.slide_master)]] += 1
        for ph in slide.placeholders:
            if ph.has_text_frame:
                txt = ph.text_frame.text or ""
                if txt.strip():
                    idx = ph.placeholder_format.idx
                    text_lengths[lname][idx].append(len(txt))
                    para_counts[lname][idx].append(len(ph.text_frame.paragraphs))

    inventory = {
        "template_path": str(template),
        "slide_size_emu": {"width": prs.slide_width, "height": prs.slide_height},
        "slide_size_in": {"width": _emu_to_in(prs.slide_width), "height": _emu_to_in(prs.slide_height)},
        "themes": read_themes(template),
        "superba_master_index": superba_idx,
        "masters": masters_out,
        "layouts": layouts_out,
        "example_slides": {
            "count": len(prs.slides._sldIdLst),
            "layout_usage": dict(sorted(usage.items(), key=lambda kv: -kv[1])),
            "master_usage": dict(layout_master_usage),
            "text_char_lengths": {ln: {str(i): v for i, v in idxs.items()} for ln, idxs in text_lengths.items()},
            "paragraph_counts": {ln: {str(i): v for i, v in idxs.items()} for ln, idxs in para_counts.items()},
        },
    }

    OUT.parent.mkdir(parents=True, exist_ok=True)
    OUT.write_text(json.dumps(inventory, indent=2, ensure_ascii=False), encoding="utf-8")

    # ---- human summary ----
    print(f"Template: {template.name}")
    print(f"Slide size: {inventory['slide_size_in']['width']}\" x {inventory['slide_size_in']['height']}\"  "
          f"({'16:9' if abs(prs.slide_width/prs.slide_height - 16/9) < 0.02 else 'other'})")
    for th in inventory["themes"]:
        print(f"\nTheme {th['part']}: fonts major={th['fonts'].get('major')!r} minor={th['fonts'].get('minor')!r}")
        keys = ["dk1", "lt1", "dk2", "lt2", "accent1", "accent2", "accent3", "accent4", "accent5", "accent6"]
        print("  colors: " + ", ".join(f"{k}={th['colors'].get(k)}" for k in keys if th['colors'].get(k)))
    print(f"\nMasters: " + "; ".join(f"#{mo['index']} font={mo['major_font']!r}" for mo in masters_out)
          + f"  ->  Superba master = #{superba_idx}")
    print(f"Example slides resolve to master(s): {dict(layout_master_usage)}")
    sup = [l for l in layouts_out if l["master_index"] == superba_idx]
    print(f"\n{len(sup)} layouts on the Superba master (#{superba_idx}):")
    for lay in sup:
        phs = ", ".join(f"[{p['idx']}]{p['type']}({p['width_in']}x{p['height_in']}\")" for p in lay["placeholders"])
        print(f"  - {lay['name']:<22}  ->  {phs or 'no placeholders'}")
    print(f"\nExample slides: {inventory['example_slides']['count']}. Layout usage:")
    for ln, c in inventory["example_slides"]["layout_usage"].items():
        print(f"  {c:>3}x  {ln}")
    print(f"\nWrote {OUT.relative_to(ROOT)}")


if __name__ == "__main__":
    tpl = Path(sys.argv[1]) if len(sys.argv) > 1 else DEFAULT_TEMPLATE
    if not tpl.exists():
        sys.exit(f"Template not found: {tpl}")
    main(tpl)
