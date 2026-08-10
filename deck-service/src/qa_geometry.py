"""Deterministic slide QA — margins, alignment, contrast and asset insertion.

`qa_gate.py` catches what only a vision model can judge (does this genuinely LOOK wrong).
This module catches what's exactly computable from the rendered deck's own shape coordinates,
fill colours and embedded image bytes, without a vision call or a rasteriser: a shape sitting
partly off the visible slide, a repeated element (icon discs, equal column boxes) out of line
with its siblings, text with insufficient contrast against what's behind it, and a
photo/icon/chart the plan called for that never actually landed on the slide (the renderer's
icon/photo helpers fail silently by design — see `_place_icon`/`_place_cropped` — so a broken
asset previously vanished without a trace).

Runs on every render, any quality level: no LLM, no LibreOffice/PowerPoint COM, just python-pptx
over the already-rendered deck. What's unambiguous and safe gets auto-fixed in place (off-canvas
shapes translated back on, misaligned siblings snapped to their row/column, low-contrast text
flipped to the deck's own ink/white); everything else is returned as a plain issue list for the
caller to log or fold into the polished-mode vision gate.
"""
from __future__ import annotations

import io

from pptx import Presentation
from pptx.enum.dml import MSO_COLOR_TYPE, MSO_FILL_TYPE
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches

from . import renderer

EMU_PER_IN = 914400

# The renderer's own repeated-element size discipline ("equal box heights, one gutter") means a
# tolerance this tight still reliably groups genuinely-matching elements without false grouping
# unrelated same-ish-size shapes together.
_ALIGN_SIZE_TOL = 0.03      # inches — shapes within this size delta count as "the same element"
_ALIGN_POS_TOL = 0.05       # inches — deviation from the row/column median that counts as a defect
_ROW_COL_CLUSTER_TOL = 0.5  # inches — same-size members within this of each other on one axis are
                            # "the same row/column" before checking fine alignment; comfortably
                            # below the ~1.5in+ gap between actual grid rows in the fixed synthetic
                            # skeleton, comfortably above any bug-scale nudge, so a multi-row/
                            # multi-column grid (e.g. icon_grid's 2x3) never gets treated as one
                            # giant misaligned group — that was a real false positive, caught by
                            # testing against a genuinely correct render before this shipped.
_OVERLAP_MIN = 0.6          # fraction of a text box's area that must sit under a panel to count
                            # as "drawn on top of it" when resolving the effective background


def _to_rgb(color) -> tuple[int, int, int]:
    hexs = str(color)
    return tuple(int(hexs[i:i + 2], 16) for i in (0, 2, 4))


def _luminance(rgb: tuple[int, int, int]) -> float:
    def chan(c):
        c = c / 255.0
        return c / 12.92 if c <= 0.03928 else ((c + 0.055) / 1.055) ** 2.4
    r, g, b = (chan(c) for c in rgb)
    return 0.2126 * r + 0.7152 * g + 0.0722 * b


def _contrast_ratio(a: tuple[int, int, int], b: tuple[int, int, int]) -> float:
    la, lb = _luminance(a), _luminance(b)
    lighter, darker = max(la, lb), min(la, lb)
    return (lighter + 0.05) / (darker + 0.05)


def _rect_in(shape):
    """(left, top, width, height) in inches, or None for a shape with no fixed geometry."""
    try:
        l, t, w, h = shape.left, shape.top, shape.width, shape.height
        if None in (l, t, w, h):
            return None
        return l / EMU_PER_IN, t / EMU_PER_IN, w / EMU_PER_IN, h / EMU_PER_IN
    except Exception:  # noqa: BLE001
        return None


def _is_chrome(shape) -> bool:
    try:
        return bool(shape.is_placeholder) and shape.placeholder_format.idx in renderer.CHROME_IDX
    except Exception:  # noqa: BLE001
        return False


def _is_full_bleed(shape, sw_in: float, sh_in: float) -> bool:
    """A picture covering nearly the whole slide is a deliberate background (cover/photo_stats/
    custom-slide splice), not a layout mistake — excluded from margin/alignment checks."""
    if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
        return False
    rect = _rect_in(shape)
    if not rect:
        return False
    _, _, w, h = rect
    return w >= sw_in * 0.85 and h >= sh_in * 0.85


def _label(shape) -> str:
    try:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            return "an image"
        if getattr(shape, "has_text_frame", False) and shape.text_frame.text.strip():
            return f'the "{shape.text_frame.text.strip()[:24]}" box'
    except Exception:  # noqa: BLE001
        pass
    return "a shape"


def _check_margins(slide, sw_in: float, sh_in: float, margin_in: float):
    """Deliberately narrow: only the unambiguous case — a shape sitting partly OFF the visible
    canvas — is checked. An earlier version also flagged shapes sitting inside the nominal 0.5in
    page margin, but that fired on a genuinely correct render (a stat card's own internal padding
    legitimately sits a little inside the page margin, with no shared left-edge discipline with
    top-level content boxes) — real layouts have too much legitimate small variation from the
    exact margin line for a general geometric rule to tell "intentional" from "wrong" safely.
    margin_in is accepted for a future, better-scoped version of that check."""
    issues, changed = [], False
    for shape in slide.shapes:
        if _is_chrome(shape) or _is_full_bleed(shape, sw_in, sh_in):
            continue
        rect = _rect_in(shape)
        if rect is None:
            continue
        l, t, w, h = rect
        r, b = l + w, t + h
        off_l, off_t = max(0.0, -l), max(0.0, -t)
        off_r, off_b = max(0.0, r - sw_in), max(0.0, b - sh_in)
        if off_l or off_t or off_r or off_b:
            new_l, new_t = l + off_l - off_r, t + off_t - off_b
            shape.left, shape.top = Inches(new_l), Inches(new_t)
            issues.append({"category": "margin", "fixed": True,
                           "detail": f"{_label(shape)} extended past the slide edge; moved back "
                                     f"fully on canvas."})
            changed = True
    return issues, changed


def _snap(members, axis: str) -> list[dict]:
    values = sorted((rect[0] if axis == "left" else rect[1]) for _, rect in members)
    median = values[len(values) // 2]
    out = []
    for shape, rect in members:
        v = rect[0] if axis == "left" else rect[1]
        if abs(v - median) <= _ALIGN_POS_TOL:
            continue
        if axis == "left":
            shape.left = Inches(median)
        else:
            shape.top = Inches(median)
        out.append({"category": "alignment", "fixed": True,
                    "detail": f"{_label(shape)} was out of {axis} alignment with "
                              f"{len(members) - 1} matching elements; snapped into line."})
    return out


def _cluster(members: list, axis: int) -> list[list]:
    """Split same-size members into "the same visual row/column" groups by chaining neighbours
    within _ROW_COL_CLUSTER_TOL of each other along `axis` (0=left, 1=top) — a multi-row or
    multi-column grid of identically-sized elements (e.g. icon_grid's 2x3) must be checked one
    row/column at a time, or the grid's own intentional row-to-row offset reads as one giant
    misaligned group."""
    ordered = sorted(members, key=lambda m: m[1][axis])
    clusters: list[list] = []
    for item in ordered:
        if clusters and item[1][axis] - clusters[-1][-1][1][axis] <= _ROW_COL_CLUSTER_TOL:
            clusters[-1].append(item)
        else:
            clusters.append([item])
    return clusters


def _check_alignment(slide):
    candidates = []
    for shape in slide.shapes:
        if _is_chrome(shape):
            continue
        rect = _rect_in(shape)
        if rect is not None:
            candidates.append((shape, rect))

    # Repeated elements (icon discs, equal column boxes, chip rows) are always drawn at one fixed
    # size by the renderer's own design system — near-identical (width, height) is a strong,
    # general proxy for "these are meant to line up."
    groups: dict[tuple[int, int], list] = {}
    for shape, rect in candidates:
        key = (round(rect[2] / _ALIGN_SIZE_TOL), round(rect[3] / _ALIGN_SIZE_TOL))
        groups.setdefault(key, []).append((shape, rect))

    issues, changed = [], False
    for members in groups.values():
        if len(members) < 3:
            continue
        for row in _cluster(members, axis=1):     # rows: clustered by top proximity
            if len(row) < 3:
                continue
            tops = sorted(rect[1] for _, rect in row)
            lefts = sorted(rect[0] for _, rect in row)
            top_spread, left_spread = tops[-1] - tops[0], lefts[-1] - lefts[0]
            if left_spread > top_spread and top_spread > _ALIGN_POS_TOL:
                fixed = _snap(row, "top")
                if fixed:
                    issues.extend(fixed)
                    changed = True
        for col in _cluster(members, axis=0):      # columns: clustered by left proximity
            if len(col) < 3:
                continue
            tops = sorted(rect[1] for _, rect in col)
            lefts = sorted(rect[0] for _, rect in col)
            top_spread, left_spread = tops[-1] - tops[0], lefts[-1] - lefts[0]
            if top_spread > left_spread and left_spread > _ALIGN_POS_TOL:
                fixed = _snap(col, "left")
                if fixed:
                    issues.extend(fixed)
                    changed = True
    return issues, changed


def _solid_fill_rgb(shape):
    try:
        fill = shape.fill
        if fill.type != MSO_FILL_TYPE.SOLID:
            return None
        color = fill.fore_color
        if color.type != MSO_COLOR_TYPE.RGB:
            return None
        return _to_rgb(color.rgb)
    except Exception:  # noqa: BLE001
        return None


def _rect_overlap_frac(inner, outer) -> float:
    il, it, iw, ih = inner
    ol, ot, ow, oh = outer
    xa, ya = max(il, ol), max(it, ot)
    xb, yb = min(il + iw, ol + ow), min(it + ih, ot + oh)
    if xb <= xa or yb <= ya:
        return 0.0
    area = iw * ih
    return ((xb - xa) * (yb - ya)) / area if area else 0.0


def _effective_bg(slide, shapes, index, rect):
    """The colour actually behind a piece of text: its own fill if it has one, else the topmost
    solid-filled shape drawn earlier (lower z-order) that this text visually sits on top of, else
    the slide's own background. None when nothing resolvable (e.g. text over a photo)."""
    own = _solid_fill_rgb(shapes[index])
    if own:
        return own
    best = None
    for j in range(index):
        crect = _rect_in(shapes[j])
        if not crect or _rect_overlap_frac(rect, crect) < _OVERLAP_MIN:
            continue
        rgb = _solid_fill_rgb(shapes[j])
        if rgb:
            best = rgb  # keep overwriting -> the last (topmost-under) match wins
    if best:
        return best
    try:
        return _solid_fill_rgb(slide.background)
    except Exception:  # noqa: BLE001
        return None


def _check_contrast(slide, ink_rgb, white_rgb):
    issues, changed = [], False
    shapes = list(slide.shapes)
    for i, shape in enumerate(shapes):
        if not getattr(shape, "has_text_frame", False):
            continue
        text = shape.text_frame.text
        if not text or not text.strip():
            continue
        rect = _rect_in(shape)
        if rect is None:
            continue
        bg = _effective_bg(slide, shapes, i, rect)
        if bg is None:
            continue  # over a photo/gradient we can't resolve -> unverifiable, not a false flag
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if not run.text.strip():
                    continue
                try:
                    fg = _to_rgb(run.font.color.rgb)
                except Exception:  # noqa: BLE001
                    continue  # inherited/theme colour on a native placeholder -> out of scope
                size_pt = run.font.size.pt if run.font.size else 14
                large = size_pt >= 18 or (size_pt >= 14 and run.font.bold)
                need = 3.0 if large else 4.5
                ratio = _contrast_ratio(fg, bg)
                if ratio >= need:
                    continue
                ink_ratio = _contrast_ratio(_to_rgb(ink_rgb), bg)
                white_ratio = _contrast_ratio(_to_rgb(white_rgb), bg)
                pick, picked_ratio = ((ink_rgb, ink_ratio) if ink_ratio >= white_ratio
                                      else (white_rgb, white_ratio))
                if picked_ratio > ratio:
                    run.font.color.rgb = pick
                    changed = True
                    issues.append({"category": "contrast", "fixed": True,
                                   "detail": f"{_label(shape)} text was {ratio:.1f}:1 against its "
                                             f"background (needs {need:.1f}:1); switched to a "
                                             f"higher-contrast colour ({picked_ratio:.1f}:1)."})
                else:
                    issues.append({"category": "contrast", "fixed": False,
                                   "detail": f"{_label(shape)} text is {ratio:.1f}:1 against its "
                                             f"background (needs {need:.1f}:1)."})
    return issues, changed


def _iter_pictures(shapes):
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            yield shape
        elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_pictures(shape.shapes)


def _expects_icon(spec: dict) -> bool:
    if (spec.get("benefit") or "none") not in ("none", ""):
        return True
    # Different layouts name their icon-bearing list differently (columns for e.g. `comparison`,
    # items for icon_grid/pillars/key_points/numbered_cards) — check both rather than assume one.
    for key in ("columns", "items"):
        for entry in spec.get(key) or []:
            if not isinstance(entry, dict):
                continue
            if (entry.get("icon") or entry.get("icon_generic") or "none") not in ("none", ""):
                return True
    return False


def _check_assets(slide, spec: dict):
    issues = []
    pics = list(_iter_pictures(slide.shapes))
    if not pics:
        if _expects_icon(spec):
            issues.append({"category": "asset", "fixed": False,
                           "detail": "the plan calls for a benefit/column icon here, but no image "
                                     "was inserted (icon resolution likely failed silently)."})
        if spec.get("asset_id"):
            issues.append({"category": "asset", "fixed": False,
                           "detail": "the plan calls for a photo here, but no image was inserted."})
    for pic in pics:
        try:
            data = pic.image.blob
            if not data:
                raise ValueError("empty image")
            from PIL import Image
            with Image.open(io.BytesIO(data)) as im:
                im.verify()
        except Exception:  # noqa: BLE001
            issues.append({"category": "asset", "fixed": False,
                           "detail": f"{_label(pic)} could not be decoded — likely corrupted."})
    for shape in slide.shapes:
        if not getattr(shape, "has_chart", False):
            continue
        try:
            has_data = any(any(v is not None for v in s.values) for s in shape.chart.plots[0].series)
        except Exception:  # noqa: BLE001
            has_data = False
        if not has_data:
            issues.append({"category": "asset", "fixed": False,
                           "detail": "a chart on this slide has no data points."})
    return issues


def review_and_fix(pptx_bytes: bytes, plan: dict,
                    slide_map: list[int | None] | None = None) -> tuple[bytes, list[dict]]:
    """Run every deterministic check over the rendered deck, auto-fixing what's unambiguous.
    Returns the (possibly modified) pptx bytes and the full issue list — fixed and unfixed,
    slide-numbered (1-based) — for the caller to log or hand to the vision gate. Never raises:
    a check that can't resolve something (an odd shape, a colour it can't read) just skips it.

    slide_map: `renderer.render_deck(..., return_slide_map=True)`'s second return value — the
    plan['slides'] index each FINAL rendered slide came from, or None for a slide with no plan
    entry (the benefits overview, an "always" team slide, an appendix figure). render_deck always
    splices in at least the benefits slide, so rendered-slide index and plan-slide index diverge
    from that point on in every real deck — pass this whenever you have it. Without it, the asset
    check falls back to naive 1:1 indexing (only correct for a plan rendered with no such
    splices, e.g. an ad hoc single-layout test); margin/alignment/contrast are unaffected either
    way since they never look at the plan."""
    prs = Presentation(io.BytesIO(pptx_bytes))
    sw_in, sh_in = prs.slide_width / EMU_PER_IN, prs.slide_height / EMU_PER_IN
    margin_in = renderer._MARGIN
    ink_rgb, white_rgb = renderer._INKC, renderer._WHITE
    slides_spec = plan.get("slides", [])

    issues: list[dict] = []
    changed = False
    for i, slide in enumerate(prs.slides):
        if slide_map is not None:
            plan_idx = slide_map[i] if i < len(slide_map) else None
            spec = slides_spec[plan_idx] if plan_idx is not None and plan_idx < len(slides_spec) else {}
        else:
            spec = slides_spec[i] if i < len(slides_spec) else {}
        m_issues, m_changed = _check_margins(slide, sw_in, sh_in, margin_in)
        a_issues, a_changed = _check_alignment(slide)
        c_issues, c_changed = _check_contrast(slide, ink_rgb, white_rgb)
        asset_issues = _check_assets(slide, spec)
        for issue in (*m_issues, *a_issues, *c_issues, *asset_issues):
            issue["slide"] = i + 1
            issues.append(issue)
        changed = changed or m_changed or a_changed or c_changed

    if not changed:
        return pptx_bytes, issues
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue(), issues
