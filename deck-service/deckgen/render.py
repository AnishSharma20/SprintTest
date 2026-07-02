"""python-pptx renderer: draw each planned slide with brand-styled, auto-fitting frames.

Every text frame uses ``word_wrap`` + ``TEXT_TO_FIT_SHAPE`` so variable AI text wraps
and shrinks to its box instead of overflowing the slide. All content lives on the
slide; nothing is hidden in speaker notes (notes carry the review detail only).
"""
from __future__ import annotations

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Emu, Inches, Pt

from . import config as C
from .layouts import LAYOUTS

EMU_IN = 914400


def _new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(C.SLIDE_W_IN)
    prs.slide_height = Inches(C.SLIDE_H_IN)
    return prs


def _bg(slide, color: RGBColor) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _text(slide, l, t, w, h, text, *, font=C.BODY_FONT, size=18, color=C.GREEN,
          bold=False, italic=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
          wrap=True, fit=True):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    if fit:
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb


def _bullets(slide, l, t, w, h, items, *, size=20, color=C.GREEN, gap=10):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.margin_left = tf.margin_right = Emu(0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        dot = p.add_run()
        dot.text = "•  "
        dot.font.name = C.BODY_FONT
        dot.font.size = Pt(size)
        dot.font.color.rgb = C.RED
        r = p.add_run()
        r.text = item
        r.font.name = C.BODY_FONT
        r.font.size = Pt(size)
        r.font.color.rgb = color
    return tb


def _pic(slide, path, l, t, w):
    """Insert a picture scaled to width `w` (inches); height auto-preserves aspect."""
    if path.exists():
        return slide.shapes.add_picture(str(path), Inches(l), Inches(t), width=Inches(w))
    return None


def _benefit_icon(slide, benefit: str):
    """Place the benefit hexagon icon top-right, if we have that icon."""
    if benefit in C.AVAILABLE_BENEFITS:
        _pic(slide, C.ICONS / f"{benefit}.png", 11.75, 0.62, 0.95)


def _rrect(slide, l, t, w, h, color: RGBColor):
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def _accent(slide, l, t, color=C.RED):
    _rrect(slide, l, t, 0.7, 0.11, color)


def _footer(slide, dark=True):
    """Real Superba + Aker logos, bottom corners — white on dark, green on light."""
    variant = "white" if dark else "green"
    sup = C.LOGOS / f"superba_{variant}.png"
    aker = C.LOGOS / f"aker_{variant}.png"
    if sup.exists():
        _pic(slide, sup, 0.8, 6.92, 1.7)
    else:
        _text(slide, 0.8, 6.9, 4, 0.4, "SUPERBA Krill", font=C.HEAD_FONT, size=13,
              color=(C.WHITE if dark else C.GREEN), bold=True, italic=True, fit=False)
    if aker.exists():
        _pic(slide, aker, 11.0, 6.92, 1.55)
    else:
        _text(slide, 8.5, 6.9, 4, 0.4, "Aker BioMarine", font=C.BODY_FONT, size=11,
              color=C.TURQ_D, align=PP_ALIGN.RIGHT, fit=False)


def _badge(slide, l, t, w, text):
    _rrect(slide, l, t, w, 0.6, C.TURQ_D)
    _text(slide, l, t + 0.05, w, 0.5, text, font=C.BODY_FONT, size=15, color=C.GREEN,
          bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, fit=False)


def _get(fields: dict, key: str) -> str:
    return (fields.get(key) or "").strip()


# ---- one function per layout -------------------------------------------------

def _cover(slide, f):
    _bg(slide, C.GREEN)
    # Superba wordmark as the masthead (falls back to a red accent bar if missing)
    if (C.LOGOS / "superba_white.png").exists():
        _pic(slide, C.LOGOS / "superba_white.png", 0.8, 0.8, 3.0)
    else:
        _accent(slide, 0.8, 1.05)
    _text(slide, 0.8, 2.0, 11.7, 1.9, _get(f, "TITLE"), font=C.HEAD_FONT, size=46,
          color=C.WHITE, bold=True, italic=True)
    _text(slide, 0.8, 3.9, 11, 1.2, _get(f, "SUBTITLE"), font=C.BODY_FONT, size=22, color=C.TURQ_L)
    # Aker logo bottom-right only (Superba is already the masthead)
    if (C.LOGOS / "aker_white.png").exists():
        _pic(slide, C.LOGOS / "aker_white.png", 11.0, 6.92, 1.55)


def _section(slide, f):
    _bg(slide, C.GREEN)
    _text(slide, 0.8, 2.6, 6, 0.4, _get(f, "KICKER"), font=C.BODY_FONT, size=16,
          color=C.RED, bold=True, fit=False)
    _accent(slide, 0.8, 3.0)
    _text(slide, 0.8, 3.2, 11.7, 2.0, _get(f, "SECTION_TITLE"), font=C.HEAD_FONT, size=44,
          color=C.WHITE, bold=True, italic=True)
    _footer(slide)


def _content(slide, f):
    _bg(slide, C.POLAR)
    _accent(slide, 0.8, 0.85)
    _text(slide, 0.8, 1.1, 11.7, 1.3, _get(f, "TITLE"), font=C.HEAD_FONT, size=34,
          color=C.GREEN, bold=True, italic=True)
    items = [_get(f, k) for k in ("BULLET_1", "BULLET_2", "BULLET_3", "BULLET_4", "BULLET_5")]
    items = [x for x in items if x]
    if items:
        _bullets(slide, 0.8, 2.6, 11.7, 3.9, items, size=21)
    _footer(slide, dark=False)


def _stat(slide, f):
    _bg(slide, C.GREEN)
    _text(slide, 0.8, 0.85, 11.7, 1.1, _get(f, "TITLE"), font=C.HEAD_FONT, size=32,
          color=C.WHITE, bold=True, italic=True)
    cards = []
    for n, col in ((1, C.ALT_RED), (2, C.SEA), (3, C.TURQ_D)):
        v = _get(f, f"STAT_{n}_VALUE")
        if v:
            cards.append((v, _get(f, f"STAT_{n}_LABEL"), col))
    if cards:
        gap = 0.35
        total_w = 11.7
        cw = (total_w - gap * (len(cards) - 1)) / len(cards)
        for i, (val, lab, col) in enumerate(cards):
            l = 0.8 + i * (cw + gap)
            _rrect(slide, l, 2.3, cw, 2.5, col)
            _text(slide, l + 0.25, 2.6, cw - 0.5, 1.3, val, font=C.HEAD_FONT, size=54,
                  color=C.WHITE, bold=True, italic=True, fit=False)
            _text(slide, l + 0.25, 3.95, cw - 0.5, 0.75, lab, font=C.BODY_FONT, size=16, color=C.WHITE)
    src = _get(f, "SOURCE")
    if src:
        _text(slide, 0.8, 5.15, 11.7, 0.5, src, font=C.BODY_FONT, size=15, color=C.TURQ_L, fit=False)
    _footer(slide)


def _evidence(slide, f):
    _bg(slide, C.POLAR)
    _text(slide, 0.8, 0.7, 5, 0.4, "EVIDENCE", font=C.BODY_FONT, size=14, color=C.TURQ_D,
          bold=True, fit=False)
    _accent(slide, 0.8, 1.05)
    _text(slide, 0.8, 1.3, 11.7, 1.2, _get(f, "HEADLINE"), font=C.HEAD_FONT, size=34,
          color=C.GREEN, bold=True, italic=True)
    _text(slide, 0.8, 2.6, 11.7, 1.4, _get(f, "CLAIM"), font=C.BODY_FONT, size=20, color=C.GREEN)
    badges = [_get(f, k) for k in ("BADGE_1", "BADGE_2", "BADGE_3")]
    badges = [b for b in badges if b]
    x = 0.8
    for b in badges:
        w = max(1.8, 0.28 + 0.14 * len(b))
        _badge(slide, x, 4.3, w, b)
        x += w + 0.3
    src = _get(f, "SOURCE")
    if src:
        _text(slide, 0.8, 5.3, 11.7, 0.5, src, font=C.BODY_FONT, size=14, color=C.SEA, fit=False)
    _footer(slide, dark=False)


def _two_col(slide, f):
    _bg(slide, C.POLAR)
    _accent(slide, 0.8, 0.85)
    _text(slide, 0.8, 1.1, 11.7, 1.1, _get(f, "TITLE"), font=C.HEAD_FONT, size=32,
          color=C.GREEN, bold=True, italic=True)
    cols = [("LEFT_TITLE", "LEFT_BODY", 0.8), ("RIGHT_TITLE", "RIGHT_BODY", 6.75)]
    for tkey, bkey, l in cols:
        _rrect(slide, l, 2.5, 5.75, 3.6, C.WHITE)
        _text(slide, l + 0.35, 2.8, 5.05, 0.9, _get(f, tkey), font=C.HEAD_FONT, size=22,
              color=C.SEA, bold=True, italic=True)
        _text(slide, l + 0.35, 3.75, 5.05, 2.1, _get(f, bkey), font=C.BODY_FONT, size=18, color=C.GREEN)
    _footer(slide, dark=False)


def _ending(slide, f):
    _bg(slide, C.GREEN)
    _text(slide, 0.8, 1.6, 5, 0.4, "LET'S TALK", font=C.BODY_FONT, size=16, color=C.RED,
          bold=True, fit=False)
    _accent(slide, 0.8, 2.0)
    _text(slide, 0.8, 2.25, 11.7, 1.7, _get(f, "CTA_TITLE"), font=C.HEAD_FONT, size=40,
          color=C.WHITE, bold=True, italic=True)
    _text(slide, 0.8, 4.1, 11.7, 0.7, _get(f, "CONTACT"), font=C.BODY_FONT, size=20, color=C.TURQ_L)
    dis = _get(f, "DISCLAIMER")
    if dis:
        _text(slide, 0.8, 5.4, 11.7, 1.0, dis, font=C.BODY_FONT, size=13, color=C.TURQ_D)
    _footer(slide)


_RENDERERS = {
    "cover": _cover, "section": _section, "content": _content, "stat": _stat,
    "evidence": _evidence, "two_col": _two_col, "ending": _ending,
}


def render_pptx(plan: dict) -> bytes:
    """Draw the whole plan to a native, editable .pptx and return its bytes."""
    import io
    prs = _new_prs()
    blank = prs.slide_layouts[6]
    for slide_plan in plan["slides"]:
        layout = slide_plan.get("layout")
        draw = _RENDERERS.get(layout)
        if draw is None:
            continue
        slide = prs.slides.add_slide(blank)
        draw(slide, slide_plan.get("fields", {}))
        if layout in ("content", "evidence", "stat", "two_col", "section"):
            _benefit_icon(slide, (slide_plan.get("benefit") or "").lower())
        notes = (slide_plan.get("notes") or "").strip()
        if notes:
            slide.notes_slide.notes_text_frame.text = notes
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def wording_document(plan: dict) -> str:
    """Science-review doc built from the plan: every slide's on-slide text + full notes."""
    lines = [f"# Wording review — {plan.get('deck_title', 'Superba deck')}", ""]
    for n, slide in enumerate(plan.get("slides", []), start=1):
        lines.append(f"## Slide {n} — {slide.get('layout', 'content')}")
        for key, val in (slide.get("fields") or {}).items():
            val = (val or "").strip()
            if val:
                lines.append(f"- **{key}:** {val}")
        notes = (slide.get("notes") or "").strip()
        if notes:
            lines.append("")
            lines.append(f"  _Evidence / notes:_ {notes}")
        lines.append("")
    return "\n".join(lines).strip() + "\n"
