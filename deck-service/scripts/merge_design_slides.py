"""Fold the verbatim source slides (ingredient, benefits) INTO template.pptx, so all design the
renderer uses lives in ONE file. The renderer finds them by text marker (`_find_design_slide`) and
splices them; template.pptx's own slides are deleted per-render, so these source slides only exist
to be copied from.

Idempotent-ish: skips a slide whose marker is already present. Run after (re)building the standalone
assets with build_ingredient_slide.py / build_benefits_slide.py.

    python scripts/merge_design_slides.py
"""
from __future__ import annotations

import copy
import io
import sys
from pathlib import Path

from pptx import Presentation
from pptx.oxml.ns import qn

ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(ROOT))
from src import config, renderer  # noqa: E402

_R = (qn("r:embed"), qn("r:link"), qn("r:id"))
SOURCES = [("assets/ingredient_slide.pptx", "CELLULAR NUTRIENT"),
           ("assets/benefits_slide.pptx", "PROVEN HEALTH BENEFITS")]


def _has_marker(prs, marker):
    for s in prs.slides:
        for sh in s.shapes:
            if sh.has_text_frame and marker in sh.text_frame.text.upper():
                return True
    return False


def _add(tpl, blank, asset_rel):
    src = Presentation(str(config.resolve_asset(asset_rel))).slides[0]
    dst = tpl.slides.add_slide(blank)
    for ph in list(dst.shapes):
        ph._element.getparent().remove(ph._element)
    src_bg = src._element.find(qn("p:cSld")).find(qn("p:bg"))
    if src_bg is not None:
        dcs = dst._element.find(qn("p:cSld"))
        for old in dcs.findall(qn("p:bg")):
            dcs.remove(old)
        dcs.insert(0, copy.deepcopy(src_bg))
    spTree = dst.shapes._spTree
    rmap = dict(src.part.rels.items())
    for shp in src.shapes:
        el = copy.deepcopy(shp._element)
        ok = True
        for node in el.iter():
            for a in _R:
                if a in node.attrib:
                    rel = rmap.get(node.get(a))
                    if rel is None:
                        continue
                    try:
                        if rel.is_external:
                            node.set(a, dst.part.rels.get_or_add_ext_rel(rel.reltype, rel.target_ref))
                        else:
                            _, nid = dst.part.get_or_add_image_part(io.BytesIO(rel._target.blob))
                            node.set(a, nid)
                    except Exception:  # noqa: BLE001
                        ok = False
        if ok:
            spTree.append(el)


def main():
    tpl = Presentation(str(config.template_path()))
    dark, _ = renderer._master_indices()
    blank = renderer._blank_layout(tpl, dark)
    added = 0
    for rel, marker in SOURCES:
        if _has_marker(tpl, marker):
            print(f"  already present: {marker}")
            continue
        _add(tpl, blank, rel)
        added += 1
        print(f"  merged: {rel}")
    if added:
        tpl.save(str(config.template_path()))
    print(f"template.pptx now has {len(Presentation(str(config.template_path())).slides)} slides")


if __name__ == "__main__":
    main()
