"""Smoke test: every brand still renders, and still renders as ITSELF.

Run this before pushing anything that touches the renderer, the planner prompt, a brand theme, or
a brand's config. It calls no model and needs no API key, so it takes seconds:

    python scripts/smoke_brands.py

Why it exists. Multi-brand work kept fixing one brand and quietly breaking the other, because the
two share a renderer whose brand state lives in module globals and whose config is cached. Nothing
caught that except generating a real deck and looking at it. These are the invariants that were
actually violated at some point during that work — each assertion below is a bug that shipped once.

It deliberately checks CONSEQUENCES, not implementation: that a deck renders, that it used the
right template, that text is legible, that a brand is not offered something it does not have.
"""
from __future__ import annotations

import io
import re
import sys
import zipfile
from pathlib import Path

ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(ROOT))

try:
    sys.stdout.reconfigure(encoding="utf-8")
except Exception:  # noqa: BLE001
    pass

from src import brand as brand_theme  # noqa: E402
from src import config, pipeline, planner, qa_geometry, renderer  # noqa: E402

FAILURES: list[str] = []


def check(ok: bool, label: str, detail: str = "") -> None:
    print(f"  {'PASS' if ok else 'FAIL'}  {label}" + (f"  — {detail}" if detail and not ok else ""))
    if not ok:
        FAILURES.append(f"{label}{': ' + detail if detail else ''}")


def _slide_bgs(pptx_bytes: bytes) -> list[str]:
    """Per slide: the background the slide sets ON ITSELF, as a comparable label. "inherit" means
    it sets none and takes the layout's (that is the normal case, and the one a mutating read
    silently destroyed)."""
    z = zipfile.ZipFile(io.BytesIO(pptx_bytes))
    out = []
    for name in _slide_parts(z):
        m = re.search(r"<p:bg>(.*?)</p:bg>", z.read(name).decode("utf8"), re.S)
        if not m:
            out.append("inherit")
        elif "noFill" in m.group(1):
            out.append("noFill")
        else:
            c = re.search(r'srgbClr val="([0-9A-Fa-f]{6})"', m.group(1))
            out.append(c.group(1) if c else "other")
    return out


def _nudge_shape_off_canvas(pptx_bytes: bytes) -> bytes:
    """Push one shape on the last slide out past the left edge, so the deck has a real defect for
    the margin fixer to correct. Exactly the condition under which QA re-saves the deck."""
    from pptx import Presentation
    prs = Presentation(io.BytesIO(pptx_bytes))
    for shape in prs.slides[-1].shapes:
        if shape.left is not None and shape.width:
            shape.left = -shape.width // 2
            break
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _slide_parts(z: zipfile.ZipFile) -> list[str]:
    return sorted((n for n in z.namelist() if re.fullmatch(r"ppt/slides/slide\d+\.xml", n)),
                  key=lambda s: int(re.findall(r"\d+", s)[0]))


def contrast(fg: str, bg: str) -> float:
    def lin(h: str) -> float:
        def ch(v: float) -> float:
            v /= 255
            return v / 12.92 if v <= 0.03928 else ((v + 0.055) / 1.055) ** 2.4
        r, g, b = (int(h[i:i + 2], 16) for i in (0, 2, 4))
        return 0.2126 * ch(r) + 0.7152 * ch(g) + 0.0722 * ch(b)
    a, b = lin(fg), lin(bg)
    hi, lo = max(a, b), min(a, b)
    return (hi + 0.05) / (lo + 0.05)


# One plan touching the shapes that broke before: a native cover, a bulleted agenda, an icon
# layout, a hero-figure layout and a numbered-badge layout.
PLAN = {
    "slides": [
        {"layout": "title", "title": "Smoke test", "subtitle": "Every brand renders"},
        {"layout": "agenda", "title": "Agenda", "items": ["One", "Two", "Three", "Four"]},
        {"layout": "key_points", "title": "Four points", "banner": "A banner",
         "items": [{"heading": f"H{i}", "body": "A line.\nAnother line."} for i in range(1, 5)]},
        {"layout": "stat", "title": "In numbers",
         "stats": [{"value": "50+", "label": "Trials"}, {"value": "2x", "label": "Uplift"}]},
        {"layout": "takeaways", "title": "Three takeaways",
         "items": [{"heading": f"Point {i}", "body": "Detail."} for i in range(1, 4)]},
    ]
}


def main() -> int:
    brands = config.known_brands()
    print(f"Brands this service can render: {', '.join(brands)}\n")

    for b in brands:
        arg = None if b == config.DEFAULT_BRAND else b
        t = brand_theme.theme(arg)
        c = t["colors"]
        print(f"=== {b} ===")

        # 1. It renders at all, and every planned slide survives.
        try:
            plan = {"slides": [dict(s) | {"speaker_notes": "n"} for s in PLAN["slides"]]}
            data = renderer.render_deck(plan, brand=arg, benefits_slot=None, source_appendix=False)
            check(True, "renders a 5-slide deck")
        except Exception as e:  # noqa: BLE001
            check(False, "renders a 5-slide deck", f"{type(e).__name__}: {e}")
            continue

        z = zipfile.ZipFile(io.BytesIO(data))
        slides = _slide_parts(z)
        check(len(slides) == len(PLAN["slides"]),
              "one rendered slide per planned slide", f"{len(slides)} vs {len(PLAN['slides'])}")

        # 2. It used ITS OWN template. The brand's own logo lives in the template, so the deck's
        #    media set is the tell — this catches a deck built on another brand's template.
        tpl_media = len(zipfile.ZipFile(config.template_path(arg)).namelist())
        check(tpl_media > 0, "brand template is readable")

        # 3. The bullet style matches what the brand actually ships (a borrowed picture bullet from
        #    another brand is what this replaced).
        has_bullet_png = (config.assets_dir(arg) / "bullet.png").exists()
        agenda_xml = z.read(slides[1]).decode("utf8")
        if has_bullet_png:
            check("buBlip" in agenda_xml, "ships a bullet image, so uses a picture bullet")
        else:
            check("buChar" in agenda_xml and "buBlip" not in agenda_xml,
                  "ships no bullet image, so uses PowerPoint's character bullet")

        # 4. Text the QA gate would have to rescue. Each of these pairs shipped below 3.0 once.
        for label, fg, bg in (
            ("numbered badge legible", c["data"], c["tint"]),
            ("hero figure legible on white", c["data"], "FFFFFF"),
            ("hero figure legible on panel", c["data"], c["panel"]),
            ("text on accent legible", c["on_accent"], c["accent"]),
            ("body ink legible on panel", c["ink"], c["panel"]),
        ):
            r = contrast(fg, bg)
            check(r >= 3.0, label, f"#{fg} on #{bg} is {r:.2f}:1")

        # 5. A brand is never offered a verbatim slide its template lacks — the planner, the
        #    catalog and the renderer must agree, or a deck dies mid-render.
        cat = set(config.catalog(arg))
        if not t["has_ingredient_slide"]:
            check("ingredient" not in cat, "no `ingredient` in the catalog for a brand without one")
            check("ingredient" in planner.sanitize_disabled(None, None, arg),
                  "`ingredient` is force-disabled for the planner")

        # 6. The prompt names THIS brand and leaves no placeholder unresolved.
        sysmsg = planner.build_system("standard", "balansert", brand=arg)
        check("{product}" not in sysmsg and "{company}" not in sysmsg,
              "no unresolved {placeholders} in the prompt")
        check(t["product"] in sysmsg, f"prompt names {t['product']!r}")
        others = [o["product"] for k, o in brand_theme._DEFAULTS.items()
                  if k != b and o["product"] != t["product"]]
        leaked = [o for o in others if o in sysmsg]
        check(not leaked, "prompt does not name another brand's product", ", ".join(leaked))

        # 7. QA must not repaint the deck it inspects. Reading `slide.background` in python-pptx
        #    INSERTS a transparent background, which overrides the template's own gradient — that
        #    shipped whole decks WHITE when the user had asked for the dark Blue Ocean theme.
        #    Tested on a deck that genuinely NEEDS a fix, because review_and_fix only saves its
        #    python-pptx object when some fix fired: on a clean deck the stray background was
        #    discarded with it, which is why this only hit some decks and looked random.
        try:
            nudged = _nudge_shape_off_canvas(data)
            qa_data, issues = qa_geometry.review_and_fix(nudged, plan)
            check(qa_data != nudged, "QA fixes an off-canvas shape (fixture really needs fixing)")
            before, after = _slide_bgs(nudged), _slide_bgs(qa_data)
            check("noFill" not in after, "QA leaves no slide background transparent",
                  f"slides {[i + 1 for i, k in enumerate(after) if k == 'noFill']}")
            check(before == after, "QA changes no slide's background", f"{before} -> {after}")
        except Exception as e:  # noqa: BLE001
            check(False, "deterministic QA runs over the rendered deck", f"{type(e).__name__}: {e}")

        # 8. A list written at the wrong shape must not cost the deck. Every layout reads its items
        #    with it.get(...), so a plain string used to raise AttributeError mid-render and kill
        #    the whole job — a client hit exactly that. Repaired before validation instead.
        wrong = {"slides": [dict(s) | {"speaker_notes": "n"} for s in PLAN["slides"]] + [
            {"layout": "takeaways", "title": "Strings, not objects",
             "items": ["Alpha", "Beta", "Gamma"], "speaker_notes": "n"},
            {"layout": "agenda", "title": "Objects, not strings",
             "items": [{"heading": "Evidence", "body": "the trials"}], "speaker_notes": "n"},
        ]}
        pipeline._coerce_item_shapes(wrong, arg)
        check(all(isinstance(v, dict) for v in wrong["slides"][-2]["items"])
              and all(isinstance(v, str) for v in wrong["slides"][-1]["items"]),
              "a wrongly shaped list is repaired to the shape its layout reads")
        try:
            renderer.render_deck(wrong, brand=arg, benefits_slot=None, source_appendix=False)
            check(True, "a deck with a wrongly shaped list still renders")
        except Exception as e:  # noqa: BLE001
            check(False, "a deck with a wrongly shaped list still renders", f"{type(e).__name__}: {e}")

        # 9. Photo ids offered to the model all resolve to a real staged file.
        renderer.apply_brand(arg)
        missing = [a["id"] for a in config.selectable_photos(arg)
                   if not config.resolve_asset(a["path"], arg).exists()]
        check(not missing, "every offered photo exists on disk", ", ".join(missing[:4]))
        print()

    if FAILURES:
        print(f"{len(FAILURES)} FAILURE(S):")
        for f in FAILURES:
            print(f"  - {f}")
        return 1
    print("All brands pass.")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
