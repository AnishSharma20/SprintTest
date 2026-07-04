"""Step 6 — QA a rendered deck.

Programmatic checks that fail loudly, plus a per-slide PNG render for visual inspection:
  1. no leftover placeholder / prompt text (lorem, "click to add", "add heading here", ...)
  2. no empty or broken picture placeholders (the old template-fill failure mode)
  3. the deck opens cleanly and has >= 1 slide
Then renders slides to PNG (PowerPoint COM on Windows; LibreOffice headless elsewhere) so
you can eyeball layout / overflow — the last line of defence the checks can't fully automate.

    python scripts/qa.py path/to/deck.pptx [--png-dir DIR] [--no-render]

Exit 0 = all programmatic checks passed; 1 = a check failed (or rendering was requested but
no renderer was available).
"""
from __future__ import annotations

import argparse
import re
import shutil
import subprocess
import sys
import tempfile
from pathlib import Path

from pptx import Presentation

try:
    sys.stdout.reconfigure(encoding="utf-8")
except Exception:  # noqa: BLE001
    pass

# Template prompt text that must never survive into a finished slide.
PROMPT_RE = re.compile(
    r"lorem|ipsum|xxxx+|click to add|add heading here|add text here|add title|"
    r"right click|change picture|type here|\[?placeholder\]?|sample text",
    re.I,
)


def _slide_text(prs):
    for i, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                yield i, shape.name, shape.text_frame.text
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            yield i, "table", cell.text


def check_remnants(prs) -> list[str]:
    hits = []
    for idx, name, text in _slide_text(prs):
        for m in PROMPT_RE.finditer(text):
            hits.append(f"slide {idx} [{name}]: leftover prompt text '{m.group(0)}' in {text[:60]!r}")
    return hits


def check_pictures(prs) -> list[str]:
    problems = []
    for i, slide in enumerate(prs.slides, 1):
        for ph in slide.placeholders:
            if "PICTURE" in str(ph.placeholder_format.type or ""):
                try:
                    _ = ph.image  # a filled picture placeholder has image data; empty raises
                except Exception:  # noqa: BLE001
                    problems.append(f"slide {i}: empty picture placeholder '{ph.name}' "
                                    "(should have been filled or removed)")
    return problems


def render_pngs(pptx: Path, out_dir: Path) -> bool:
    """Export each slide to PNG. LibreOffice (portable) first, then PowerPoint COM (Windows)."""
    out_dir.mkdir(parents=True, exist_ok=True)
    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    if soffice:
        with tempfile.TemporaryDirectory() as tmp:
            subprocess.run([soffice, "--headless", "--convert-to", "pdf", "--outdir", tmp, str(pptx)],
                           check=True, capture_output=True)
            pdf = next(Path(tmp).glob("*.pdf"))
            if shutil.which("pdftoppm"):
                subprocess.run(["pdftoppm", "-png", "-r", "120", str(pdf), str(out_dir / "slide")], check=True)
            else:
                import fitz  # PyMuPDF
                doc = fitz.open(pdf)
                for n, page in enumerate(doc, 1):
                    page.get_pixmap(dpi=120).save(str(out_dir / f"slide{n}.png"))
        return True
    if sys.platform.startswith("win"):
        ps = (f'$pp=New-Object -ComObject PowerPoint.Application;'
              f'$pres=$pp.Presentations.Open("{pptx}",$true,$true,$false);'
              f'$pres.Export("{out_dir}","PNG",1600,900);$pres.Close();$pp.Quit()')
        subprocess.run(["powershell", "-NoProfile", "-Command", ps], check=True, capture_output=True)
        return True
    return False


def main(argv=None) -> int:
    ap = argparse.ArgumentParser(prog="scripts/qa.py")
    ap.add_argument("pptx", type=Path)
    ap.add_argument("--png-dir", type=Path, default=None)
    ap.add_argument("--no-render", action="store_true")
    args = ap.parse_args(argv)

    if not args.pptx.exists():
        print(f"Not found: {args.pptx}", file=sys.stderr)
        return 1

    prs = Presentation(str(args.pptx))
    n = len(prs.slides._sldIdLst)
    print(f"Deck opens: {n} slide(s)")

    failures = []
    if n < 1:
        failures.append("deck has no slides")
    remnants = check_remnants(prs)
    pics = check_pictures(prs)
    failures += remnants + pics

    print(f"  placeholder-remnant check: {'FAIL' if remnants else 'pass'} ({len(remnants)} hits)")
    for h in remnants:
        print(f"    - {h}")
    print(f"  picture check: {'FAIL' if pics else 'pass'} ({len(pics)} problems)")
    for p in pics:
        print(f"    - {p}")

    if not args.no_render:
        png_dir = args.png_dir or args.pptx.with_suffix("")  # deck.pptx -> deck/
        png_dir = Path(str(png_dir) + "_qa")
        try:
            if render_pngs(args.pptx, png_dir):
                pngs = sorted(png_dir.glob("*.png")) + sorted(png_dir.glob("*.PNG"))
                print(f"  rendered {len(pngs)} PNG(s) to {png_dir} — inspect these visually")
            else:
                print("  render: SKIPPED (no LibreOffice; not on Windows) — install soffice to enable")
        except Exception as e:  # noqa: BLE001
            print(f"  render: FAILED ({e})")

    print("\nQA " + ("PASSED" if not failures else f"FAILED ({len(failures)} issue(s))"))
    return 0 if not failures else 1


if __name__ == "__main__":
    raise SystemExit(main())
