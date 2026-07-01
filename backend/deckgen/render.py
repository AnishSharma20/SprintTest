"""
render.py — deterministic renderer: deck JSON -> on-brand .pptx

No LLM here. Opens the real Superba template, picks existing layouts per section
kind, populates placeholders by idx, alternates light/dark masters for visual
rhythm. Same input always yields the same deck.

Usage:
    python render.py deck.json output.pptx [template.pptx]
"""
import sys
import json
from pptx import Presentation
from pptx.util import Pt
from pptx.oxml.ns import qn

TEMPLATE_DEFAULT = "Superba_refresh_power_point_template.pptx"

# Master index: 0 = dark set (deep-sea gradients, light text)
#               1 = light set (pale turquoise/white, dark text)
DARK, LIGHT = 0, 1

# Layout targets per section kind. (master_index, layout_name)
# Chosen from layouts confirmed present in BOTH masters so light/dark both work.
LAYOUT_MAP = {
    "cover":          (DARK,  "Title Slide 1"),
    "section_header": (DARK,  "Section Header 1"),
    "benefit_claim":  (LIGHT, "Text With Picture 1"),   # title + body + picture
    "benefit_alt":    (LIGHT, "Two Columns"),           # alt composition for rhythm
    "stat":           (DARK,  "Section Header 1"),       # big statement on dark
    "summary":        (LIGHT, "Two Columns"),
    "closing":        (DARK,  "Title Slide 1"),
}

# Placeholder idx by role, per layout (from live inspection of the template).
PLACEHOLDERS = {
    "Title Slide 1":      {"title": 0, "subtitle": 1},
    "Section Header 1":   {"title": 0, "body": 13},
    "Text With Picture 1":{"title": 0, "body": 14, "body2": 16, "picture": 17},
    "Two Columns":        {"title": 0, "body_l": 14, "body_r": 22, "pic_l": 20, "pic_r": 21},
}


def _layout(prs, master_idx, name):
    master = prs.slide_masters[master_idx]
    for l in master.slide_layouts:
        if l.name == name:
            return l
    raise KeyError(f"layout {name!r} not found in master {master_idx}")


def _ph(slide, idx):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            return ph
    return None


def _set_text(ph, lines, bold_first=False):
    """Write lines into a text placeholder as separate paragraphs (real bullets)."""
    if ph is None:
        return
    tf = ph.text_frame
    tf.clear()
    if isinstance(lines, str):
        lines = [lines]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = line
        if bold_first and i == 0:
            run.font.bold = True


def _remove_ph(slide, idx):
    """Delete an unused placeholder so empty prompt frames don't render."""
    ph = _ph(slide, idx)
    if ph is not None:
        ph._element.getparent().remove(ph._element)


def _add_slide(prs, master_idx, layout_name):
    layout = _layout(prs, master_idx, layout_name)
    return prs.slides.add_slide(layout)


def _clear_slides(prs):
    lst = prs.slides._sldIdLst
    for sldId in list(lst):
        prs.part.drop_rel(sldId.get(qn("r:id")))
        lst.remove(sldId)


def _benefit_body(section):
    """Compose the claim body: claims, then evidence line, then EFSA/source."""
    lines = list(section.get("claims", []))
    tc = section.get("trial_count", 0)
    if tc:
        lines.append(f"Backed by {tc} human clinical trial{'s' if tc != 1 else ''}.")
    if section.get("efsa_approved"):
        lines.append("EFSA-approved claim.")
    src = section.get("source", "")
    if src:
        lines.append(f"Source: {src}")
    return lines


def render(deck, out_path, template_path):
    prs = Presentation(template_path)
    _clear_slides(prs)

    # --- Cover ---
    mi, ln = LAYOUT_MAP["cover"]
    s = _add_slide(prs, mi, ln)
    p = PLACEHOLDERS[ln]
    _set_text(_ph(s, p["title"]), deck["deck_title"])
    _set_text(_ph(s, p["subtitle"]), deck.get("deck_subtitle", ""))

    # --- Sections ---
    benefit_i = 0
    for sec in deck["sections"]:
        kind = sec["kind"]

        if kind == "section_header":
            mi, ln = LAYOUT_MAP["section_header"]
            s = _add_slide(prs, mi, ln)
            p = PLACEHOLDERS[ln]
            _set_text(_ph(s, p["title"]), sec["section_title"])
            _set_text(_ph(s, p["body"]), sec.get("claims", []))

        elif kind == "benefit_claim":
            # Alternate between two layouts for rhythm.
            key = "benefit_claim" if benefit_i % 2 == 0 else "benefit_alt"
            benefit_i += 1
            mi, ln = LAYOUT_MAP[key]
            s = _add_slide(prs, mi, ln)
            p = PLACEHOLDERS[ln]
            _set_text(_ph(s, p["title"]), sec["section_title"])
            body_lines = _benefit_body(sec)
            if ln == "Text With Picture 1":
                _set_text(_ph(s, p["body"]), body_lines)
                _remove_ph(s, p["body2"])   # unused second body frame
            else:  # Two Columns
                # split claims left, evidence right
                claims = sec.get("claims", [])
                mid = (len(claims) + 1) // 2
                _set_text(_ph(s, p["body_l"]), claims[:mid] or [""])
                right = claims[mid:] + [l for l in body_lines if l not in claims]
                _set_text(_ph(s, p["body_r"]), right or [""])
                _remove_ph(s, p["pic_l"])
                _remove_ph(s, p["pic_r"])

        elif kind == "stat":
            mi, ln = LAYOUT_MAP["stat"]
            s = _add_slide(prs, mi, ln)
            p = PLACEHOLDERS[ln]
            _set_text(_ph(s, p["title"]), sec["section_title"])
            stat_lines = [f"{st['value']}  —  {st['label']}" for st in sec.get("stats", [])]
            _set_text(_ph(s, p["body"]), stat_lines or sec.get("claims", []))

        elif kind == "summary":
            mi, ln = LAYOUT_MAP["summary"]
            s = _add_slide(prs, mi, ln)
            p = PLACEHOLDERS[ln]
            _set_text(_ph(s, p["title"]), sec["section_title"])
            claims = sec.get("claims", [])
            mid = (len(claims) + 1) // 2
            _set_text(_ph(s, p["body_l"]), claims[:mid] or [""])
            _set_text(_ph(s, p["body_r"]), claims[mid:] or [""])
            _remove_ph(s, p["pic_l"])
            _remove_ph(s, p["pic_r"])

    # --- Closing ---
    mi, ln = LAYOUT_MAP["closing"]
    s = _add_slide(prs, mi, ln)
    p = PLACEHOLDERS[ln]
    _set_text(_ph(s, p["title"]), "Let\u2019s innovate together")
    _set_text(_ph(s, p["subtitle"]), "www.superbakrill.com")

    prs.save(out_path)
    return out_path


if __name__ == "__main__":
    deck_json = sys.argv[1]
    out = sys.argv[2] if len(sys.argv) > 2 else "output.pptx"
    tpl = sys.argv[3] if len(sys.argv) > 3 else TEMPLATE_DEFAULT
    with open(deck_json) as f:
        deck = json.load(f)
    render(deck, out, tpl)
    print("wrote", out)
